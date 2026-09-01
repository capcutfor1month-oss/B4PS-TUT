"""Integration tests for the top-level `apply_scene_aware_edit` pilot -
synthetic fixtures for the general behaviors, plus one real-deck
integration test against the same Pinned Messages/Message Actions
benchmark that exposed the original Slice 3 formatting-loss gap. The
real file is only ever copied into `tmp_path`; it is never itself a
mutation source or target in any test."""

from __future__ import annotations

import hashlib
import shutil
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Pt

from presentation_editing_intelligence.pilot import apply_scene_aware_edit

_REPO_ROOT = Path(__file__).resolve().parents[2]
_REAL_DECK = (
    _REPO_ROOT
    / "documentation-artifacts"
    / "masterslide"
    / "old"
    / "MASTER Complete Bridge4PS Desktop-Browser Feature Tutorials.pptx"
)


def _file_hash(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()


def _build_simple_deck(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(2_000_000), Emu(400_000))
    a.text_frame.paragraphs[0].text = "Bold label: "
    a.text_frame.paragraphs[0].runs[0].font.bold = True
    a.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    run2 = a.text_frame.paragraphs[0].add_run()
    run2.text = "plain body text"
    run2.font.size = Pt(12)
    b = slide.shapes.add_textbox(Emu(500_000), Emu(1_500_000), Emu(2_000_000), Emu(400_000))
    b.text_frame.text = "Unrelated shape - must never change"
    prs.save(str(path))
    return str(path), a.shape_id, b.shape_id


def test_applies_a_formatting_preserving_edit_and_reports_no_unrelated_mutation(tmp_path):
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Bold label: plain body text, extended",
        shape_id=target_id,
    )

    assert result.status == "applied"
    assert result.protected_structure_ok is True

    prs = Presentation(output)
    target_shape = list(prs.slides[0].shapes)[0]
    assert target_shape.text_frame.paragraphs[0].runs[0].font.bold is True
    assert target_shape.text_frame.text == "Bold label: plain body text, extended"

    other_shape = list(prs.slides[0].shapes)[1]
    assert other_shape.text_frame.text == "Unrelated shape - must never change"


def test_source_deck_is_never_modified(tmp_path):
    source, target_id, _ = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")
    before_hash = _file_hash(source)

    apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Bold label: plain body text!",
        shape_id=target_id,
    )

    assert _file_hash(source) == before_hash


def test_shape_id_guard_mismatch_is_unresolved_never_retargeted(tmp_path):
    source, target_id, _ = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Bold label: plain body text!",
        shape_id=target_id + 999,  # deliberately wrong
    )

    assert result.status == "unresolved"
    assert "shape_id" in result.reason


def test_ambiguous_multi_run_change_stops_before_any_mutation(tmp_path):
    source, target_id, _ = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Completely different wording entirely",
        shape_id=target_id,
    )

    assert result.status == "unresolved"
    assert result.mutation_result is None
    import os
    assert not os.path.exists(output)


@pytest.mark.skipif(not _REAL_DECK.exists(), reason="real MasterSlide deck not present in working tree")
class TestRealPinnedMessagesBenchmark:
    """The exact same real target Documentation Intelligence Slice 3
    proved structurally, and that founder visual review found lost
    formatting on: slide index 74, shape index 12, shape id 2428."""

    def test_scene_aware_edit_preserves_formatting_and_relationships(self, tmp_path):
        source = str(tmp_path / "source.pptx")
        shutil.copyfile(str(_REAL_DECK), source)
        output = str(tmp_path / "benchmark_output.pptx")
        before_real_hash = _file_hash(_REAL_DECK)

        old_text = (
            "Pin:\nMessages bookmarked for all members of the channel to "
            "view in the Pinned Messages list."
        )
        new_text = old_text + " (approved test update)"

        result = apply_scene_aware_edit(
            source, output, slide_index=74, shape_index=12,
            expected_old_text=old_text, new_text=new_text,
            shape_id=2428,
        )

        assert result.status == "applied"
        assert result.protected_structure_ok is True
        assert result.safety_report.ok

        # The real file on disk must be byte-identical to before.
        assert _file_hash(_REAL_DECK) == before_real_hash

        prs = Presentation(output)
        shape = list(prs.slides[74].shapes)[12]
        tf = shape.text_frame

        # Bold "Pin:" label - untouched.
        assert tf.paragraphs[0].runs[0].text == "Pin:"
        assert tf.paragraphs[0].runs[0].font.bold is True

        # Bold+underlined "for all members" emphasis run - untouched.
        emphasis_run = tf.paragraphs[1].runs[1]
        assert emphasis_run.text == "for all members"
        assert emphasis_run.font.bold is True
        assert emphasis_run.font.underline is True

        # Only the trailing run grew, with the approved suffix, still Calibri.
        trailing_run = tf.paragraphs[1].runs[2]
        assert trailing_run.text.endswith("(approved test update)")
        assert trailing_run.font.name == "Calibri"

        # Full text is exactly the approved replacement.
        assert tf.text == new_text

        # Shape geometry itself is unchanged (this edit needed no reflow).
        assert (shape.left, shape.top, shape.width, shape.height) == (5308964, 866941, 3474844, 646290)

        # A real relationship this benchmark slide actually has was
        # correctly identified before the edit was ever applied.
        separates = [r for r in result.relationships if r.relation_type == "separates"]
        assert any(r.target_shape_id == 2428 for r in separates)

    def test_wrong_shape_id_guard_refuses_the_real_benchmark_target(self, tmp_path):
        source = str(tmp_path / "source.pptx")
        shutil.copyfile(str(_REAL_DECK), source)
        output = str(tmp_path / "output.pptx")

        result = apply_scene_aware_edit(
            source, output, slide_index=74, shape_index=12,
            expected_old_text="anything", new_text="anything else",
            shape_id=999999,  # not the real shape id 2428
        )
        assert result.status == "unresolved"
        import os
        assert not os.path.exists(output)
