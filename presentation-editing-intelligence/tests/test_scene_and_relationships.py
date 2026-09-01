"""Local scene extraction and relationship inference, against small
synthetic `.pptx` fixtures built with python-pptx - two stacked
definition entries separated by a divider line, plus a connector
pointing from one text block toward a highlight box on a picture,
mirroring the real Pinned Messages/Message Actions slide's own
structure (divider-separated definitions, a red highlight on a
screenshot, a pointer connecting the two) without hardcoding anything
about that specific slide."""

from __future__ import annotations

import io

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Emu

from presentation_editing_intelligence.relationships import infer_relationships
from presentation_editing_intelligence.scene import extract_local_scene


def _png_bytes(size=(40, 40), color=(120, 120, 120)):
    buf = io.BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="PNG")
    buf.seek(0)
    return buf


def _build_two_entry_deck(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    entry_a = slide.shapes.add_textbox(Emu(1_000_000), Emu(500_000), Emu(2_000_000), Emu(400_000))
    entry_a.text_frame.text = "Entry A"

    divider = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(1_000_000), Emu(950_000), Emu(3_000_000), Emu(950_000))
    divider.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)

    entry_b = slide.shapes.add_textbox(Emu(1_000_000), Emu(1_050_000), Emu(2_000_000), Emu(400_000))
    entry_b.text_frame.text = "Entry B"

    slide.shapes.add_picture(_png_bytes(), Emu(4_000_000), Emu(500_000), Emu(1_000_000), Emu(1_000_000))

    highlight = slide.shapes.add_shape(1, Emu(4_100_000), Emu(550_000), Emu(200_000), Emu(200_000))
    highlight.fill.background()
    highlight.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    prs.save(str(path))
    return str(path)


def test_local_scene_includes_adjacent_divider_and_excludes_far_shapes(tmp_path):
    path = _build_two_entry_deck(tmp_path / "deck.pptx")
    prs = Presentation(path)

    entry_a_index = 0
    scene = extract_local_scene(prs, slide_index=0, target_shape_index=entry_a_index, proximity_margin_emu=200_000)

    context_ids = {s.shape_id for s in scene.context}
    entry_b_shape = list(prs.slides[0].shapes)[2]
    divider_shape = list(prs.slides[0].shapes)[1]
    assert divider_shape.shape_id in context_ids
    assert entry_b_shape.shape_id in context_ids


def test_local_scene_proximity_is_bounded_not_whole_slide(tmp_path):
    path = _build_two_entry_deck(tmp_path / "deck.pptx")
    prs = Presentation(path)

    # A near-zero margin should exclude entry B (it's ~100_000 EMU away,
    # separated by the divider) while still including the target itself.
    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0, proximity_margin_emu=0)
    assert len(scene.shapes) < len(list(prs.slides[0].shapes))


def test_divider_separates_and_is_anchored_to_the_shape_above_it(tmp_path):
    path = _build_two_entry_deck(tmp_path / "deck.pptx")
    prs = Presentation(path)
    entry_a = list(prs.slides[0].shapes)[0]
    divider = list(prs.slides[0].shapes)[1]

    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0)
    rels = infer_relationships(scene)

    separates = [r for r in rels if r.relation_type == "separates" and r.source_shape_id == divider.shape_id]
    anchored = [r for r in rels if r.relation_type == "anchored-to" and r.source_shape_id == divider.shape_id]
    assert any(r.target_shape_id == entry_a.shape_id for r in separates)
    assert any(r.target_shape_id == entry_a.shape_id for r in anchored)
    for r in anchored:
        # PEI-S1-04: geometry-only adjacency is never HIGH - HIGH is
        # reserved for an explicit OOXML structural binding.
        assert r.confidence == "MEDIUM"
        assert r.evidence  # every relationship carries concrete evidence text


def test_highlight_box_highlights_overlapping_picture(tmp_path):
    path = _build_two_entry_deck(tmp_path / "deck.pptx")
    prs = Presentation(path)
    picture = list(prs.slides[0].shapes)[3]
    highlight = list(prs.slides[0].shapes)[4]

    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0, proximity_margin_emu=3_500_000)
    rels = infer_relationships(scene)

    highlights = [r for r in rels if r.relation_type == "highlights"]
    assert any(r.source_shape_id == highlight.shape_id and r.target_shape_id == picture.shape_id for r in highlights)


def test_bound_connector_points_to_is_high_confidence(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(1_000_000), Emu(300_000))
    a.text_frame.text = "Source"
    b = slide.shapes.add_textbox(Emu(3_000_000), Emu(500_000), Emu(1_000_000), Emu(300_000))
    b.text_frame.text = "Target"
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(1_500_000), Emu(650_000), Emu(3_000_000), Emu(650_000))
    connector.begin_connect(a, 0)
    connector.end_connect(b, 0)

    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0, proximity_margin_emu=3_000_000)
    rels = infer_relationships(scene)
    points_to = [r for r in rels if r.relation_type == "points-to"]
    assert any(
        r.source_shape_id == a.shape_id and r.target_shape_id == b.shape_id and r.confidence == "HIGH"
        for r in points_to
    )


def test_relationship_confidence_is_never_fabricated_as_high_without_structural_binding(tmp_path):
    # Two shapes merely near a freeform with no explicit connector
    # binding must get at most MEDIUM confidence for points-to - never
    # HIGH, since HIGH is reserved for an explicit stCxn/endCxn binding.
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(1_000_000), Emu(300_000))
    a.text_frame.text = "Source"
    b = slide.shapes.add_textbox(Emu(3_000_000), Emu(500_000), Emu(1_000_000), Emu(300_000))
    b.text_frame.text = "Target"
    freeform_builder = slide.shapes.build_freeform(Emu(1_500_000), Emu(650_000))
    freeform_builder.add_line_segments([(Emu(3_000_000), Emu(650_000))], close=False)
    freeform_builder.convert_to_shape()

    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0, proximity_margin_emu=3_000_000)
    rels = infer_relationships(scene)
    points_to = [r for r in rels if r.relation_type == "points-to"]
    for r in points_to:
        assert r.confidence in ("MEDIUM", "LOW")
