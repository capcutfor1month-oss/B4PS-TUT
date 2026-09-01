"""Adversarial regressions for the Codex Audit 1 findings PEI-S1-01
through PEI-S1-06 (PEI-S1-07 is the governance record itself, not code).
Each test is named after the finding it proves fixed."""

from __future__ import annotations

import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Emu, Pt

from presentation_editing_intelligence.pilot import apply_scene_aware_edit
from presentation_editing_intelligence.relationships import Relationship, infer_relationships
from presentation_editing_intelligence.safety import check_edit
from presentation_editing_intelligence.scene import (
    Geometry,
    LocalScene,
    SceneShape,
    extract_all_shapes_on_slide,
    extract_local_scene,
)
from presentation_editing_intelligence.reflow import ReflowPlan, plan_reflow
from presentation_editing_intelligence.text_edit import RunSpan, plan_single_run_replacement

_REAL_DECK_PATH = (
    Path(__file__).resolve().parents[2]
    / "documentation-artifacts"
    / "masterslide"
    / "old"
    / "MASTER Complete Bridge4PS Desktop-Browser Feature Tutorials.pptx"
)


def _shape(index, shape_id, left, top, width, height, is_line=False, text=None):
    return SceneShape(
        shape_index=index, shape_id=shape_id, shape_type="TEXT_BOX (17)" if not is_line else "LINE (9)",
        name="s%d" % shape_id, geometry=Geometry(left, top, width, height), z_order=index,
        text=text, line_color_rgb=None, line_width_emu=None, fill_color_rgb=None,
        is_line_shape=is_line, is_freeform=False, group_shape_index=None, connector_endpoints=None,
    )


def _build_simple_deck(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(2_000_000), Emu(400_000))
    a.text_frame.paragraphs[0].text = "Bold label: "
    a.text_frame.paragraphs[0].runs[0].font.bold = True
    a.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    run2 = a.text_frame.paragraphs[0].add_run()
    run2.text = "plain body text"
    run2.font.size = Pt(12)
    b = slide.shapes.add_textbox(Emu(500_000), Emu(1_500_000), Emu(2_000_000), Emu(400_000))
    b.text_frame.text = "Unrelated shape - must never change"
    prs.save(str(path))
    return str(path), a.shape_id, b.shape_id


# --------------------------------------------------------------------------
# PEI-S1-02: unauthorized post-plan mutation -> no published artifact,
# never "applied"
# --------------------------------------------------------------------------

def test_unauthorized_mutation_after_publish_is_rejected_and_output_removed(tmp_path, monkeypatch):
    """Simulates the engine somehow publishing a result that touched an
    unrelated shape (a bug in a lower layer, or a mid-flight tamper) by
    monkeypatching the mutation call to also corrupt the unrelated shape
    after the real engine call succeeds. Protected-structure verification
    must catch it, delete the output, and the pilot must never report
    'applied'."""
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    import presentation_editing_intelligence.pilot as pilot_module

    real_apply = pilot_module.apply_scene_aware_text_edit

    def tampering_apply(*args, **kwargs):
        result = real_apply(*args, **kwargs)
        # Tamper with the already-published output directly - exactly
        # the class of bug protected-structure verification exists to
        # catch (some lower layer changed something outside the plan).
        prs = Presentation(result.output_path)
        other_shape = list(prs.slides[0].shapes)[1]
        other_shape.text_frame.text = "TAMPERED - this must be caught"
        prs.save(result.output_path)
        return result

    monkeypatch.setattr(pilot_module, "apply_scene_aware_text_edit", tampering_apply)

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Bold label: plain body text, extended",
        shape_id=target_id,
    )

    assert result.status != "applied"
    assert result.protected_structure_ok is False
    assert not os.path.exists(output)


def test_unauthorized_formatting_change_on_untouched_run_is_rejected(tmp_path, monkeypatch):
    """Even a change that keeps every run's *text* the same but alters an
    untouched run's formatting must be rejected."""
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    import presentation_editing_intelligence.pilot as pilot_module
    real_apply = pilot_module.apply_scene_aware_text_edit

    def tampering_apply(*args, **kwargs):
        result = real_apply(*args, **kwargs)
        prs = Presentation(result.output_path)
        target_shape = list(prs.slides[0].shapes)[0]
        # Run 0 ("Bold label: ") was never part of the approved edit -
        # flipping its bold flag off must be caught.
        target_shape.text_frame.paragraphs[0].runs[0].font.bold = False
        prs.save(result.output_path)
        return result

    monkeypatch.setattr(pilot_module, "apply_scene_aware_text_edit", tampering_apply)

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Bold label: plain body text, extended",
        shape_id=target_id,
    )

    assert result.status != "applied"
    assert not os.path.exists(output)


# --------------------------------------------------------------------------
# PEI-S1-03: ambiguous zero-width run insertion
# --------------------------------------------------------------------------

def test_adjacent_identical_formatting_insertion_resolves():
    sig = (True, False, False, Pt(12), "Calibri", None)
    runs = [
        RunSpan(0, 0, "A", format_signature=sig),
        RunSpan(0, 1, "B", format_signature=sig),
    ]
    plan = plan_single_run_replacement(runs, "AB", "AXB")
    assert plan.status == "resolved"


def test_adjacent_different_formatting_insertion_is_unresolved():
    sig_a = (True, False, False, Pt(12), "Calibri", None)
    sig_b = (False, False, False, Pt(12), "Calibri", None)
    runs = [
        RunSpan(0, 0, "A", format_signature=sig_a),
        RunSpan(0, 1, "B", format_signature=sig_b),
    ]
    plan = plan_single_run_replacement(runs, "AB", "AXB")
    assert plan.status == "unresolved"
    assert "formatting" in plan.reason


def test_adjacent_unknown_formatting_insertion_is_unresolved():
    runs = [RunSpan(0, 0, "A"), RunSpan(0, 1, "B")]  # format_signature defaults to None
    plan = plan_single_run_replacement(runs, "AB", "AXB")
    assert plan.status == "unresolved"


def test_insertion_at_start_or_end_with_single_source_is_still_deterministic():
    sig = (False, False, False, Pt(12), "Calibri", None)
    runs = [RunSpan(0, 0, "World", format_signature=sig)]
    plan = plan_single_run_replacement(runs, "World", "Hello World")
    assert plan.status == "resolved"
    assert plan.edits[0].new_text == "Hello World"


# --------------------------------------------------------------------------
# PEI-S1-04: relationship confidence calibration
# --------------------------------------------------------------------------

def test_geometry_only_evidence_never_produces_high_confidence(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(2_000_000), Emu(400_000))
    a.text_frame.text = "Entry A"
    divider = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(500_000), Emu(950_000), Emu(2_500_000), Emu(950_000))
    divider.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
    b = slide.shapes.add_textbox(Emu(500_000), Emu(1_050_000), Emu(2_000_000), Emu(400_000))
    b.text_frame.text = "Entry B"

    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0)
    rels = infer_relationships(scene)

    geometry_only_types = {"separates", "anchored-to", "depends-on", "highlights", "aligned-with"}
    for r in rels:
        if r.relation_type in geometry_only_types:
            assert r.confidence != "HIGH", "geometry-only relationship %r must never be HIGH" % (r,)


def test_explicit_connector_binding_may_still_be_high(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(1_000_000), Emu(300_000))
    a.text_frame.text = "Source"
    b = slide.shapes.add_textbox(Emu(3_000_000), Emu(500_000), Emu(1_000_000), Emu(300_000))
    b.text_frame.text = "Target"
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(1_500_000), Emu(650_000), Emu(3_000_000), Emu(650_000))
    connector.begin_connect(a, 0)
    connector.end_connect(b, 0)

    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0, proximity_margin_emu=3_000_000)
    rels = infer_relationships(scene)
    points_to = [r for r in rels if r.relation_type == "points-to"]
    assert any(r.confidence == "HIGH" for r in points_to)


# --------------------------------------------------------------------------
# PEI-S1-05: group-support truthfulness (deferred capability)
# --------------------------------------------------------------------------

def test_group_shape_index_is_never_populated(tmp_path):
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    prs = Presentation(source)
    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0)
    for s in scene.shapes:
        assert s.group_shape_index is None


def test_moves_with_relationship_type_is_never_emitted(tmp_path):
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    prs = Presentation(source)
    scene = extract_local_scene(prs, slide_index=0, target_shape_index=0)
    rels = infer_relationships(scene)
    assert not any(r.relation_type == "moves-with" for r in rels)


# --------------------------------------------------------------------------
# PEI-S1-06: collision safety beyond the initial local scene
# --------------------------------------------------------------------------

def test_collision_with_a_shape_outside_local_scene_is_detected():
    # A target that grows, a divider directly below it that reflow will
    # move, and a THIRD shape far enough away to be excluded from the
    # local reasoning scene but sitting exactly where the divider's
    # movement would land it.
    target = _shape(0, 100, left=0, top=0, width=1_000_000, height=500_000, text="Target")
    divider = _shape(1, 101, left=0, top=550_000, width=1_000_000, height=0, is_line=True)
    scene = LocalScene(slide_index=0, target_shape_index=0, shapes=[target, divider])
    relationships = [
        Relationship(source_shape_id=101, target_shape_id=100, relation_type="depends-on",
                     confidence="MEDIUM", evidence="test"),
    ]
    grown = Geometry(target.geometry.left, target.geometry.top, target.geometry.width,
                      target.geometry.height + 2_000_000)  # divider will move 2,000,000 EMU down
    plan = plan_reflow(scene, relationships, grown)
    assert 101 in {m.shape_id for m in plan.moves}

    # This shape was never within the default local-scene proximity
    # margin of the target, so it is NOT in `scene.shapes` - only in the
    # full-slide collision candidate set.
    far_away_shape = _shape(5, 500, left=0, top=550_000 + 2_000_000, width=1_000_000, height=500_000,
                             text="Far away, outside local scene")
    assert far_away_shape.shape_id not in {s.shape_id for s in scene.shapes}

    report_without_full_slide_index = check_edit(
        scene, relationships, grown, plan, slide_width_emu=10_000_000, slide_height_emu=10_000_000,
    )
    assert report_without_full_slide_index.ok  # proves the bug is real: no candidate, no detection

    report_with_full_slide_index = check_edit(
        scene, relationships, grown, plan, slide_width_emu=10_000_000, slide_height_emu=10_000_000,
        collision_candidates=[target, divider, far_away_shape],
    )
    assert not report_with_full_slide_index.ok
    assert any(i.check == "swept_path_collision" for i in report_with_full_slide_index.critical)


def test_extract_all_shapes_on_slide_is_unfiltered(tmp_path):
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    prs = Presentation(source)
    all_shapes = extract_all_shapes_on_slide(prs, 0)
    assert {s.shape_id for s in all_shapes} == {target_id, other_id}


def test_pilot_wires_full_slide_collision_candidates_by_default(tmp_path):
    """End-to-end: the pilot itself must pass the full-slide index into
    `check_edit`, not just the local scene - proven by patching
    `check_edit` and inspecting what it was actually called with."""
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    import presentation_editing_intelligence.pilot as pilot_module
    captured = {}
    real_check_edit = pilot_module.check_edit

    def spying_check_edit(*args, **kwargs):
        captured["collision_candidates"] = kwargs.get("collision_candidates")
        return real_check_edit(*args, **kwargs)

    pilot_module.check_edit = spying_check_edit
    try:
        apply_scene_aware_edit(
            source, output, slide_index=0, shape_index=0,
            expected_old_text="Bold label: plain body text",
            new_text="Bold label: plain body text, extended",
            shape_id=target_id,
        )
    finally:
        pilot_module.check_edit = real_check_edit

    assert captured["collision_candidates"] is not None
    assert len(captured["collision_candidates"]) == 2  # both shapes on the slide


# --------------------------------------------------------------------------
# PEI-S1-02.A (re-audit): edited-run formatting is protected too
# --------------------------------------------------------------------------

def test_authorized_text_change_plus_unauthorized_formatting_on_edited_run_is_rejected(tmp_path, monkeypatch):
    """The edited run's TEXT changes exactly as approved, but its bold
    flag is also (unauthorized) flipped - Slice 1 has no approved
    formatting mutation, so this must be caught even though the run in
    question is the very one that was approved to change."""
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    import presentation_editing_intelligence.pilot as pilot_module
    real_apply = pilot_module.apply_scene_aware_text_edit

    def tampering_apply(*args, **kwargs):
        result = real_apply(*args, **kwargs)
        prs = Presentation(result.output_path)
        target_shape = list(prs.slides[0].shapes)[0]
        edited_run = target_shape.text_frame.paragraphs[0].runs[1]  # the run named in run_edits
        assert edited_run.text == "plain body text, extended"  # text change did apply as approved
        edited_run.font.bold = True  # unauthorized formatting change
        prs.save(result.output_path)
        return result

    monkeypatch.setattr(pilot_module, "apply_scene_aware_text_edit", tampering_apply)

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Bold label: plain body text, extended",
        shape_id=target_id,
    )

    assert result.status != "applied"
    assert result.protected_structure_ok is False
    assert "formatting" in result.reason
    assert not os.path.exists(output)


# --------------------------------------------------------------------------
# PEI-S1-02.B (re-audit): cleanup failure must never be silently hidden
# --------------------------------------------------------------------------

def test_cleanup_failure_after_verification_failure_is_reported_not_hidden(tmp_path, monkeypatch):
    """Fault-injects an `os.remove` failure at the exact moment the pilot
    tries to delete an unsafe, already-published output. Both the
    original verification failure AND the cleanup failure must be
    visible in the result - and the result must never claim the unsafe
    artifact was removed when it was not."""
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    import presentation_editing_intelligence.pilot as pilot_module
    real_apply = pilot_module.apply_scene_aware_text_edit

    def tampering_apply(*args, **kwargs):
        result = real_apply(*args, **kwargs)
        prs = Presentation(result.output_path)
        other_shape = list(prs.slides[0].shapes)[1]
        other_shape.text_frame.text = "TAMPERED"
        prs.save(result.output_path)
        return result

    monkeypatch.setattr(pilot_module, "apply_scene_aware_text_edit", tampering_apply)

    real_remove = os.remove

    def failing_remove(path):
        if path == output:
            raise OSError("simulated permission failure for this test")
        return real_remove(path)

    monkeypatch.setattr(pilot_module.os, "remove", failing_remove)

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Bold label: plain body text, extended",
        shape_id=target_id,
    )

    assert result.status != "applied"
    assert result.protected_structure_ok is False
    # The unsafe artifact genuinely was NOT removed - proven both by the
    # result's own honest field and by the file actually still existing.
    assert result.unsafe_output_removed is False
    assert result.unsafe_output_cleanup_error is not None
    assert "simulated permission failure" in result.unsafe_output_cleanup_error
    assert os.path.exists(output)  # still there - never falsely claimed removed
    assert "verification failed" in result.reason
    assert "cleanup" in result.reason.lower()


def test_cleanup_succeeds_normally_when_not_fault_injected(tmp_path, monkeypatch):
    """Sanity check: the normal (non-fault-injected) cleanup path still
    actually removes the file and reports removal truthfully."""
    source, target_id, other_id = _build_simple_deck(tmp_path / "source.pptx")
    output = str(tmp_path / "output.pptx")

    import presentation_editing_intelligence.pilot as pilot_module
    real_apply = pilot_module.apply_scene_aware_text_edit

    def tampering_apply(*args, **kwargs):
        result = real_apply(*args, **kwargs)
        prs = Presentation(result.output_path)
        other_shape = list(prs.slides[0].shapes)[1]
        other_shape.text_frame.text = "TAMPERED"
        prs.save(result.output_path)
        return result

    monkeypatch.setattr(pilot_module, "apply_scene_aware_text_edit", tampering_apply)

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Bold label: plain body text",
        new_text="Bold label: plain body text, extended",
        shape_id=target_id,
    )

    assert result.status != "applied"
    assert result.unsafe_output_removed is True
    assert result.unsafe_output_cleanup_error is None
    assert not os.path.exists(output)


# --------------------------------------------------------------------------
# PEI-S1-08 (re-audit 2): explicit refusal whenever fit in the EXISTING
# target container cannot be conservatively established - independent of
# PowerPoint `auto_size` mode. Codex's own reproduction: a fixed-height
# (`MSO_AUTO_SIZE.NONE`) box received a replacement estimated at 2->272
# lines, target height unchanged, and was still published with only a
# `possible_text_overflow` *warning* - this must now be a hard refusal.
# --------------------------------------------------------------------------

def _build_narrow_box_deck(tmp_path, auto_size, short_text="Short."):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(1_000_000), Emu(300_000))
    box.text_frame.word_wrap = True
    if auto_size is not None:
        box.text_frame.auto_size = auto_size
    box.text_frame.paragraphs[0].text = short_text
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    source = str(tmp_path / "source.pptx")
    prs.save(source)
    return source, box.shape_id


def test_shape_to_fit_textbox_requiring_growth_is_refused(tmp_path):
    """Requirement 2: SHAPE_TO_FIT_TEXT does not rescue an edit that
    exceeds the existing container - Slice 1 never resizes it."""
    from pptx.enum.text import MSO_AUTO_SIZE
    source, target_id = _build_narrow_box_deck(tmp_path, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    output = str(tmp_path / "output.pptx")
    long_replacement = "Short. " + ("This is a much longer replacement sentence. " * 8)

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Short.", new_text=long_replacement, shape_id=target_id,
    )

    assert result.status != "applied"
    assert "target resizing is not implemented" in result.reason
    assert not os.path.exists(output)
    assert result.mutation_result is None  # no mutation was ever attempted


def test_fixed_height_none_autosize_textbox_requiring_growth_is_refused(tmp_path):
    """Requirement 1, and Codex's exact reported defect: a fixed-height
    box (`MSO_AUTO_SIZE.NONE` - no autofit at all, PowerPoint will not
    rescue this by growing the box) receiving a massively longer
    replacement must be refused, not published with only a warning."""
    from pptx.enum.text import MSO_AUTO_SIZE
    source, target_id = _build_narrow_box_deck(tmp_path, MSO_AUTO_SIZE.NONE)
    output = str(tmp_path / "output.pptx")
    # Mirrors Codex's own reproduction shape (~2 lines -> ~272 lines).
    huge_replacement = "Short. " + ("This is a much longer replacement sentence. " * 60)

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Short.", new_text=huge_replacement, shape_id=target_id,
    )

    assert result.status != "applied"
    assert "target resizing is not implemented" in result.reason
    assert not os.path.exists(output)
    assert result.mutation_result is None
    assert result.safety_report is None  # refused before any safety check ran


def test_third_autosize_value_still_refuses_conservatively_on_overflow(tmp_path):
    """Requirement 4: an `auto_size` value this slice does not
    special-case anywhere - `MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE` ("shrink
    text on overflow"), a genuinely different third value from `NONE` and
    `SHAPE_TO_FIT_TEXT` - must still refuse an edit that cannot be
    conservatively established to fit. The Slice 1 rule applies
    regardless of `auto_size`, known or not, because the check never
    branches on `auto_size` at all - covered here with an explicitly-set,
    real third enum value, not merely by skipping the assignment.

    (Corrected per Codex PEI-S1-07: skipping the `auto_size` assignment
    on a fresh python-pptx textbox does NOT produce an "unset" value -
    `add_textbox`'s own real default is `SHAPE_TO_FIT_TEXT` itself,
    confirmed directly against python-pptx; a prior version of this test
    incorrectly claimed the opposite and was, in substance, a duplicate
    of the `SHAPE_TO_FIT_TEXT` case above, not a distinct probe.)"""
    from pptx.enum.text import MSO_AUTO_SIZE
    source, target_id = _build_narrow_box_deck(tmp_path, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
    output = str(tmp_path / "output.pptx")
    huge_replacement = "Short. " + ("This is a much longer replacement sentence. " * 60)

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Short.", new_text=huge_replacement, shape_id=target_id,
    )

    assert result.status != "applied"
    assert "target resizing is not implemented" in result.reason
    assert not os.path.exists(output)


def test_replacement_that_clearly_fits_existing_container_is_allowed(tmp_path):
    """Requirement 3: an edit that plainly does not need more space than
    the existing container already provides must proceed normally, for
    every `auto_size` value - proven here for all three real, explicitly-
    set enum values (`NONE`, `SHAPE_TO_FIT_TEXT`, `TEXT_TO_FIT_SHAPE`),
    each set directly rather than by relying on skipping the assignment."""
    from pptx.enum.text import MSO_AUTO_SIZE
    for auto_size in (MSO_AUTO_SIZE.NONE, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE):
        source, target_id = _build_narrow_box_deck(tmp_path, auto_size, short_text="Short.")
        output = str(tmp_path / ("output_%d.pptx" % int(auto_size)))

        result = apply_scene_aware_edit(
            source, output, slide_index=0, shape_index=0,
            expected_old_text="Short.", new_text="Short!", shape_id=target_id,
        )

        assert result.status == "applied", "auto_size=%r: %r" % (auto_size, result.reason)
        assert os.path.exists(output)


# --------------------------------------------------------------------------
# PEI-S1-08 (final repair): conservative mixed-format fit estimation
# --------------------------------------------------------------------------

def _build_mixed_format_deck(tmp_path, first_run_size_pt=8, edited_run_size_pt=40,
                              geometry=(Emu(500_000), Emu(500_000), Emu(2_000_000), Emu(400_000)),
                              edited_run_has_size=True):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(*geometry)
    box.text_frame.word_wrap = True
    para = box.text_frame.paragraphs[0]
    para.text = "Label: "
    para.runs[0].font.size = Pt(first_run_size_pt)
    edited_run = para.add_run()
    edited_run.text = "value"
    if edited_run_has_size:
        edited_run.font.size = Pt(edited_run_size_pt)
    source = str(tmp_path / "source.pptx")
    prs.save(source)
    return source, box.shape_id


def test_mixed_font_8pt_first_run_40pt_edited_run_refuses_on_long_replacement(tmp_path):
    """Codex's exact reproduction: an 8pt first run must not cause the
    much larger 40pt edited run's actual space requirement to be
    underestimated. A long replacement, sized for the true (40pt)
    conservative estimate, must be refused - not silently sized against
    the smaller 8pt run."""
    source, target_id = _build_mixed_format_deck(tmp_path, first_run_size_pt=8, edited_run_size_pt=40)
    output = str(tmp_path / "output.pptx")
    long_replacement = "value " * 30  # long enough at 40pt to overflow a small box, not at 8pt

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Label: value", new_text="Label: " + long_replacement, shape_id=target_id,
    )

    assert result.status != "applied"
    assert "target resizing is not implemented" in result.reason
    assert not os.path.exists(output)


def test_mixed_fonts_where_conservative_calculation_proves_fit_is_allowed(tmp_path):
    """Mixed run sizes, but a box generously sized relative to even the
    largest (conservative) font and a trivial replacement - conservative
    calculation must still allow an edit that plainly fits."""
    source, target_id = _build_mixed_format_deck(
        tmp_path, first_run_size_pt=8, edited_run_size_pt=14,
        geometry=(Emu(500_000), Emu(500_000), Emu(4_000_000), Emu(2_000_000)),
    )
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Label: value", new_text="Label: item", shape_id=target_id,
    )

    assert result.status == "applied", result.reason
    assert os.path.exists(output)


def test_unresolvable_edited_run_font_size_refuses(tmp_path):
    """The edited run has no run-level size AND the paragraph has no
    default size set either - an effective size cannot be resolved at
    either deterministic level, so this must refuse outright, not
    silently assume a default (e.g. 12pt) and estimate against it."""
    source, target_id = _build_mixed_format_deck(
        tmp_path, first_run_size_pt=12, edited_run_has_size=False,
    )
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Label: value", new_text="Label: item", shape_id=target_id,
    )

    assert result.status != "applied"
    assert "cannot conservatively resolve" in result.reason
    assert not os.path.exists(output)


def test_paragraph_spacing_that_materially_increases_required_height_causes_refusal(tmp_path):
    """A paragraph with an explicit `space_after` set. The edit itself
    grows the text from 1 wrapped line to 2 - a growth that, by line-wrap
    height alone, would still comfortably fit this (deliberately roomier)
    box. Only once the paragraph's own `space_after` (a fixed per-
    paragraph overhead, deterministically available, now accounted for)
    is added on top does the true required height exceed the container -
    proving spacing is genuinely factored into the decision, not that
    line-count growth alone was already going to refuse regardless of it.
    The `before` state (1 line + spacing) still comfortably fits, so the
    empirical floor does not neutralize this case."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Roomy enough that 2 plain-text lines fit fine on their own (no
    # spacing) - only the added space_after pushes the 2-line case over.
    box = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(2_000_000), Emu(600_000))
    box.text_frame.word_wrap = True
    para = box.text_frame.paragraphs[0]
    para.text = "Short."
    para.runs[0].font.size = Pt(12)
    para.space_after = Emu(200_000)
    source = str(tmp_path / "source.pptx")
    prs.save(source)
    target_id = box.shape_id
    output = str(tmp_path / "output.pptx")

    # Long enough to wrap to a second line in this box's width - the
    # combination of 2 lines' height plus the paragraph's own space_after
    # is what exceeds the container; 2 lines alone would not have.
    two_line_replacement = "This sentence needs two full lines definitely"

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Short.", new_text=two_line_replacement, shape_id=target_id,
    )

    assert result.status != "applied"
    assert "target resizing is not implemented" in result.reason
    assert not os.path.exists(output)


def test_real_benchmark_slide_still_applies_with_conservative_formatting(tmp_path):
    """Requirement 5: the existing real benchmark (uniform 12pt runs on
    slide 74 shape 2428) must still apply under the new conservative,
    mixed-format-aware estimation - not only under the old single-run
    estimate."""
    if not _REAL_DECK_PATH.exists():
        import pytest as _pytest
        _pytest.skip("real MasterSlide deck not present in working tree")
    source = str(tmp_path / "source.pptx")
    import shutil as _shutil
    _shutil.copyfile(str(_REAL_DECK_PATH), source)
    output = str(tmp_path / "output.pptx")
    old_text = (
        "Pin:\nMessages bookmarked for all members of the channel to "
        "view in the Pinned Messages list."
    )
    new_text = old_text + " (approved test update)"

    result = apply_scene_aware_edit(
        source, output, slide_index=74, shape_index=12,
        expected_old_text=old_text, new_text=new_text, shape_id=2428,
    )

    assert result.status == "applied", result.reason
    assert result.protected_structure_ok is True


# PEI-S1-08 (second repair): absolute Length line-spacing type handling.
#
# python-pptx returns an absolute `line_spacing` (e.g. after `paragraph.
# line_spacing = Pt(12)`) as a `Centipoints`/`Emu` instance - every such
# `Length` subclass is *itself* an `int` subclass. Checking
# `isinstance(line_spacing, (int, float))` before checking `isinstance(
# line_spacing, Length)` therefore silently misclassifies every absolute
# value as a multiplier, multiplying the EMU line height by e.g. 152400
# instead of using 152400 EMU directly - producing an astronomically
# large, always-refusing estimate regardless of how well the text
# actually fits.

def _one_paragraph_deck(tmp_path, box_h, box_w, text, font_size_pt=12, line_spacing=None):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(box_w), Emu(box_h))
    box.text_frame.word_wrap = True
    para = box.text_frame.paragraphs[0]
    para.text = text
    para.runs[0].font.size = Pt(font_size_pt)
    if line_spacing is not None:
        para.line_spacing = line_spacing
    source = str(tmp_path / "source.pptx")
    prs.save(source)
    return source, box.shape_id


def test_absolute_line_spacing_that_clearly_fits_is_allowed_not_treated_as_multiplier(tmp_path):
    """`Pt(12)` absolute line spacing (~152,400 EMU) on a short single
    line that obviously fits a 300,000 EMU box. Under the pre-repair bug
    this was multiplied against the base line height instead of used
    directly, producing a ~27.8 billion EMU estimate and a false refusal
    no matter how roomy the box was. Must apply."""
    source, target_id = _one_paragraph_deck(tmp_path, box_h=300_000, box_w=2_000_000, text="Hi", line_spacing=Pt(12))
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Hi", new_text="Ho", shape_id=target_id,
    )

    assert result.status == "applied", result.reason


def test_absolute_line_spacing_that_genuinely_overflows_is_refused(tmp_path):
    """`Pt(30)` absolute per-line spacing (~381,000 EMU/line). The
    single-line `before` state comfortably fits its container (so the
    empirical floor does not neutralize the case); the edit grows the
    text to 2 wrapped lines, and 2 lines at the real absolute per-line
    height genuinely exceeds the container. Must refuse."""
    source, target_id = _one_paragraph_deck(
        tmp_path, box_h=500_000, box_w=2_000_000, text="Hi", line_spacing=Pt(30),
    )
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Hi",
        new_text="This sentence needs two lines for sure definitely",
        shape_id=target_id,
    )

    assert result.status != "applied"
    assert "target resizing is not implemented" in result.reason
    assert not os.path.exists(output)


def test_numeric_multiplier_line_spacing_still_behaves_as_a_multiplier(tmp_path):
    """A plain float `line_spacing = 2.0` (not a `Length`) must still be
    treated as a multiple of the base line height, not misread as an
    absolute EMU value (which would collapse to ~2 EMU/line and always
    fit). The single-line `before` state fits; growing to 2 lines at
    double line height genuinely overflows the container - a result only
    possible if the multiplier is actually still being applied."""
    source, target_id = _one_paragraph_deck(
        tmp_path, box_h=400_000, box_w=2_000_000, text="Hi", line_spacing=2.0,
    )
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Hi",
        new_text="This sentence needs two lines for sure definitely",
        shape_id=target_id,
    )

    assert result.status != "applied"
    assert "target resizing is not implemented" in result.reason
    assert not os.path.exists(output)


def test_space_before_still_contributes_to_required_height_after_the_line_spacing_fix(tmp_path):
    """Regression guard: the `Length`-vs-multiplier branch reorder in
    `_paragraph_conservative_line_height_emu` must not disturb
    `_paragraph_spacing_emu`'s separate, unrelated handling of
    `space_before`/`space_after`. Mirrors the existing `space_after`
    adversarial case but exercises `space_before` instead, to prove both
    are still independently accounted for."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(500_000), Emu(500_000), Emu(2_000_000), Emu(600_000))
    box.text_frame.word_wrap = True
    para = box.text_frame.paragraphs[0]
    para.text = "Short."
    para.runs[0].font.size = Pt(12)
    para.space_before = Emu(200_000)
    source = str(tmp_path / "source.pptx")
    prs.save(source)
    target_id = box.shape_id
    output = str(tmp_path / "output.pptx")

    result = apply_scene_aware_edit(
        source, output, slide_index=0, shape_index=0,
        expected_old_text="Short.",
        new_text="This sentence needs two full lines definitely",
        shape_id=target_id,
    )

    assert result.status != "applied"
    assert "target resizing is not implemented" in result.reason
    assert not os.path.exists(output)
