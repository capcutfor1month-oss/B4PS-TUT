"""Scene-aware single-slide text-edit pilot: the one public entry point
this slice provides, composing every other module in this package.

Given an already-approved `(old_text, new_text)` edit for one exact
target shape, this:

1. extracts the local scene (`scene.py`);
2. infers relationships in it (`relationships.py`);
3. plans a formatting-preserving run edit (`text_edit.py`) - stopping
   with `status="unresolved"` rather than guessing if the edit cannot be
   safely localized to one run;
4. resolves a conservative effective font size across every run in the
   target text (`safety.resolve_conservative_font_size_pt` - the maximum
   resolvable size, never just the first run's; refuses if any run's size
   cannot be resolved at all), then estimates the replacement's required
   height, accounting for each paragraph's own line spacing and before/
   after spacing where deterministically available
   (`safety.estimate_required_height_emu`), and checks it against the
   target's **existing** container height - independent of PowerPoint
   `auto_size` mode, since Slice 1 never resizes the target container in
   any code path, regardless of what `auto_size` says (Codex PEI-S1-08).
   If fit cannot be conservatively established, this returns `status=
   "unresolved"` right here - before planning any reflow, before any
   safety check, before any mutation or publication. Otherwise, plans a
   bounded local reflow for shapes evidenced as depending on the target
   (`reflow.py`) - the target's own geometry is always unchanged, so this
   plan is a no-op unless a future slice adds real target resizing;
5. runs deterministic pre/post safety checks (`safety.py`), including
   swept-path collision detection (start bbox through final bbox, not
   just final-vs-final - PEI-S1-06) against every shape on the target
   slide, not only the local reasoning scene - stopping before
   publishing on any critical issue;
6. applies the edit and any reflow as one atomic output, entirely
   through the Safe PPT Engine's own audited transaction boundary,
   including the text target's own identity guard validated on the
   staged copy inside that same transaction (`text_mutation.py` /
   PEI-S1-01);
7. runs comprehensive protected-structure verification
   (`protected_structure.py` / PEI-S1-02) - geometry, identity, ordering,
   type, text, and run-level formatting for every shape outside the
   approved plan, INCLUDING the edited run(s) themselves (Slice 1 has no
   approved formatting mutation - only text may change on an edited
   run). If this fails, the just-published output is removed and
   `"applied"` is never reported; if that removal itself fails, both
   failures are preserved and reported, never silently discarded.

This module never decides whether `new_text` is correct wording - exactly
like Slice 3's `apply_approved_replacement`, it assumes the replacement
text is already human-approved before this is ever called.
"""

from __future__ import annotations

import dataclasses
import os
from typing import List, Optional

from pptx import Presentation

from ._safe_ppt_engine_import import MutationError
from . import protected_structure
from .reflow import ReflowPlan, plan_reflow
from .relationships import Relationship, infer_relationships
from .safety import (
    SafetyReport,
    check_edit,
    estimate_line_count,
    estimate_required_height_emu,
    resolve_conservative_font_size_pt,
)
from .scene import LocalScene, extract_all_shapes_on_slide, extract_local_scene
from .text_edit import EditPlan, RunSpan, plan_single_run_replacement
from .text_mutation import MutationResult, apply_scene_aware_text_edit


@dataclasses.dataclass
class PilotResult:
    status: str  # "applied" | "unresolved" | "unsafe"
    reason: Optional[str] = None
    scene: Optional[LocalScene] = None
    relationships: List[Relationship] = dataclasses.field(default_factory=list)
    edit_plan: Optional[EditPlan] = None
    reflow_plan: Optional[ReflowPlan] = None
    safety_report: Optional[SafetyReport] = None
    mutation_result: Optional[MutationResult] = None
    protected_structure_ok: Optional[bool] = None
    unsafe_output_removed: Optional[bool] = None
    unsafe_output_cleanup_error: Optional[str] = None


def _run_format_signature(run) -> tuple:
    font = run.font
    color = None
    try:
        if font.color is not None and font.color.type is not None:
            color = (str(font.color.type), str(font.color.rgb) if font.color.type == 1 else str(font.color.theme_color))
    except Exception:
        color = None
    return (font.bold, font.italic, font.underline, font.size, font.name, color)


def _run_spans_for(shape) -> List[RunSpan]:
    spans: List[RunSpan] = []
    for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
        for r_idx, run in enumerate(paragraph.runs):
            spans.append(RunSpan(
                paragraph_index=p_idx, run_index=r_idx, text=run.text,
                format_signature=_run_format_signature(run),
            ))
    return spans


def _file_identity(path: str):
    try:
        st = os.stat(path)
    except OSError:
        return None
    return (st.st_dev, st.st_ino)


@dataclasses.dataclass(frozen=True)
class _CleanupOutcome:
    attempted: bool
    removed: bool
    detail: Optional[str] = None


def _remove_if_still_the_published_file(path: str, published_identity) -> _CleanupOutcome:
    """Removes `path` only if it is still, by POSIX device+inode
    identity, the exact file this run just published - the same
    ownership-race-safe pattern Documentation Intelligence Slice 3's own
    cleanup uses. If something else has since replaced it, it is left
    completely untouched, and that is reported (`attempted=False`), not
    conflated with a successful removal.

    Codex finding PEI-S1-02.B (re-audit): a cleanup failure here must
    never be silently discarded. This function always returns a
    `_CleanupOutcome` - `removed=True` only when `os.remove` actually
    succeeded; `removed=False` with `detail` set whenever cleanup was
    attempted but failed, so the caller can report both the original
    verification failure AND the fact that the unsafe artifact may still
    be on disk requiring manual cleanup - never silently claim it was
    removed when it was not."""
    current_identity = _file_identity(path)
    if current_identity is None or current_identity != published_identity:
        return _CleanupOutcome(attempted=False, removed=False)
    try:
        os.remove(path)
    except OSError as exc:
        return _CleanupOutcome(
            attempted=True, removed=False,
            detail="failed to remove unsafe output %r after protected-structure verification "
                   "failed (%s) - this artifact was NOT actually removed and may require "
                   "manual cleanup" % (path, exc),
        )
    return _CleanupOutcome(attempted=True, removed=True)


def apply_scene_aware_edit(
    source_deck_path: str,
    output_deck_path: str,
    slide_index: int,
    shape_index: int,
    expected_old_text: str,
    new_text: str,
    shape_id: Optional[int] = None,
    overwrite: bool = False,
) -> PilotResult:
    prs = Presentation(source_deck_path)
    scene = extract_local_scene(prs, slide_index, shape_index)
    target_scene_shape = scene.target

    if shape_id is not None and target_scene_shape.shape_id != shape_id:
        return PilotResult(
            status="unresolved",
            reason="shape_id guard mismatch: expected %r, found %r" % (shape_id, target_scene_shape.shape_id),
            scene=scene,
        )

    relationships = infer_relationships(scene)

    slide = prs.slides[slide_index]
    target_shape = list(slide.shapes)[shape_index]
    run_spans = _run_spans_for(target_shape)

    edit_plan = plan_single_run_replacement(run_spans, expected_old_text, new_text)
    if edit_plan.status != "resolved":
        return PilotResult(
            status="unresolved",
            reason=edit_plan.reason,
            scene=scene,
            relationships=relationships,
            edit_plan=edit_plan,
        )

    tf = target_shape.text_frame
    paragraphs = list(tf.paragraphs)  # stable count/formatting before and after - see estimate_required_height_emu

    # Codex PEI-S1-08 (final repair): fit estimation must use a
    # conservative *effective* font size across every run participating
    # in the target text, not whichever run happens to be first with an
    # explicit size - Codex's own reproduction used an 8pt first run and
    # a 40pt edited run, and the prior estimate silently used 8pt.
    # "uncertain fit => refuse" applies to formatting resolution itself:
    # if any run's effective size cannot be resolved deterministically
    # (no run-level size and no paragraph-level default), this refuses
    # outright rather than guessing a font size to estimate with.
    conservative_font_size_pt = resolve_conservative_font_size_pt(tf)
    if conservative_font_size_pt is None:
        return PilotResult(
            status="unresolved",
            reason=(
                "cannot conservatively resolve the effective font size for one or more runs "
                "in the target shape (no run-level or paragraph-level size available) - "
                "refusing rather than estimating fit with an unresolved size"
            ),
            scene=scene,
            relationships=relationships,
            edit_plan=edit_plan,
        )

    paragraph_texts_before = expected_old_text.split("\n")
    paragraph_texts_after = new_text.split("\n")
    if len(paragraph_texts_before) != len(paragraphs) or len(paragraph_texts_after) != len(paragraphs):
        # Slice 1 never adds/removes paragraphs (text_edit.py refuses any
        # edit that would) - a mismatch here means an assumption this
        # module depends on does not hold; refuse rather than estimate
        # against paragraph formatting that may not correspond to the text.
        return PilotResult(
            status="unresolved",
            reason=(
                "paragraph count in the requested text does not match the target shape's "
                "current paragraph structure - cannot conservatively estimate fit"
            ),
            scene=scene,
            relationships=relationships,
            edit_plan=edit_plan,
        )

    # The true Slice 1 rule (Codex PEI-S1-08): if fit in the EXISTING
    # target container cannot be conservatively established, refuse -
    # independent of PowerPoint `auto_size` mode (`NONE`,
    # `SHAPE_TO_FIT_TEXT`, or anything else/unknown are all treated
    # identically; this check never inspects `auto_size` at all) and now
    # accounting for mixed run formatting (the resolved conservative font
    # size above) and, where deterministically available, each
    # paragraph's own line spacing and before/after spacing
    # (`estimate_required_height_emu`). This slice never resizes the
    # target container in any code path - `target_new_geometry` is always
    # the target's own unchanged geometry.
    target_new_geometry = target_scene_shape.geometry
    margin_left = tf.margin_left or 0
    margin_right = tf.margin_right or 0
    margin_top = tf.margin_top or 0
    margin_bottom = tf.margin_bottom or 0

    required_height_before = estimate_required_height_emu(
        paragraphs, paragraph_texts_before, target_scene_shape.geometry, conservative_font_size_pt,
        margin_left_emu=margin_left, margin_right_emu=margin_right,
    )
    required_height_after = estimate_required_height_emu(
        paragraphs, paragraph_texts_after, target_scene_shape.geometry, conservative_font_size_pt,
        margin_left_emu=margin_left, margin_right_emu=margin_right,
    )
    usable_height_emu = target_scene_shape.geometry.height - margin_top - margin_bottom

    # `required_height_before` is an empirical floor on real capacity:
    # this exact container, at its exact existing height, already holds
    # the existing text today (it is real, currently-published content) -
    # regardless of what the heuristic alone would estimate. Without this
    # floor, the approximate heuristic can be stricter than reality and
    # falsely refuse a same-length or trivially-extended edit on a slide
    # that already renders correctly (confirmed against the real
    # benchmark, which needs exactly this). A genuinely large increase
    # (Codex's own 8pt/40pt mixed-format reproduction, or the earlier
    # 2->272-line reproduction) clears any reasonable floor regardless.
    safely_available_height = max(usable_height_emu, required_height_before)
    if required_height_after > safely_available_height:
        return PilotResult(
            status="unresolved",
            reason=(
                "replacement may exceed existing target container; target resizing is not "
                "implemented in Slice 1 (estimated %d EMU required, existing container "
                "conservatively fits approximately %d EMU)"
                % (required_height_after, safely_available_height)
            ),
            scene=scene,
            relationships=relationships,
            edit_plan=edit_plan,
        )

    lines_before = estimate_line_count(
        expected_old_text, target_scene_shape.geometry, conservative_font_size_pt,
        margin_left_emu=margin_left, margin_right_emu=margin_right,
    )
    lines_after = estimate_line_count(
        new_text, target_scene_shape.geometry, conservative_font_size_pt,
        margin_left_emu=margin_left, margin_right_emu=margin_right,
    )

    reflow_plan = plan_reflow(scene, relationships, target_new_geometry)

    # PEI-S1-06: collision detection uses every shape on the target
    # slide, not only the local reasoning scene.
    full_slide_shapes = extract_all_shapes_on_slide(prs, slide_index)

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    safety_report = check_edit(
        scene, relationships, target_new_geometry, reflow_plan,
        slide_width_emu=slide_width, slide_height_emu=slide_height,
        estimated_line_count_before=lines_before, estimated_line_count_after=lines_after,
        collision_candidates=full_slide_shapes,
    )

    if not safety_report.ok:
        return PilotResult(
            status="unsafe",
            reason="; ".join(i.detail for i in safety_report.critical),
            scene=scene,
            relationships=relationships,
            edit_plan=edit_plan,
            reflow_plan=reflow_plan,
            safety_report=safety_report,
        )

    try:
        mutation_result = apply_scene_aware_text_edit(
            source_deck_path,
            output_deck_path,
            slide_index,
            shape_index,
            edit_plan.edits,
            reflow_moves=reflow_plan.moves,
            expected_shape_id=shape_id if shape_id is not None else target_scene_shape.shape_id,
            overwrite=overwrite,
        )
    except MutationError as exc:
        return PilotResult(
            status="unresolved",
            reason=str(exc),
            scene=scene,
            relationships=relationships,
            edit_plan=edit_plan,
            reflow_plan=reflow_plan,
            safety_report=safety_report,
        )

    # PEI-S1-02: comprehensive protected-structure verification is the
    # final gate. A failure here removes the just-published output (if it
    # is still, by file identity, the exact file this call published) and
    # this function never reports "applied".
    published_identity = _file_identity(output_deck_path)
    structure_result = protected_structure.verify(
        source_deck_path, output_deck_path, slide_index, shape_index,
        edit_plan.edits, reflow_plan.moves, expected_target_full_text=new_text,
    )

    if not structure_result.ok:
        cleanup = _remove_if_still_the_published_file(output_deck_path, published_identity)
        # Codex PEI-S1-02.B (re-audit): the primary verification failure
        # and any secondary cleanup failure are both preserved explicitly
        # - never silently swallowed, never claimed removed if it was not.
        reason = "protected-structure verification failed: %s" % structure_result.reason
        if cleanup.attempted and not cleanup.removed:
            reason += (
                "; ADDITIONALLY, cleanup of the unsafe output failed (%s) - the unsafe "
                "artifact may still exist at %r and requires manual cleanup"
                % (cleanup.detail, output_deck_path)
            )
        return PilotResult(
            status="unsafe",
            reason=reason,
            scene=scene,
            relationships=relationships,
            edit_plan=edit_plan,
            reflow_plan=reflow_plan,
            safety_report=safety_report,
            mutation_result=mutation_result,
            protected_structure_ok=False,
            unsafe_output_removed=cleanup.removed if cleanup.attempted else None,
            unsafe_output_cleanup_error=cleanup.detail,
        )

    return PilotResult(
        status="applied",
        scene=scene,
        relationships=relationships,
        edit_plan=edit_plan,
        reflow_plan=reflow_plan,
        safety_report=safety_report,
        mutation_result=mutation_result,
        protected_structure_ok=True,
    )
