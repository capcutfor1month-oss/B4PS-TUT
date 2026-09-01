"""Deterministic pre/post safety checks.

Every check here is a plain geometric or structural comparison - nothing
is inferred from a rendered image, because none is rendered. These
checks catch what is deterministically detectable from OOXML facts
alone; they are not a substitute for native Microsoft PowerPoint visual
review, which remains the actual acceptance authority for how the slide
looks (see the governance spec's "native PowerPoint visual authority"
requirement).

A `critical` issue means: do not publish this edit. A `warning` is
recorded but does not by itself block the edit - it is reported so a
human reviewer sees it before deciding.
"""

from __future__ import annotations

import dataclasses
from typing import Dict, List, Optional

from pptx.util import Length

from .reflow import ReflowPlan
from .relationships import Relationship
from .scene import Geometry, LocalScene, SceneShape

# Very rough, explicitly-approximate average glyph width as a fraction of
# font size (em), used only to flag *possible* overflow for human review -
# never to silently resize a shape based on a guessed exact line count.
# Native PowerPoint's own text layout is the real authority.
_AVG_GLYPH_WIDTH_EM = 0.52
_LINE_HEIGHT_EM = 1.2
_EMU_PER_INCH = 914400
_EMU_PER_POINT = _EMU_PER_INCH / 72


@dataclasses.dataclass(frozen=True)
class SafetyIssue:
    severity: str  # "critical" | "warning"
    check: str
    detail: str


@dataclasses.dataclass(frozen=True)
class SafetyReport:
    issues: List[SafetyIssue]

    @property
    def ok(self) -> bool:
        return not any(i.severity == "critical" for i in self.issues)

    @property
    def critical(self) -> List[SafetyIssue]:
        return [i for i in self.issues if i.severity == "critical"]

    @property
    def warnings(self) -> List[SafetyIssue]:
        return [i for i in self.issues if i.severity == "warning"]


def _final_geometry(shape: SceneShape, target_shape_id: int, target_new: Geometry,
                     reflow: ReflowPlan) -> Geometry:
    if shape.shape_id == target_shape_id:
        return target_new
    for move in reflow.moves:
        if move.shape_id == shape.shape_id:
            g = shape.geometry
            return Geometry(left=g.left, top=g.top + move.delta_top_emu, width=g.width, height=g.height)
    return shape.geometry


def _swept_geometry(original: Geometry, final: Geometry) -> Geometry:
    """A conservative bounding rectangle covering a shape's entire
    movement from its original position to its final one (Codex
    PEI-S1-06) - not an animation path, just the union of the two
    bounding boxes. This slice's reflow only ever changes `top`, so this
    reduces to extending the box vertically to span both positions;
    `left`/`width` are taken from whichever box is wider, so the
    conservative rectangle never under-covers even if that ever changes.
    A shape that never moved has `original == final`, so its swept
    region is just its own unchanged bounding box - safe to use
    unconditionally for every shape, not only moved ones."""
    left = min(original.left, final.left)
    right = max(original.right, final.right)
    top = min(original.top, final.top)
    bottom = max(original.bottom, final.bottom)
    return Geometry(left=left, top=top, width=right - left, height=bottom - top)


def check_edit(
    scene: LocalScene,
    relationships: List[Relationship],
    target_new_geometry: Geometry,
    reflow: ReflowPlan,
    slide_width_emu: Optional[int] = None,
    slide_height_emu: Optional[int] = None,
    estimated_line_count_before: Optional[int] = None,
    estimated_line_count_after: Optional[int] = None,
    collision_candidates: Optional[List[SceneShape]] = None,
) -> SafetyReport:
    """`collision_candidates` (Codex PEI-S1-06): every shape on the target
    slide, not only the ones the local *reasoning* scene included -
    proximity-bounded scene extraction can exclude a shape that a planned
    movement's swept path would still reach. Relationship-derived checks
    (alignment breakage, relationship displacement) stay scoped to
    `scene`/`relationships`, since those need evidence this bounded
    reasoning scope actually produced; only the overlap/collision check
    below is widened to the full-slide candidate set. Omitting
    `collision_candidates` falls back to `scene.shapes` only (the
    pre-PEI-S1-06 behavior) - callers should always supply the full-slide
    index in production use (see `pilot.py`)."""
    issues: List[SafetyIssue] = []
    target = scene.target
    moved_ids = {target.shape_id} | {m.shape_id for m in reflow.moves}

    finals: Dict[int, Geometry] = {
        s.shape_id: _final_geometry(s, target.shape_id, target_new_geometry, reflow)
        for s in scene.shapes
    }

    # 1. Overlap/collision: any pair with no overlap before but overlap
    # after - checked over the full-slide candidate set (PEI-S1-06), not
    # only the local reasoning scene. A candidate outside the local scene
    # never moves (nothing in this slice's reflow can move a shape the
    # relationship graph never evidenced a dependency for), so its
    # "final" geometry is simply its own original geometry.
    shapes_by_id: Dict[int, SceneShape] = {s.shape_id: s for s in scene.shapes}
    for s in (collision_candidates or []):
        shapes_by_id.setdefault(s.shape_id, s)
        finals.setdefault(s.shape_id, s.geometry)

    ids = list(shapes_by_id.keys())

    # 1a. Final-position collision: any pair with no overlap before but
    # overlap in their final resting positions. Checked for every pair,
    # moved or not - this is what catches a reflow plan whose *end state*
    # is itself wrong (e.g. two dependents both moved into the same
    # place).
    for i in range(len(ids)):
        for j in range(i + 1, len(ids)):
            a_id, b_id = ids[i], ids[j]
            a, b = shapes_by_id[a_id], shapes_by_id[b_id]
            before_overlap = a.geometry.intersects(b.geometry, margin=0)
            after_overlap = finals[a_id].intersects(finals[b_id], margin=0)
            if after_overlap and not before_overlap:
                issues.append(SafetyIssue(
                    severity="critical",
                    check="overlap",
                    detail="shape id %d and shape id %d did not overlap before this edit "
                           "but would overlap after it" % (a_id, b_id),
                ))

    # 1b. Swept-path collision (Codex PEI-S1-06): for every MOVED shape,
    # check its swept region (start bbox unioned with final bbox) against
    # every STATIONARY candidate's own fixed geometry - not against other
    # moved shapes. Two shapes that are both part of this edit's own
    # coordinated reflow (e.g. a divider being pushed down by exactly the
    # amount the target grows) are *expected* to stay adjacent throughout
    # a synchronized move; comparing their independent swept envelopes
    # against each other produces a false collision even when a
    # synchronized transition never actually overlaps them (the checked-
    # in adversarial test for this is `test_reflow_and_safety.py`'s
    # correctly-reflowed-edit case). A moved shape passing through a
    # genuinely stationary, uninvolved shape and finishing beyond it is
    # exactly what this check exists to catch - that shape's position
    # never changes, so there is no coordination to account for.
    for moved_id in moved_ids:
        if moved_id not in shapes_by_id:
            continue
        moved_swept = _swept_geometry(shapes_by_id[moved_id].geometry, finals[moved_id])
        for other_id in ids:
            if other_id == moved_id or other_id in moved_ids:
                continue
            stationary = shapes_by_id[other_id]
            before_overlap = shapes_by_id[moved_id].geometry.intersects(stationary.geometry, margin=0)
            swept_overlap = moved_swept.intersects(stationary.geometry, margin=0)
            if swept_overlap and not before_overlap:
                issues.append(SafetyIssue(
                    severity="critical",
                    check="swept_path_collision",
                    detail="shape id %d's movement path (start bbox through final bbox) would "
                           "pass through stationary shape id %d, even though their final "
                           "positions alone might not collide" % (moved_id, other_id),
                ))

    # 2. Out-of-bounds (only checked when slide dimensions are supplied).
    if slide_width_emu is not None and slide_height_emu is not None:
        for shape_id in moved_ids:
            g = finals[shape_id]
            if g.left < 0 or g.top < 0 or g.right > slide_width_emu or g.bottom > slide_height_emu:
                issues.append(SafetyIssue(
                    severity="critical",
                    check="out_of_bounds",
                    detail="shape id %d would fall outside the slide bounds after this edit" % shape_id,
                ))

    # 3. Changed relative spacing between moved dependents and their
    # immediate chain neighbor (the shape each move's shift was computed
    # relative to) - the reflow planner is supposed to preserve gaps
    # exactly, so any drift here means the plan itself has a bug.
    for move in reflow.moves:
        shape = shapes_by_id.get(move.shape_id)
        if shape is None:
            continue
        # Nearest shape above, by original geometry, horizontally overlapping.
        above_candidates = [
            s for s in scene.shapes
            if s.shape_id != shape.shape_id
            and s.geometry.horizontal_overlap(shape.geometry)
            and s.geometry.bottom <= shape.geometry.top
        ]
        if not above_candidates:
            continue
        nearest_above = max(above_candidates, key=lambda s: s.geometry.bottom)
        gap_before = shape.geometry.top - nearest_above.geometry.bottom
        gap_after = finals[shape.shape_id].top - finals[nearest_above.shape_id].bottom
        if abs(gap_after - gap_before) > 1:  # 1 EMU: exact-arithmetic tolerance only
            issues.append(SafetyIssue(
                severity="warning",
                check="relative_spacing_changed",
                detail="gap between shape id %d and shape id %d changed from %d to %d EMU"
                       % (shape.shape_id, nearest_above.shape_id, gap_before, gap_after),
            ))

    # 4. Alignment breakage: any `aligned-with` pair must still be aligned.
    for r in relationships:
        if r.relation_type != "aligned-with":
            continue
        if r.source_shape_id not in finals or r.target_shape_id not in finals:
            continue
        left_before_a = shapes_by_id[r.source_shape_id].geometry.left
        left_before_b = shapes_by_id[r.target_shape_id].geometry.left
        left_after_a = finals[r.source_shape_id].left
        left_after_b = finals[r.target_shape_id].left
        was_aligned = abs(left_before_a - left_before_b) <= 40_000
        now_aligned = abs(left_after_a - left_after_b) <= 40_000
        if was_aligned and not now_aligned:
            issues.append(SafetyIssue(
                severity="critical",
                check="alignment_breakage",
                detail="shape id %d and shape id %d were left-aligned before this edit "
                       "and would no longer be aligned after it" % (r.source_shape_id, r.target_shape_id),
            ))

    # 5. Text overflow/wrap heuristic (approximate; see module docstring).
    if estimated_line_count_before is not None and estimated_line_count_after is not None:
        if estimated_line_count_after > estimated_line_count_before:
            issues.append(SafetyIssue(
                severity="warning",
                check="possible_text_overflow",
                detail="estimated wrapped line count increased from %d to %d; verify in native "
                       "PowerPoint that this still fits the shape's rendered box"
                       % (estimated_line_count_before, estimated_line_count_after),
            ))

    # 6 & 7. Relationship target displacement / connector-pointer consistency:
    # a `points-to`, `highlights`, or `belongs-to` relationship where exactly
    # one side moved and the other did not is a pointer that would become
    # misdirected or stranded.
    for r in relationships:
        if r.relation_type not in ("points-to", "highlights", "belongs-to"):
            continue
        source_moved = r.source_shape_id in moved_ids
        target_moved = r.target_shape_id in moved_ids
        if source_moved != target_moved:
            issues.append(SafetyIssue(
                severity="critical",
                check="relationship_target_displacement",
                detail="%s relationship between shape id %d and shape id %d would be broken - "
                       "only one side of the relationship moves in this plan"
                       % (r.relation_type, r.source_shape_id, r.target_shape_id),
            ))

    return SafetyReport(issues=issues)


def estimate_line_count(text: str, geometry: Geometry, font_size_pt: float,
                         margin_left_emu: int = 0, margin_right_emu: int = 0) -> int:
    """A deliberately approximate, deterministic wrap-line estimate -
    average-glyph-width based, not real font metrics or a layout engine.
    Used only to flag *possible* overflow for human review; never treated
    as ground truth. See module docstring."""
    usable_width_emu = geometry.width - margin_left_emu - margin_right_emu
    if usable_width_emu <= 0 or font_size_pt <= 0:
        return max(1, text.count("\n") + 1)
    char_width_emu = _AVG_GLYPH_WIDTH_EM * font_size_pt * _EMU_PER_POINT
    chars_per_line = max(1, int(usable_width_emu / char_width_emu))
    total_lines = 0
    for paragraph in text.split("\n"):
        if not paragraph:
            total_lines += 1
            continue
        total_lines += -(-len(paragraph) // chars_per_line)  # ceil division
    return max(1, total_lines)


def max_lines_that_fit(geometry: Geometry, font_size_pt: float,
                        margin_top_emu: int = 0, margin_bottom_emu: int = 0) -> int:
    """How many wrapped lines fit in the shape's **existing** container
    height, at its current size - the fixed-height-overflow counterpart
    to `estimate_line_count`'s width-based wrap estimate (Codex
    PEI-S1-08, re-audit 2). Same explicit caveat: an approximate,
    deterministic, average-line-height heuristic, not real font metrics
    or a layout engine, and never treated as ground truth - it exists
    only to decide, conservatively, whether Slice 1 can proceed without
    ever resizing the container. Independent of PowerPoint `auto_size`
    mode: Slice 1 never resizes the target regardless of what
    `auto_size` says, so how many lines the *current* box height can
    hold is what actually matters, not whether PowerPoint would
    eventually auto-grow it on open."""
    usable_height_emu = geometry.height - margin_top_emu - margin_bottom_emu
    line_height_emu = max(1, int(_LINE_HEIGHT_EM * max(font_size_pt, 1) * _EMU_PER_POINT))
    if usable_height_emu <= 0:
        return 0
    return max(0, int(usable_height_emu / line_height_emu))


def resolve_conservative_font_size_pt(text_frame) -> Optional[float]:
    """The largest effective font size across every run in `text_frame`
    (Codex PEI-S1-08, final repair - the prior estimate used whichever
    run happened to be first with an explicit size, silently
    underestimating fit for text mixing a small "label" run with a much
    larger edited run - reproduced by Codex with an 8pt first run and a
    40pt edited run).

    Each run's size is resolved from its own explicit `Font.size`, or -
    when the run itself does not set one - its paragraph's own explicit
    default `Font.size`. Those are the only two levels of size
    inheritance python-pptx exposes deterministically; this function does
    not walk placeholder/layout/master-level inheritance. Returns `None`
    the moment ANY run's size cannot be resolved at either level -
    "uncertain fit => refuse" applies to the formatting resolution itself,
    not only to the resulting width/height comparison: one unresolvable
    run makes the whole estimate unsafe, not just that run's own
    contribution to it.

    The maximum (never the first, never an average) is used deliberately:
    a larger font is both wider per character and taller per line, so the
    maximum is the conservative choice for both the width-based wrap
    estimate and the height-based capacity estimate that use this
    result."""
    sizes: List[float] = []
    for paragraph in text_frame.paragraphs:
        paragraph_default = paragraph.font.size
        for run in paragraph.runs:
            resolved = run.font.size if run.font.size is not None else paragraph_default
            if resolved is None:
                return None
            sizes.append(resolved.pt)
    return max(sizes) if sizes else None


def _paragraph_conservative_line_height_emu(paragraph, base_line_height_emu: int) -> int:
    """Conservative per-line height for one paragraph: the font-size-
    derived base line height, overridden by the paragraph's own explicit
    `line_spacing` when set - a `Length` (e.g. `Pt(12)`; python-pptx
    returns a `Centipoints`/`Emu` instance for absolute line spacing) is
    an absolute per-line height, used directly in EMU; a plain float/int
    that is NOT a `Length` is a multiple of the base height. `Length`
    MUST be checked before the generic numeric check below - every
    python-pptx `Length` subclass (including `Centipoints`) is itself an
    `int` subclass, so checking `isinstance(x, (int, float))` first would
    silently misclassify every absolute value as a multiplier (Codex
    PEI-S1-08, second finding). Left as the unmodified base height when
    `line_spacing` is unset - single-line spacing is python-pptx/
    PowerPoint's own real default in that case, not an unknown; only
    font-size resolution failure (see `resolve_conservative_font_size_pt`)
    triggers an outright refusal in this module - an absent,
    deterministically-default property does not."""
    line_spacing = paragraph.line_spacing
    if line_spacing is None:
        return base_line_height_emu
    if isinstance(line_spacing, Length):
        return max(1, int(line_spacing))
    if isinstance(line_spacing, (int, float)) and not isinstance(line_spacing, bool):
        return max(1, int(base_line_height_emu * line_spacing))
    return base_line_height_emu


def _paragraph_spacing_emu(paragraph) -> int:
    """Deterministically-available paragraph `space_before` +
    `space_after`, in EMU - 0 when a given one is unset (PowerPoint's own
    real default, not an unknown)."""
    total = 0
    for value in (paragraph.space_before, paragraph.space_after):
        if value is not None:
            try:
                total += int(value)
            except (TypeError, ValueError):
                pass
    return total


def estimate_required_height_emu(
    paragraphs, paragraph_texts: List[str], geometry: Geometry, font_size_pt: float,
    margin_left_emu: int = 0, margin_right_emu: int = 0,
) -> int:
    """Conservative total required height for `paragraph_texts` (one
    string per paragraph, in paragraph order) rendered at `font_size_pt`
    within `geometry`'s width - the mixed-formatting- and spacing-aware
    replacement for treating a whole text frame as one flat wrapped-line
    count (Codex PEI-S1-08, final repair: "account conservatively for
    paragraph line spacing and before/after spacing where deterministically
    available").

    `paragraphs` supplies the real `Paragraph` objects - used only for
    their `line_spacing`/`space_before`/`space_after` facts, never their
    own `.text` - so the *same* paragraph formatting can be combined with
    two different `paragraph_texts` lists (the text frame's current
    content, and the approved replacement) to get a directly-comparable
    before/after estimate. Slice 1 never adds or removes paragraphs, so
    paragraph count and per-paragraph formatting are stable across an
    edit - `paragraphs` and `paragraph_texts` are expected to be the same
    length; a caller that cannot guarantee this must not call this
    function (see `pilot.py`'s own length check, which refuses instead)."""
    usable_width_emu = geometry.width - margin_left_emu - margin_right_emu
    if usable_width_emu <= 0 or font_size_pt <= 0:
        chars_per_line = 1
    else:
        char_width_emu = _AVG_GLYPH_WIDTH_EM * font_size_pt * _EMU_PER_POINT
        chars_per_line = max(1, int(usable_width_emu / char_width_emu))
    base_line_height_emu = max(1, int(_LINE_HEIGHT_EM * font_size_pt * _EMU_PER_POINT))

    total_height = 0
    for paragraph, text in zip(paragraphs, paragraph_texts):
        lines = max(1, -(-len(text) // chars_per_line)) if text else 1
        line_height = _paragraph_conservative_line_height_emu(paragraph, base_line_height_emu)
        total_height += lines * line_height
        total_height += _paragraph_spacing_emu(paragraph)
    return total_height
