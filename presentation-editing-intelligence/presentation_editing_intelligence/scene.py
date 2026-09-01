"""Local scene extraction.

Builds a compact, deterministic representation of the shapes near one
edit target on one slide - never the whole slide indiscriminately, and
never the whole deck. Every fact here comes from OOXML/python-pptx
structure directly (geometry, shape type, line/fill color, connector
endpoints, text) - nothing here is inferred or guessed; inference is
`relationships.py`'s job, one layer up, operating on this scene.

Scope: this looks only at the target slide's own top-level shapes (not
into other slides, not recursively unpacking every shape in every group
elsewhere in the deck) - `inspect_deck()`-style whole-deck loading is
deliberately not used here.
"""

from __future__ import annotations

import dataclasses
from typing import List, Optional, Tuple

from pptx.presentation import Presentation as PresentationType


@dataclasses.dataclass(frozen=True)
class Geometry:
    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    def intersects(self, other: "Geometry", margin: int = 0) -> bool:
        return not (
            self.right + margin < other.left
            or other.right + margin < self.left
            or self.bottom + margin < other.top
            or other.bottom + margin < self.top
        )

    def horizontal_overlap(self, other: "Geometry") -> bool:
        return self.left < other.right and other.left < self.right


@dataclasses.dataclass(frozen=True)
class ConnectorEndpoints:
    """Raw `<a:stCxn>`/`<a:endCxn>` connection-site bindings, when the
    shape is a real bound connector - `shape_id` of the connected shape
    on each end, or `None` if that end is unbound (a free endpoint, still
    positioned in space but not structurally attached to anything)."""

    start_shape_id: Optional[int]
    end_shape_id: Optional[int]


@dataclasses.dataclass(frozen=True)
class SceneShape:
    """One shape's deterministic facts, as needed for relationship
    inference and reflow planning - not a general-purpose shape dump."""

    shape_index: int
    shape_id: int
    shape_type: str
    name: str
    geometry: Geometry
    z_order: int
    text: Optional[str]
    line_color_rgb: Optional[str]
    line_width_emu: Optional[int]
    fill_color_rgb: Optional[str]
    is_line_shape: bool  # prstGeom line/straightConnector/bentConnector/curvedConnector
    is_freeform: bool
    # Codex PEI-S1-05: group-member extraction is NOT implemented in this
    # slice (it would require re-basing each member's coordinates through
    # its group's chOff/chExt transform, a materially larger scope) - this
    # field is always None, never populated, and no relationship in this
    # slice (in particular, no `moves-with`) is derived from it. Deferred
    # to a future slice; not claimed as working capability anywhere.
    group_shape_index: Optional[int]
    connector_endpoints: Optional[ConnectorEndpoints]


@dataclasses.dataclass(frozen=True)
class LocalScene:
    """The compact local scene around one edit target."""

    slide_index: int
    target_shape_index: int
    shapes: List[SceneShape]  # includes the target itself, always first

    @property
    def target(self) -> SceneShape:
        return self.shapes[0]

    @property
    def context(self) -> List[SceneShape]:
        return self.shapes[1:]


_LINE_GEOM_PRESETS = {"line", "straightConnector1", "bentConnector2", "bentConnector3",
                      "bentConnector4", "bentConnector5", "curvedConnector2",
                      "curvedConnector3", "curvedConnector4", "curvedConnector5"}


def _rgb_of(color_format) -> Optional[str]:
    try:
        if color_format is not None and color_format.type is not None:
            return str(color_format.rgb)
    except Exception:
        return None
    return None


def _line_facts(shape) -> Tuple[Optional[str], Optional[int], bool]:
    try:
        line = shape.line
    except Exception:
        return None, None, False
    color = _rgb_of(getattr(line, "color", None))
    width = None
    try:
        width = line.width.emu if line.width is not None else None
    except Exception:
        width = None
    return color, width, False


def _fill_facts(shape) -> Optional[str]:
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # MSO_FILL.SOLID
            return _rgb_of(fill.fore_color)
    except Exception:
        return None
    return None


def _geom_preset(shape) -> Optional[str]:
    try:
        geom = shape._element.spPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom"
        )
        if geom is not None:
            return geom.get("prst")
    except Exception:
        return None
    return None


def _connector_endpoints_from_element(shape_element) -> Optional[ConnectorEndpoints]:
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    if shape_element.tag != "{http://schemas.openxmlformats.org/presentationml/2006/main}cxnSp":
        return None
    nv = shape_element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}nvCxnSpPr"
    )
    if nv is None:
        return None
    cnv = nv.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cNvCxnSpPr")
    if cnv is None:
        return None
    st = cnv.find(ns + "stCxn")
    end = cnv.find(ns + "endCxn")
    start_id = int(st.get("id")) if st is not None and st.get("id") is not None else None
    end_id = int(end.get("id")) if end is not None and end.get("id") is not None else None
    if start_id is None and end_id is None:
        return None
    return ConnectorEndpoints(start_shape_id=start_id, end_shape_id=end_id)


def _shape_text(shape) -> Optional[str]:
    if getattr(shape, "has_text_frame", False):
        try:
            return shape.text_frame.text
        except Exception:
            return None
    return None


def _build_scene_shape(shape, shape_index: int, z_order: int, group_shape_index: Optional[int]) -> SceneShape:
    geometry = Geometry(
        left=shape.left or 0,
        top=shape.top or 0,
        width=shape.width or 0,
        height=shape.height or 0,
    )
    line_color, line_width, _ = _line_facts(shape)
    fill_color = _fill_facts(shape)
    preset = _geom_preset(shape)
    is_line = bool(preset and preset in _LINE_GEOM_PRESETS)
    is_freeform = bool(preset == "custom" or shape.shape_type is not None and str(shape.shape_type).startswith("FREEFORM"))
    endpoints = _connector_endpoints_from_element(shape._element)

    return SceneShape(
        shape_index=shape_index,
        shape_id=shape.shape_id,
        shape_type=str(shape.shape_type),
        name=shape.name,
        geometry=geometry,
        z_order=z_order,
        text=_shape_text(shape),
        line_color_rgb=line_color,
        line_width_emu=line_width,
        fill_color_rgb=fill_color,
        is_line_shape=is_line,
        is_freeform=is_freeform,
        group_shape_index=group_shape_index,
        connector_endpoints=endpoints,
    )


def extract_all_shapes_on_slide(prs: PresentationType, slide_index: int) -> List[SceneShape]:
    """Every top-level shape on one slide, unfiltered - the full,
    deterministic geometry index Codex finding PEI-S1-06 requires as the
    collision-safety candidate set (distinct from the local *reasoning*
    scene `extract_local_scene` returns, which stays proximity-bounded).
    Building this is cheap - one slide's shapes, not the whole deck - and
    reuses exactly the same per-shape fact extraction `extract_local_scene`
    itself uses, so the two are always consistent with each other."""
    if not 0 <= slide_index < len(prs.slides):
        raise IndexError("slide_index %d out of range" % slide_index)
    shapes = list(prs.slides[slide_index].shapes)
    return [_build_scene_shape(shape, i, z_order=i, group_shape_index=None) for i, shape in enumerate(shapes)]


def extract_local_scene(
    prs: PresentationType,
    slide_index: int,
    target_shape_index: int,
    proximity_margin_emu: int = 1_000_000,
) -> LocalScene:
    """Extracts the target shape plus every other top-level shape on the
    same slide whose bounding box lies within `proximity_margin_emu` of
    the target's bounding box (a bounded, deterministic proximity rule -
    not "everything on the slide", not "everything below"), plus any
    shape whose real connector endpoint (`<a:stCxn>`/`<a:endCxn>`)
    structurally references the target's `shape_id` regardless of
    distance (a genuine structural relationship always belongs in the
    local scene, even if visually far away).

    `proximity_margin_emu` defaults to roughly 1.1 inches (914400 EMU per
    inch) - enough to capture immediately adjacent dividers/pointers/
    highlights without pulling in unrelated, distant slide content. This
    *reasoning* scope is intentionally narrower than the collision-safety
    candidate set `extract_all_shapes_on_slide` provides - relationship
    inference stays local; collision detection (see `safety.py`) does not
    have to (PEI-S1-06).
    """
    scene_shapes = extract_all_shapes_on_slide(prs, slide_index)
    if not 0 <= target_shape_index < len(scene_shapes):
        raise IndexError("shape_index %d out of range" % target_shape_index)

    target = scene_shapes[target_shape_index]
    target_geom = target.geometry

    selected = [target]
    for s in scene_shapes:
        if s.shape_index == target.shape_index:
            continue
        connected = (
            s.connector_endpoints is not None
            and target.shape_id in (s.connector_endpoints.start_shape_id, s.connector_endpoints.end_shape_id)
        )
        near = target_geom.intersects(s.geometry, margin=proximity_margin_emu)
        if connected or near:
            selected.append(s)

    return LocalScene(slide_index=slide_index, target_shape_index=target_shape_index, shapes=selected)
