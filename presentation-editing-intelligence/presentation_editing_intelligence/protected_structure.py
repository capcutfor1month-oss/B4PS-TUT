"""Protected-structure verification: the final gate before a scene-aware
edit's result may be reported as `"applied"`.

Codex finding PEI-S1-02 (FAIL, repaired): the pilot could previously
detect an unrelated mutation via `unrelated_shapes_unchanged=False` and
still keep the published output and return `status="applied"` - the flag
was informational, not enforced. It also only ever compared flattened
`text`, ignoring geometry/identity/ordering/type changes on unrelated
shapes and all run-level formatting on the target shape's own untouched
runs.

This module is the single point that decides whether a published result
may stand. A successful result may contain ONLY:

1. the approved target text run changes;
2. the approved geometry/reflow changes;
3. nothing else.

Everything else - slide count/order, every shape's identity (`shape_id`),
type, geometry, and text; the target shape's own untouched runs'
formatting; and (where this slice's scene model has evidence for it) the
connector/group structure of shapes in the local scene - must be
byte-for-byte identical to before. Any difference is `ok=False`; the
caller (`pilot.py`) must then remove the just-published output and never
report `"applied"`.
"""

from __future__ import annotations

import dataclasses
from typing import Dict, List, Optional, Set, Tuple

from pptx import Presentation

from ._safe_ppt_engine_import import inspect_deck
from .reflow import ReflowMove
from .scene import extract_local_scene
from .text_edit import RunEdit


@dataclasses.dataclass(frozen=True)
class ProtectedStructureResult:
    ok: bool
    reason: Optional[str] = None


def _run_format_signature(run) -> tuple:
    font = run.font
    color = None
    try:
        if font.color is not None and font.color.type is not None:
            color = (str(font.color.type), str(font.color.rgb) if font.color.type == 1 else str(font.color.theme_color))
    except Exception:
        color = None
    return (font.bold, font.italic, font.underline, font.size, font.name, color)


def verify(
    source_deck_path: str,
    output_deck_path: str,
    slide_index: int,
    shape_index: int,
    run_edits: List[RunEdit],
    reflow_moves: List[ReflowMove],
    expected_target_full_text: str,
) -> ProtectedStructureResult:
    """Deterministic before/after comparison. Uses the Safe PPT Engine's
    own `inspect_deck()` (unmodified) for the deck-wide slide/shape
    identity, type, geometry, ordering, and text comparison - the same
    representation Documentation Intelligence Slice 3's own
    `_verify_mutation` relied on, now applied deck-wide and combined with
    a dedicated run-level formatting check and a local connector/group
    structural check this slice's own scene model can evidence."""
    before = inspect_deck(source_deck_path)
    after = inspect_deck(output_deck_path)

    if before["slide_count"] != after["slide_count"]:
        return ProtectedStructureResult(False, "slide count changed: expected %d, got %d"
                                         % (before["slide_count"], after["slide_count"]))

    moved_shape_indices: Set[int] = {m.shape_index for m in reflow_moves}
    moves_by_index: Dict[int, ReflowMove] = {m.shape_index: m for m in reflow_moves}
    allowed_on_target_slide = {shape_index} | moved_shape_indices

    for s_idx in range(before["slide_count"]):
        before_shapes = before["slides"][s_idx]["shapes"]
        after_shapes = after["slides"][s_idx]["shapes"]
        if len(before_shapes) != len(after_shapes):
            return ProtectedStructureResult(
                False, "shape count on slide %d changed: expected %d, got %d"
                       % (s_idx, len(before_shapes), len(after_shapes))
            )

        # Shape identity/ordering (shape_id sequence) must be identical -
        # this alone catches any insertion, deletion, or reordering.
        before_ids = [sh["shape_id"] for sh in before_shapes]
        after_ids = [sh["shape_id"] for sh in after_shapes]
        if before_ids != after_ids:
            return ProtectedStructureResult(
                False, "shape identity/order changed on slide %d: expected %r, got %r"
                       % (s_idx, before_ids, after_ids)
            )

        for sh_idx in range(len(before_shapes)):
            before_sh = before_shapes[sh_idx]
            after_sh = after_shapes[sh_idx]
            is_target = (s_idx, sh_idx) == (slide_index, shape_index)
            is_reflow_move = s_idx == slide_index and sh_idx in moved_shape_indices

            if before_sh["shape_type"] != after_sh["shape_type"]:
                return ProtectedStructureResult(
                    False, "shape type changed at slide %d shape %d (not an approved change)"
                           % (s_idx, sh_idx)
                )

            if is_target:
                if after_sh["text"] != expected_target_full_text:
                    return ProtectedStructureResult(
                        False, "target shape's text does not equal the exact approved "
                               "replacement at slide %d shape %d" % (s_idx, sh_idx)
                    )
                for attr in ("left", "top", "width", "height"):
                    if before_sh[attr] != after_sh[attr]:
                        return ProtectedStructureResult(
                            False, "target shape's own geometry (%s) changed at slide %d shape %d, "
                                   "but no geometry change for the target itself was approved"
                                   % (attr, s_idx, sh_idx)
                        )
                continue

            if is_reflow_move:
                move = moves_by_index[sh_idx]
                if after_sh["top"] != move.new_top_emu:
                    return ProtectedStructureResult(
                        False, "reflow-moved shape at slide %d shape %d does not have the "
                               "exact approved top after publication" % (s_idx, sh_idx)
                    )
                for attr in ("left", "width", "height"):
                    if before_sh[attr] != after_sh[attr]:
                        return ProtectedStructureResult(
                            False, "reflow-moved shape's %s changed at slide %d shape %d, but "
                                   "only top was approved to move" % (attr, s_idx, sh_idx)
                        )
                if before_sh["text"] != after_sh["text"]:
                    return ProtectedStructureResult(
                        False, "reflow-moved shape's text changed at slide %d shape %d - not approved"
                               % (s_idx, sh_idx)
                    )
                continue

            # Every other shape, on every slide: fully unchanged.
            for attr in ("left", "top", "width", "height", "text", "name", "has_text_frame", "has_image"):
                if before_sh[attr] != after_sh[attr]:
                    return ProtectedStructureResult(
                        False, "unauthorized change detected at slide %d shape %d (%s changed) - "
                               "not part of the approved edit plan" % (s_idx, sh_idx, attr)
                    )

    # Run-level formatting (Codex PEI-S1-02.A, re-audit): EVERY run on the
    # target shape - including the ones named in `run_edits` - must keep
    # its formatting unchanged. Slice 1 has no approved formatting
    # mutation anywhere in its plan vocabulary (`RunEdit` carries only
    # `new_text`, never a formatting change) - so an edited run's text
    # may change exactly as approved, but its `<a:rPr>` must be
    # byte-for-byte identical to before, exactly like every untouched
    # run. `inspect_deck()` only reports flattened text, not per-run
    # formatting, so this is checked separately, over every run with no
    # exemption.
    edited_runs: Set[Tuple[int, int]] = {(e.paragraph_index, e.run_index) for e in run_edits}
    edited_run_new_text = {(e.paragraph_index, e.run_index): e.new_text for e in run_edits}
    before_prs = Presentation(source_deck_path)
    after_prs = Presentation(output_deck_path)
    before_target = list(before_prs.slides[slide_index].shapes)[shape_index]
    after_target = list(after_prs.slides[slide_index].shapes)[shape_index]
    if before_target.has_text_frame and after_target.has_text_frame:
        before_paras = before_target.text_frame.paragraphs
        after_paras = after_target.text_frame.paragraphs
        if len(before_paras) != len(after_paras):
            return ProtectedStructureResult(
                False, "target shape's paragraph count changed - not an approved change"
            )
        for p_idx, (before_p, after_p) in enumerate(zip(before_paras, after_paras)):
            before_runs = before_p.runs
            after_runs = after_p.runs
            if len(before_runs) != len(after_runs):
                return ProtectedStructureResult(
                    False, "target shape's run count in paragraph %d changed - not an approved change" % p_idx
                )
            for r_idx, (before_r, after_r) in enumerate(zip(before_runs, after_runs)):
                is_edited = (p_idx, r_idx) in edited_runs
                if is_edited:
                    expected_text = edited_run_new_text[(p_idx, r_idx)]
                    if after_r.text != expected_text:
                        return ProtectedStructureResult(
                            False, "target shape's run (paragraph %d, run %d) does not contain "
                                   "the exact approved replacement text" % (p_idx, r_idx)
                        )
                else:
                    if before_r.text != after_r.text:
                        return ProtectedStructureResult(
                            False, "target shape's run (paragraph %d, run %d) text changed but was "
                                   "not part of the approved edit" % (p_idx, r_idx)
                        )
                # Formatting is checked identically for every run,
                # edited or not - no exemption for the edited run(s).
                if _run_format_signature(before_r) != _run_format_signature(after_r):
                    return ProtectedStructureResult(
                        False, "target shape's run (paragraph %d, run %d) formatting changed, "
                               "but Slice 1 has no approved formatting mutation - only the "
                               "run's text may change" % (p_idx, r_idx)
                    )

    # Connector/group structure, where this slice's scene model has
    # evidence for it: every shape's connector endpoints in the local
    # scene around the target must be identical before/after (moving a
    # shape by `top` never changes which shapes its connector endpoints
    # reference).
    before_scene = extract_local_scene(before_prs, slide_index, shape_index)
    after_scene = extract_local_scene(after_prs, slide_index, shape_index)
    before_by_id = {s.shape_id: s for s in before_scene.shapes}
    after_by_id = {s.shape_id: s for s in after_scene.shapes}
    for shape_id, before_s in before_by_id.items():
        after_s = after_by_id.get(shape_id)
        if after_s is None:
            continue  # already caught by the deck-wide identity check above
        if before_s.connector_endpoints != after_s.connector_endpoints:
            return ProtectedStructureResult(
                False, "connector endpoint structure changed for shape id %d - not approved" % shape_id
            )

    return ProtectedStructureResult(True, None)
