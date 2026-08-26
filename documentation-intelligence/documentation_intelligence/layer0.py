"""Documentation Intelligence - Layer 0: Deterministic MasterSlide Inventory.

Implements the approved Layer 0 specification
(`openspec/changes/masterslide-layer0-inventory/spec.md`). Read-only:
given one target `.pptx`, deterministically extracts a reusable, compact
structural inventory of every slide - shape geometry, groups, connector
bindings, freeform paths, picture crops, text style summaries - so later
presentation-editing agents never need to repeatedly re-read raw OOXML or
rediscover the same PowerPoint facts.

Locked boundaries this module enforces in code, not only documentation:

- Zero PPT mutation - no Safe PPT Engine mutation primitive is called or
  imported anywhere in this module. The deck is opened once, read-only,
  via the same typed `load_deck()` boundary Slice 3 already uses.
- Zero semantic/editorial judgment - this module extracts *what exists*
  (shape types, geometry, bindings, style statistics), never *what it
  means*. It never decides which editing era a slide belongs to, who
  edited it, what the "current" editorial style is, or how a slide
  should be edited or created. Any clustering/range output is explicitly
  labelled a structural candidate, never an editorial conclusion.
- Deterministic and reproducible - two calls against the same byte-
  identical `.pptx` (regardless of whether the caller spells the path
  relatively or absolutely) produce byte-identical Tier A/Tier B/
  manifest/structural-analysis JSON. No timestamps, random seeds without
  a fixed value, or caller-supplied path spelling appear anywhere in the
  canonical output.
- No visual model, no embeddings, no LLM call anywhere in this module.
- No dependency on OfficeCLI or any renderer. OOXML structure is the
  only evidence source; a native PowerPoint/PDF render pipeline is a
  separate, later concern (Presentation Exploration), never this layer.

Two-tier output, matching the validated research design:

- **Tier A** (`extract_tier_a`): one compact per-slide summary record -
  counts, histograms, the structural fingerprint, the normalized title/
  topic key. Cheap enough for the complete deck to be loaded in one
  context (range decomposition, clustering, outlier detection).
- **Tier B** (`extract_tier_b`): one full per-shape detail record per
  slide, retrieved lazily for only the slide(s) actually needed. Never
  bulk-loaded by Tier A consumers.

`generate_layer0_inventory()` is the single production entry point. It
re-verifies source integrity around the extraction (`SourceIntegrityError`
on any mismatch - see "Source-integrity transaction" below), stages the
complete output tree privately, validates it, and only then publishes it
as one atomically-swapped, exclusively Layer-0-owned tree under
`output_dir` (`StagingValidationError`/`PublicationError` on failure - see
"Atomic publication" below). A failed run never modifies a previously
published tree, and unrelated files that happen to live alongside it are
never touched.
"""

from __future__ import annotations

import contextlib
import fcntl
import hashlib
import json
import math
import os
import random
import re
import shutil
import uuid
from collections import Counter
from pathlib import Path
from typing import Optional

from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.oxml.ns import qn

from ._safe_ppt_engine_import import DeckSourceError, file_sha256, load_deck

__all__ = [
    "SCHEMA_VERSION",
    "Layer0Error",
    "SourceIntegrityError",
    "StagingValidationError",
    "PublicationError",
    "PublicationRollbackError",
    "PublicationLockError",
    "PublicationStagingError",
    "is_tree_valid",
    "extract_deck_identity",
    "extract_tier_a",
    "extract_tier_b",
    "extract_title",
    "compute_fingerprint",
    "assign_topic_occurrence_ordinals",
    "propose_ranges",
    "cluster_structural_candidates",
    "generate_layer0_inventory",
]

SCHEMA_VERSION = "1.1.0"

_REPO_ROOT = Path(__file__).resolve().parents[2]

_EMU_PER_IN = 914400
_TITLE_HEADER_ZONE_RATIO = 0.15
_DESKTOP_PREFIX_RE = re.compile(r"^\s*desktop\s*:\s*", re.IGNORECASE)
_WHITESPACE_RE = re.compile(r"\s+")
_TIER_B_SLIDE_FILE_RE = re.compile(r"^slide_\d{4}\.json$")

_FINGERPRINT_KEYS = (
    "sp_rate",
    "pic_rate",
    "grp_rate",
    "cxn_rate",
    "bound_rate",
    "crop_rate",
    "freeform_rate",
    "shape_count_norm",
)

_OWNED_TOP_LEVEL_FILES = ("manifest.json", "tier_a_summary.json", "structural_analysis.json")
_TIER_B_DIRNAME = "tier_b"
_COMMIT_MARKER_NAME = "_committed.json"


# --------------------------------------------------------------------------
# typed errors - distinct from the Safe PPT Engine's `SafeDeckError`
# hierarchy, since this module performs no mutation. `SourceIntegrityError`
# is also a `DeckSourceError` so existing "except DeckSourceError" callers
# continue to catch it, matching Slice 3's own `ShapeIdMismatchError`
# pattern of narrow typed subclasses.
# --------------------------------------------------------------------------

class Layer0Error(Exception):
    """Base for Layer 0's own typed errors."""


class SourceIntegrityError(DeckSourceError, Layer0Error):
    """The source deck's bytes changed between the pre-extraction hash and
    the pre-publication re-hash (L0-01) - Layer 0 refuses to publish an
    inventory that may not describe the bytes it was actually extracted
    from. Raised before any staging directory is created; nothing is
    ever published when this is raised."""


class StagingValidationError(Layer0Error):
    """The staged Layer 0 tree failed its own internal consistency check
    before publication was attempted - `output_dir` (any previously
    published tree) was never touched."""


class PublicationError(Layer0Error):
    """Publishing the staged tree into `output_dir` failed. Any partially
    applied change is rolled back to the previous state before this is
    raised, except where explicitly noted (best-effort post-publish
    cleanup of now-stale previous artifacts). Covers: `output_dir`
    existing as a non-directory, staging-directory creation failure,
    commit failure (with successful rollback), and post-success cleanup
    failure. No raw `OSError`/`FileExistsError` escapes this module's
    publication path - see `PublicationRollbackError`/`PublicationLockError`
    for the two narrower cases that need their own semantics."""


class PublicationRollbackError(PublicationError):
    """Publication failed AND the automatic rollback that normally
    restores the previous tree could not fully complete. The previous,
    still-valid artifacts are retained under their `.prev-<token>` backup
    names and are deliberately NOT deleted in this case - this is the
    narrow, documented exception to "a failed publication always leaves
    the previous tree looking exactly as it did before": the previous
    tree's content survives, but under a backup name rather than its
    canonical name, and manual recovery (renaming `<name>.prev-<token>`
    back to `<name>`) is required before `output_dir` can be trusted
    again. `generate_layer0_inventory` never silently treats this as a
    successful publish."""


class PublicationLockError(PublicationError):
    """Another Layer 0 generation currently owns the exclusive
    single-writer lock for this `output_dir` (see `_acquire_publish_lock`).
    Local, single-host, single-filesystem advisory locking only, via
    `fcntl.flock` on a dedicated lock file - this provides no distributed
    or network-filesystem safety (`flock` semantics are unreliable or
    unsupported on some network filesystems, e.g. classic NFSv3), and is
    POSIX-only. Fails fast (non-blocking) rather than waiting, by
    deliberate choice for simplicity."""


class PublicationStagingError(Layer0Error):
    """Writing the staged tree itself failed for a filesystem reason
    (disk full, permission denied, an unwritable staging directory,
    etc.) - never a raw `OSError`. This occurs entirely within the
    private staging directory, strictly before publication into
    `output_dir` is ever attempted - any previously published tree is
    left completely untouched."""


# --------------------------------------------------------------------------
# small geometry/XML helpers
# --------------------------------------------------------------------------

def _emu_to_in(value):
    return round(value / _EMU_PER_IN, 3) if value is not None else None


def _classify_tag(shape):
    return shape._element.tag.split("}")[-1]


def _get_hidden(shape):
    cNvPr = shape._element.find(f".//{qn('p:cNvPr')}")
    return cNvPr is not None and cNvPr.get("hidden") == "1"


def _get_rotation(shape):
    """The shape's own local rotation, as stored in its `<a:xfrm rot=".."/>`.
    Does not compose inherited rotation from any ancestor group - that
    composition is a separate, more involved transform this module does
    not attempt (structural-fact scope only, see `_resolve_geometry_chain`
    for the geometry-only transform this module does perform)."""
    try:
        return round(shape.rotation, 1)
    except Exception:
        return None


def _parse_nonneg_int(raw):
    if raw is None:
        return None
    try:
        value = int(raw)
    except (TypeError, ValueError):
        return None
    return value if value >= 0 else None


def _endpoint_status(cxn_element):
    """One connector endpoint's structural validity. `ABSENT` - no
    `stCxn`/`endCxn` element at all (an ordinary unbound line, not an
    anomaly). `VALID` - the element is present with a well-formed
    non-negative integer `id` and `idx`. `INVALID` - the element is
    present but `id`/`idx` is missing or not a well-formed non-negative
    integer - surfaced explicitly rather than silently treated as either
    a clean binding or a clean absence (L0-06)."""
    if cxn_element is None:
        return "ABSENT", None, None
    raw_id = cxn_element.get("id")
    raw_idx = cxn_element.get("idx")
    target_id = _parse_nonneg_int(raw_id)
    idx = _parse_nonneg_int(raw_idx)
    if target_id is None or idx is None:
        return "INVALID", raw_id, raw_idx
    return "VALID", raw_id, raw_idx


def _connector_facts(shape):
    """Only meaningful for `cxnSp`. `stCxn`/`endCxn` live under
    `p:nvCxnSpPr/p:cNvCxnSpPr` - a grandchild of the connector element,
    not a direct child - so this must search descendants. A direct-
    children-only search silently reports every connector as unbound;
    see `TestConnectorBindingDescendantSearch` for the regression.

    Aggregate `status` is `INVALID` if either endpoint is malformed
    (surfaced even if the other endpoint is clean - a malformed binding
    is real evidence of an anomaly, never silently downgraded to a plain
    absence), else `VALID` if at least one endpoint is a well-formed
    binding, else `ABSENT` (an ordinary unbound line/connector shape)."""
    if _classify_tag(shape) != "cxnSp":
        return None
    el = shape._element
    stCxn = el.find(f".//{qn('a:stCxn')}")
    endCxn = el.find(f".//{qn('a:endCxn')}")
    start_status, start_id, start_idx = _endpoint_status(stCxn)
    end_status, end_id, end_idx = _endpoint_status(endCxn)
    statuses = {start_status, end_status}
    if "INVALID" in statuses:
        aggregate = "INVALID"
    elif "VALID" in statuses:
        aggregate = "VALID"
    else:
        aggregate = "ABSENT"
    return {
        "status": aggregate,
        "start_endpoint": {
            "status": start_status,
            "target_shape_id": start_id,
            "connection_idx": start_idx,
        },
        "end_endpoint": {
            "status": end_status,
            "target_shape_id": end_id,
            "connection_idx": end_idx,
        },
    }


def _freeform_facts(shape):
    custGeom = shape._element.find(f".//{qn('a:custGeom')}")
    if custGeom is None:
        return None
    path = custGeom.find(f".//{qn('a:path')}")
    if path is None:
        return {"moveTo": 0, "lnTo": 0, "cubicBezTo": 0}
    counts = Counter(child.tag.split("}")[-1] for child in path)
    return {
        "moveTo": counts.get("moveTo", 0),
        "lnTo": counts.get("lnTo", 0),
        "cubicBezTo": counts.get("cubicBezTo", 0),
    }


def _picture_facts(shape):
    if _classify_tag(shape) != "pic":
        return None
    srcRect = shape._element.find(f".//{qn('a:srcRect')}")
    facts = {
        "has_crop": srcRect is not None,
        "crop_l": srcRect.get("l") if srcRect is not None else None,
        "crop_t": srcRect.get("t") if srcRect is not None else None,
        "crop_r": srcRect.get("r") if srcRect is not None else None,
        "crop_b": srcRect.get("b") if srcRect is not None else None,
        "source_width_px": None,
        "source_height_px": None,
        "image_sha1": None,
    }
    try:
        image = shape.image
        facts["source_width_px"], facts["source_height_px"] = image.size
        facts["image_sha1"] = image.sha1
    except Exception:
        pass
    return facts


def _fill_facts(shape):
    try:
        fill = shape.fill
        ftype = str(fill.type) if fill.type is not None else None
        color = None
        if ftype and "SOLID" in ftype:
            try:
                color = str(fill.fore_color.rgb)
            except Exception:
                try:
                    color = f"theme:{fill.fore_color.theme_color}"
                except Exception:
                    color = None
        return {"type": ftype, "color": color}
    except Exception:
        return {"type": None, "color": None}


def _line_facts(shape):
    try:
        line = shape.line
        color = None
        try:
            color = str(line.color.rgb)
        except Exception:
            pass
        width = _emu_to_in(line.width) if line.width else None
        return {"color": color, "width_in": width}
    except Exception:
        return {"color": None, "width_in": None}


def _text_facts(shape):
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    runs = [r for p in paragraphs for r in p.runs]
    total_len = sum(len(r.text) for r in runs)
    if not runs:
        return {
            "paragraph_count": len(paragraphs),
            "run_count": 0,
            "text_length": total_len,
            "dominant_font": None,
            "dominant_size_pt": None,
            "dominant_color": None,
            "bold_run_ratio": 0.0,
            "italic_run_ratio": 0.0,
            "underline_run_ratio": 0.0,
            "alignment": str(paragraphs[0].alignment) if paragraphs and paragraphs[0].alignment else None,
            "autofit": _autofit_mode(tf),
            "word_wrap": tf.word_wrap,
        }
    fonts, sizes, colors = Counter(), Counter(), Counter()
    bold_ct = italic_ct = underline_ct = 0
    for r in runs:
        f = r.font
        if f.name:
            fonts[f.name] += 1
        if f.size:
            sizes[f.size.pt] += 1
        try:
            if f.color and f.color.type is not None:
                colors[str(f.color.rgb)] += 1
        except Exception:
            pass
        bold_ct += bool(f.bold)
        italic_ct += bool(f.italic)
        underline_ct += bool(f.underline)
    return {
        "paragraph_count": len(paragraphs),
        "run_count": len(runs),
        "text_length": total_len,
        "dominant_font": fonts.most_common(1)[0][0] if fonts else None,
        "dominant_size_pt": sizes.most_common(1)[0][0] if sizes else None,
        "dominant_color": colors.most_common(1)[0][0] if colors else None,
        "bold_run_ratio": round(bold_ct / len(runs), 2),
        "italic_run_ratio": round(italic_ct / len(runs), 2),
        "underline_run_ratio": round(underline_ct / len(runs), 2),
        "alignment": str(paragraphs[0].alignment) if paragraphs[0].alignment else None,
        "autofit": _autofit_mode(tf),
        "word_wrap": tf.word_wrap,
    }


def _autofit_mode(text_frame):
    try:
        bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
    except Exception:
        return None
    if bodyPr is None:
        return None
    if bodyPr.find(qn("a:spAutoFit")) is not None:
        return "spAutoFit"
    if bodyPr.find(qn("a:normAutofit")) is not None:
        return "normAutofit"
    if bodyPr.find(qn("a:noAutofit")) is not None:
        return "noAutofit"
    return None


# --------------------------------------------------------------------------
# group coordinate-space resolution (L0-04)
# --------------------------------------------------------------------------

def _group_transform(shape):
    """A group shape's `off`/`ext` (its own position/size in its
    *parent's* coordinate space) and `chOff`/`chExt` (the coordinate
    space its *children's* raw coordinates are expressed in), read
    directly from `p:grpSpPr/a:xfrm`. Returns `None` if any of the four
    is missing/malformed - callers must then treat any descendant's
    slide-resolved geometry as not reliably computable, never invented.

    Also returns `None` - deliberately, not an approximation - when this
    group's own `<a:xfrm>` carries a non-default `rot` or `flipH`/`flipV`
    (L0-04, confirmed present on the real deck: slide index 32's groups
    1255/1259/1263 carry `rot="10800000" flipH="1"`). Composing rotation/
    flip into the affine chain is out of scope for this bounded structural
    extractor; rather than silently ignore it (as the pre-repair version
    did, which produced a false non-null `resolved_geometry_in` for every
    descendant of such a group), any group with rotation/flip poisons its
    own transform and therefore the resolution of everything beneath it -
    `local_geometry_in`, hierarchy, and z-order are entirely unaffected."""
    grpSpPr = shape._element.find(qn("p:grpSpPr"))
    if grpSpPr is None:
        return None
    xfrm = grpSpPr.find(qn("a:xfrm"))
    if xfrm is None:
        return None
    rot = xfrm.get("rot")
    flip_h = xfrm.get("flipH")
    flip_v = xfrm.get("flipV")
    if (rot is not None and rot != "0") or flip_h == "1" or flip_v == "1":
        return None
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    ch_off = xfrm.find(qn("a:chOff"))
    ch_ext = xfrm.find(qn("a:chExt"))
    if off is None or ext is None or ch_off is None or ch_ext is None:
        return None
    try:
        return (
            int(off.get("x")), int(off.get("y")),
            int(ext.get("cx")), int(ext.get("cy")),
            int(ch_off.get("x")), int(ch_off.get("y")),
            int(ch_ext.get("cx")), int(ch_ext.get("cy")),
        )
    except (TypeError, ValueError):
        return None


def _resolve_geometry_chain(local_l, local_t, local_w, local_h, transform_stack):
    """Resolves one shape's local (immediate-parent-space) geometry to
    slide-absolute EMU by applying each ancestor group's transform, from
    the innermost (immediate parent) outward. `transform_stack` is
    root-to-immediate-parent order, so it is walked in reverse. Returns
    `None` - never a guess - the moment any ancestor transform is
    missing or has a zero child-extent (degenerate, unresolvable)."""
    l, t, w, h = local_l, local_t, local_w, local_h
    for transform in reversed(transform_stack):
        if transform is None:
            return None
        off_x, off_y, ext_cx, ext_cy, ch_off_x, ch_off_y, ch_ext_cx, ch_ext_cy = transform
        if ch_ext_cx == 0 or ch_ext_cy == 0:
            return None
        scale_x = ext_cx / ch_ext_cx
        scale_y = ext_cy / ch_ext_cy
        l = off_x + (l - ch_off_x) * scale_x
        t = off_y + (t - ch_off_y) * scale_y
        w = w * scale_x
        h = h * scale_y
    return l, t, w, h


# --------------------------------------------------------------------------
# Tier B - per-shape detail
# --------------------------------------------------------------------------

def _shape_record(shape, parent_group_id, depth, local_z_order, transform_stack):
    tag = _classify_tag(shape)
    try:
        shape_type = str(shape.shape_type) if shape.shape_type is not None else None
    except Exception:
        shape_type = None
    placeholder_type = None
    if shape.is_placeholder:
        try:
            placeholder_type = str(shape.placeholder_format.type)
        except Exception:
            placeholder_type = "UNKNOWN"

    local_l, local_t, local_w, local_h = shape.left, shape.top, shape.width, shape.height
    resolved = None
    if None not in (local_l, local_t, local_w, local_h):
        resolved = _resolve_geometry_chain(local_l, local_t, local_w, local_h, transform_stack)

    return {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "xml_tag": tag,
        "shape_type": shape_type,
        "local_z_order": local_z_order,
        "hidden": _get_hidden(shape),
        # Always in this shape's *immediate parent's* coordinate space -
        # for a top-level (ungrouped) shape that IS slide-absolute space;
        # for a grouped shape it is the group's own `chOff`/`chExt`
        # child space, honestly labelled `local_` rather than the
        # previous, misleading generic `geometry_in` (L0-04).
        "local_geometry_in": {
            "l": _emu_to_in(local_l), "t": _emu_to_in(local_t),
            "w": _emu_to_in(local_w), "h": _emu_to_in(local_h),
        },
        # Slide-absolute geometry, resolved via the ancestor group
        # transform chain where reliably computable (pure OOXML affine
        # arithmetic - never semantic inference); `None`, never a guess,
        # when any ancestor transform is missing/degenerate.
        "resolved_geometry_in": (
            {"l": _emu_to_in(resolved[0]), "t": _emu_to_in(resolved[1]),
             "w": _emu_to_in(resolved[2]), "h": _emu_to_in(resolved[3])}
            if resolved is not None else None
        ),
        "rotation": _get_rotation(shape),
        "parent_group_id": parent_group_id,
        "group_depth": depth,
        "child_ids": None,  # filled by the caller, immediate children only
        "is_placeholder": shape.is_placeholder,
        "placeholder_type": placeholder_type,
        "fill": _fill_facts(shape) if tag in ("sp", "cxnSp") else None,
        "line": _line_facts(shape) if tag in ("sp", "cxnSp") else None,
        "connector": _connector_facts(shape),
        "freeform": _freeform_facts(shape),
        "picture": _picture_facts(shape),
        "text": _text_facts(shape),
    }


def _walk_shapes(shapes, parent_group_id, depth, transform_stack):
    out = []
    for local_z_order, shape in enumerate(shapes):
        rec = _shape_record(shape, parent_group_id, depth, local_z_order, transform_stack)
        out.append(rec)
        if _classify_tag(shape) == "grpSp":
            # Immediate children only - computed directly from this one
            # level of `shape.shapes`, independent of how deep the
            # recursive walk below goes (L0-04: the previous version
            # flattened every recursive descendant into `child_ids`).
            rec["child_ids"] = [child.shape_id for child in shape.shapes]
            child_transform = _group_transform(shape)
            child_recs = _walk_shapes(
                shape.shapes, shape.shape_id, depth + 1, transform_stack + [child_transform]
            )
            out.extend(child_recs)
    return out


def extract_tier_b(slide, slide_index):
    """Full per-shape detail record for one slide. Meant for lazy,
    single-slide (or small-batch) retrieval - never bulk-loaded across
    the whole deck. See module docstring."""
    shapes = _walk_shapes(slide.shapes, None, 0, [])
    return {"slide_index": slide_index, "shapes": shapes}


# --------------------------------------------------------------------------
# title / topic-key extraction (L0-05: recurses into groups)
# --------------------------------------------------------------------------

def normalize_title(raw_title):
    """Whitespace-collapse and normalize the `Desktop:` / `Desktop :`
    convention to a single canonical spelling, for cross-deck matching.
    The original extracted title is preserved separately by callers -
    this function only produces the *normalized* key."""
    if raw_title is None:
        return None
    collapsed = _WHITESPACE_RE.sub(" ", raw_title).strip()
    normalized = _DESKTOP_PREFIX_RE.sub("desktop: ", collapsed)
    return normalized.lower()


def _title_placeholder_candidate(slide):
    for shape in slide.shapes:
        if not shape.is_placeholder or not shape.has_text_frame:
            continue
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            continue
        if ph_type in (PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE):
            text = shape.text_frame.text.strip()
            if text:
                return text, "title_placeholder"
    return None, None


def _iter_text_candidates(shapes, transform_stack):
    """Yields `(shape, resolved_top_emu_or_None)` for every text-bearing
    shape anywhere in the (possibly nested) shape tree - not just
    top-level `slide.shapes` (L0-05: a real title can sit inside a
    group). `resolved_top` is the shape's slide-absolute top when
    reliably computable via the group transform chain; `None` when it
    is not (an unresolvable ancestor transform) - title extraction then
    treats that candidate as not header-zone-eligible, never guesses."""
    for shape in shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() and shape.top is not None:
            resolved = None
            if None not in (shape.left, shape.top, shape.width, shape.height):
                resolved = _resolve_geometry_chain(
                    shape.left, shape.top, shape.width, shape.height, transform_stack
                )
            yield shape, (resolved[1] if resolved is not None else None)
        if _classify_tag(shape) == "grpSp":
            child_transform = _group_transform(shape)
            yield from _iter_text_candidates(shape.shapes, transform_stack + [child_transform])


def extract_title(slide, slide_height_emu):
    """Deterministic title/topic extraction with a bounded fallback
    chain: (1) an actual PowerPoint title placeholder; (2) the
    `Desktop:`-convention text within the header region (top
    `_TITLE_HEADER_ZONE_RATIO` of the slide, resolved to slide-absolute
    coordinates even for grouped text - L0-05); (3) any other header-
    region text, topmost first; (4) the `Desktop:` convention anywhere
    on the slide; (5) the topmost, then leftmost, non-empty text shape
    on the whole slide. Never the first shape in raw z-order, and never
    limited to top-level shapes - the real Message Actions/Pinned
    Messages slide (index 74) has a body shape ("Copy link: …") that
    iterates before its true heading ("Desktop: Message Actions"), and
    grouped title text must be found just as reliably as top-level text.

    Returns `(raw_title, basis)`; `raw_title` is `None` only if the
    slide has no text at all.
    """
    text, basis = _title_placeholder_candidate(slide)
    if text is not None:
        return text, basis

    candidates = list(_iter_text_candidates(slide.shapes, []))
    if not candidates:
        return None, "no_title_found"

    positioned = [(s, top) for s, top in candidates if top is not None]

    header_zone = slide_height_emu * _TITLE_HEADER_ZONE_RATIO
    header_candidates = sorted(
        (item for item in positioned if item[1] < header_zone),
        key=lambda item: (item[1], item[0].left or 0),
    )
    for shape, _ in header_candidates:
        text = shape.text_frame.text.strip()
        if _DESKTOP_PREFIX_RE.match(text):
            return text, "header_region_desktop_convention"
    if header_candidates:
        return header_candidates[0][0].text_frame.text.strip(), "header_region_fallback"

    for shape, _ in sorted(positioned, key=lambda item: (item[1], item[0].left or 0)):
        text = shape.text_frame.text.strip()
        if _DESKTOP_PREFIX_RE.match(text):
            return text, "desktop_convention_anywhere"

    if positioned:
        fallback_shape = sorted(positioned, key=lambda item: (item[1], item[0].left or 0))[0][0]
        return fallback_shape.text_frame.text.strip(), "position_fallback"

    # No candidate has a reliably resolvable position (e.g. every
    # text-bearing shape sits inside a group with a degenerate/missing
    # transform) - fall back to document order rather than inventing a
    # position.
    return candidates[0][0].text_frame.text.strip(), "position_fallback"


def assign_topic_occurrence_ordinals(tier_a_records):
    """Second pass, in slide order: assigns each record's `topic_key` an
    `occurrence_ordinal` (1-indexed count of that exact normalized title
    seen so far). Repeated titles (e.g. multi-slide tutorials) are
    common in this deck and slide index alone is an unreliable cross-
    deck join key (old PPTX and newer PDF are not guaranteed to have
    matching slide counts/order) - `topic_key` + `occurrence_ordinal`
    together are the intended reconciliation key. Mutates and returns
    the same list."""
    seen = Counter()
    for rec in tier_a_records:
        key = rec["topic_key"]
        seen[key] += 1
        rec["occurrence_ordinal"] = seen[key]
    return tier_a_records


# --------------------------------------------------------------------------
# Tier A - per-slide compact summary + fingerprint
# --------------------------------------------------------------------------

def compute_fingerprint(tier_a_partial):
    """Compact, deterministic, explainable structural fingerprint from
    already-computed Tier A aggregate fields. No embeddings, no LLM, no
    visual model - pure arithmetic over counts already extracted.

    `shape_count_norm` is normalized from the *real* shape count - a
    genuinely empty slide reports `0.0`, not a small nonzero value.
    Guarded (never-zero) denominators are used only where they divide,
    never as the value being normalized (L0-07)."""
    raw_count = tier_a_partial["shape_count_total"]
    denom = max(raw_count, 1)
    hist = tier_a_partial["shape_type_histogram"]
    connector_count = tier_a_partial["connector_count"]
    picture_count = tier_a_partial["picture_count"]
    return {
        "sp_rate": round(hist.get("sp", 0) / denom, 3),
        "pic_rate": round(hist.get("pic", 0) / denom, 3),
        "grp_rate": round(hist.get("grpSp", 0) / denom, 3),
        "cxn_rate": round(hist.get("cxnSp", 0) / denom, 3),
        "bound_rate": round(tier_a_partial["bound_connector_count"] / connector_count, 3) if connector_count else 0.0,
        "crop_rate": round(tier_a_partial["cropped_picture_count"] / picture_count, 3) if picture_count else 0.0,
        "freeform_rate": round(tier_a_partial["freeform_count"] / denom, 3),
        "shape_count_norm": round(min(raw_count / 60, 1.0), 3),
    }


def extract_tier_a(slide, slide_index, slide_height_emu, slide_layout_name):
    """Compact per-slide summary. Cheap by construction: no per-shape
    list is included (that is Tier B's job) - only aggregate counts, a
    histogram, the fingerprint, and the title/topic key."""
    shapes = _walk_shapes(slide.shapes, None, 0, [])
    tag_hist = Counter(s["xml_tag"] for s in shapes)
    connectors = [s for s in shapes if s["xml_tag"] == "cxnSp"]
    bound_connectors = [c for c in connectors if c["connector"]["status"] == "VALID"]
    invalid_connectors = [c for c in connectors if c["connector"]["status"] == "INVALID"]
    pictures = [s for s in shapes if s["xml_tag"] == "pic"]
    cropped_pictures = [p for p in pictures if p["picture"]["has_crop"]]
    freeforms = [s for s in shapes if s["freeform"]]
    textboxes = [s for s in shapes if s["text"] and s["text"]["run_count"] > 0]
    hidden = [s for s in shapes if s["hidden"]]
    groups = [s for s in shapes if s["xml_tag"] == "grpSp"]
    max_group_depth = max((s["group_depth"] for s in shapes), default=0)

    unsupported = [
        {"shape_id": s["shape_id"], "kind": f"{s['xml_tag']}(unrecognized/graphicFrame)"}
        for s in shapes
        if s["xml_tag"] not in ("sp", "cxnSp", "pic", "grpSp")
    ]

    raw_title, title_basis = extract_title(slide, slide_height_emu)
    topic_key = normalize_title(raw_title)

    record = {
        "slide_index": slide_index,
        "slide_number": slide_index + 1,
        "raw_title": raw_title,
        "title_extraction_basis": title_basis,
        "topic_key": topic_key,
        # `occurrence_ordinal` is filled in by `assign_topic_occurrence_ordinals`
        # after every slide's Tier A record exists - it requires whole-deck context.
        "occurrence_ordinal": None,
        "slide_layout": slide_layout_name,
        "shape_count_total": len(shapes),
        "shape_count_top_level": sum(1 for s in shapes if s["group_depth"] == 0),
        "shape_type_histogram": dict(tag_hist),
        "group_count": len(groups),
        "max_group_depth": max_group_depth,
        "picture_count": len(pictures),
        "cropped_picture_count": len(cropped_pictures),
        "connector_count": len(connectors),
        "bound_connector_count": len(bound_connectors),
        "invalid_connector_count": len(invalid_connectors),
        "freeform_count": len(freeforms),
        "textbox_count": len(textboxes),
        "hidden_shape_count": len(hidden),
        "unsupported_structures": unsupported,
    }
    record["fingerprint"] = compute_fingerprint(record)
    return record


def extract_deck_identity(prs):
    return {
        "slide_width_in": _emu_to_in(prs.slide_width),
        "slide_height_in": _emu_to_in(prs.slide_height),
        "slide_count": len(prs.slides),
        "slide_masters": [getattr(m, "name", None) for m in prs.slide_masters],
        "layouts_used": sorted({s.slide_layout.name for s in prs.slides}),
    }


# --------------------------------------------------------------------------
# range-seam proposal - workload decomposition only
# --------------------------------------------------------------------------

def _fingerprint_distance(a, b):
    return math.sqrt(sum((a[k] - b[k]) ** 2 for k in _FINGERPRINT_KEYS))


def propose_ranges(tier_a_records, target_range_count=6, min_range_size=15, merge_window=3):
    """Deterministic workload-decomposition proposal only - NEVER an
    editorial-era, editor-transition, or "current style boundary" claim.
    Every seam is scored purely from adjacent-slide fingerprint distance.

    Seam candidates within `merge_window` slides of each other are
    merged into one (a cluster of adjacent high-distance seams is almost
    always one real discontinuity, not several), and a seam is rejected
    if accepting it would create a range smaller than `min_range_size`
    on either side, unless no other seam is available (an "unavoidable"
    tiny range is recorded explicitly via `undersized`, never hidden).

    On pathological/uniform input (e.g. every slide structurally
    identical), adjacent distances may legitimately all be zero - this
    is accepted, documented behaviour for a workload-only decomposition,
    not something this function tries to "fix" with an artificial tie-
    break.
    """
    n = len(tier_a_records)
    if n == 0:
        return {"label": "WORKLOAD_DECOMPOSITION_ONLY", "ranges": [], "accepted_seams": [],
                "forced_undersized_seams": [], "all_adjacent_distances": []}

    adjacent = []
    for i in range(1, n):
        d = _fingerprint_distance(tier_a_records[i - 1]["fingerprint"], tier_a_records[i]["fingerprint"])
        adjacent.append({"before_slide_index": i, "distance": round(d, 4)})

    by_distance = sorted(adjacent, key=lambda s: -s["distance"])
    merged = []
    for cand in by_distance:
        if any(abs(cand["before_slide_index"] - m["before_slide_index"]) < merge_window for m in merged):
            continue
        merged.append(cand)

    accepted = []
    for cand in sorted(merged, key=lambda s: -s["distance"]):
        boundary = sorted([0] + [a["before_slide_index"] for a in accepted] + [n])
        idx = cand["before_slide_index"]
        pos = 0
        while pos < len(boundary) and boundary[pos] < idx:
            pos += 1
        left_size = idx - boundary[pos - 1]
        right_size = boundary[pos] - idx
        if left_size < min_range_size or right_size < min_range_size:
            continue
        accepted.append(cand)
        if len(accepted) >= target_range_count - 1:
            break

    forced = []
    if not accepted and n > min_range_size:
        best = merged[0] if merged else None
        if best is not None:
            forced.append(best)

    final_seams = sorted(accepted + forced, key=lambda s: s["before_slide_index"])
    boundaries = [0] + [s["before_slide_index"] for s in final_seams] + [n]
    ranges = []
    for i in range(len(boundaries) - 1):
        start, end = boundaries[i], boundaries[i + 1] - 1
        ranges.append({
            "start_slide_index": start,
            "end_slide_index": end,
            "slide_count": end - start + 1,
            "undersized": (end - start + 1) < min_range_size,
        })

    return {
        "label": "WORKLOAD_DECOMPOSITION_ONLY",
        "note": "Structural seams are workload-decomposition candidates only, "
                "never editorial-era, editor-transition, or current-style boundaries.",
        "ranges": ranges,
        "accepted_seams": final_seams,
        "forced_undersized_seams": forced,
        "all_adjacent_distances": adjacent,
    }


# --------------------------------------------------------------------------
# structural candidate clustering - explicitly non-editorial
# --------------------------------------------------------------------------

def cluster_structural_candidates(tier_a_records, k=5, seed=42, iterations=30):
    """Deterministic k-means (fixed seed, pure stdlib) over the
    structural fingerprint only. Every output MUST be read as
    `STRUCTURAL CANDIDATE CLUSTER` - a reproducible grouping by shape-
    level statistics, never an editorial-style or editing-era
    conclusion. No embeddings, no LLM."""
    n = len(tier_a_records)
    if n < k:
        return {"label": "STRUCTURAL_CANDIDATE_CLUSTER", "clusters": [], "outliers_top20": []}

    rng = random.Random(seed)
    vectors = [[r["fingerprint"][key] for key in _FINGERPRINT_KEYS] for r in tier_a_records]
    centers = rng.sample(vectors, k)

    assign = [0] * n
    for _ in range(iterations):
        for i, v in enumerate(vectors):
            dists = [sum((v[j] - c[j]) ** 2 for j in range(len(v))) for c in centers]
            assign[i] = dists.index(min(dists))
        new_centers = []
        for cluster_idx in range(k):
            members = [vectors[i] for i in range(n) if assign[i] == cluster_idx]
            if not members:
                new_centers.append(centers[cluster_idx])
                continue
            new_centers.append([sum(m[j] for m in members) / len(members) for j in range(len(members[0]))])
        centers = new_centers

    clusters = []
    for cluster_idx in range(k):
        members = [i for i in range(n) if assign[i] == cluster_idx]
        if not members:
            continue
        centroid = {key: round(centers[cluster_idx][j], 3) for j, key in enumerate(_FINGERPRINT_KEYS)}
        clusters.append({
            "label": "STRUCTURAL_CANDIDATE_CLUSTER",
            "cluster_id": cluster_idx,
            "member_slide_indices": members,
            "centroid_fingerprint": centroid,
        })

    outliers = []
    for i, v in enumerate(vectors):
        c = centers[assign[i]]
        dist = math.sqrt(sum((v[j] - c[j]) ** 2 for j in range(len(v))))
        outliers.append({"slide_index": i, "distance_from_own_centroid": round(dist, 4), "cluster_id": assign[i]})
    outliers.sort(key=lambda o: -o["distance_from_own_centroid"])

    return {
        "label": "STRUCTURAL_CANDIDATE_CLUSTER",
        "note": "Feature-based structural groupings only, never editorial-style "
                "or editing-era conclusions.",
        "k": k,
        "seed": seed,
        "clusters": clusters,
        "outliers_top20": outliers[:20],
    }


# --------------------------------------------------------------------------
# source integrity + deterministic identifiers
# --------------------------------------------------------------------------

def _hash_source_or_raise(path):
    """Typed-error-preserving hash. `file_sha256` alone raises a bare
    `FileNotFoundError`/`OSError` - normalized here to `DeckSourceError`
    so both the pre-extraction and pre-publication hash calls (L0-01)
    fail the same typed way `load_deck()` already does."""
    try:
        return file_sha256(path)
    except FileNotFoundError as exc:
        raise DeckSourceError("PPTX not found: %s" % path) from exc
    except OSError as exc:
        raise DeckSourceError("Could not read source for hashing: %s (%s)" % (path, exc)) from exc


def _normalize_source_identifier(source_deck_path):
    """A stable logical identifier for the source, independent of
    whether the caller spelled the path relatively or absolutely
    (L0-02): the resolved real path, expressed relative to the
    repository root when the source lives inside it (the canonical
    deck always does), else the resolved absolute path as a last
    resort. Two different spellings of the same physical file always
    resolve to the identical string."""
    resolved = Path(source_deck_path).resolve()
    try:
        return resolved.relative_to(_REPO_ROOT).as_posix()
    except ValueError:
        return resolved.as_posix()


# --------------------------------------------------------------------------
# staged build + validation + atomic, owned-tree publication (L0-03)
# --------------------------------------------------------------------------

def _mkdir_or_raise(path, what):
    """`Path.mkdir(parents=True, exist_ok=True)` still raises a raw
    `FileExistsError` when the path exists but is not a directory (e.g.
    the caller passed an `output_dir` that is actually a regular file) -
    normalized here to a typed `PublicationError` so no raw filesystem
    exception escapes the publication path."""
    try:
        path.mkdir(parents=True, exist_ok=True)
    except FileExistsError as exc:
        raise PublicationError(
            "Could not create %s at %s - a filesystem entry already exists "
            "there and is not a directory: %s" % (what, path, exc)
        ) from exc
    except OSError as exc:
        raise PublicationError("Could not create %s at %s: %s" % (what, path, exc)) from exc


@contextlib.contextmanager
def _acquire_publish_lock(output_dir):
    """Narrow, single-host, single-writer protection for one Layer 0
    `output_dir`: an exclusive, non-blocking `fcntl.flock` on a dedicated
    `.{name}.lock` file next to `output_dir` (never inside it - it is not
    one of Layer 0's owned tree artifacts, and its own name pattern
    `.lock` is distinct from `.staging-`/`.prev-`, so it can never
    collide with them). Held for the entire staging+validate+publish
    sequence, so at most one generation is ever actively producing output
    for a given tree at a time - a second concurrent call for the same
    `output_dir` fails fast with `PublicationLockError` rather than
    waiting or interleaving.

    The lock file itself is intentionally never deleted (deleting and
    recreating lock files is a well-known source of races: a second
    process could open-then-stat a file another process is about to
    unlink and replace, sees a *different* inode than the one it just
    locked, and the two "holders" then genuinely stop excluding each
    other) - only the flock itself is released. This has no classic
    stale-PID-file problem: `flock` is kernel-managed and is
    automatically released the moment the holding process exits for any
    reason, including a crash, so a dead writer can never permanently
    wedge this lock. POSIX-only (`fcntl`); provides no distributed or
    network-filesystem guarantee - see `PublicationLockError`.
    """
    _mkdir_or_raise(output_dir.parent, "output_dir's parent directory")
    lock_path = output_dir.parent / f".{output_dir.name}.lock"
    try:
        fd = os.open(str(lock_path), os.O_CREAT | os.O_RDWR)
    except OSError as exc:
        raise PublicationError(
            "Could not open Layer 0 publish lock file %s: %s" % (lock_path, exc)
        ) from exc
    try:
        try:
            fcntl.flock(fd, fcntl.LOCK_EX | fcntl.LOCK_NB)
        except OSError as exc:
            raise PublicationLockError(
                "Another Layer 0 generation currently owns the publish lock "
                "for %s (lock file: %s): %s" % (output_dir, lock_path, exc)
            ) from exc
        yield
    finally:
        try:
            fcntl.flock(fd, fcntl.LOCK_UN)
        finally:
            os.close(fd)


def _write_staged_tree(staging_dir, manifest, tier_a_records, tier_b_records, range_proposal, cluster_result):
    """Writes the complete staged tree. Every filesystem operation here
    is wrapped in one boundary: any `OSError` (JSON write, directory
    creation, an individual Tier B file write, ...) is normalized to a
    typed `PublicationStagingError` - no raw `OSError` escapes for an
    expected staging-time filesystem failure. `output_dir` is never
    touched by this function; the caller's `finally` block removes
    whatever partial content ends up in `staging_dir` regardless of
    success or failure here."""
    try:
        tier_b_dir = staging_dir / _TIER_B_DIRNAME
        tier_b_dir.mkdir(parents=True, exist_ok=True)
        (staging_dir / "tier_a_summary.json").write_text(
            json.dumps(tier_a_records, indent=2), encoding="utf-8"
        )
        for rec in tier_b_records:
            path = tier_b_dir / f"slide_{rec['slide_index']:04d}.json"
            path.write_text(json.dumps(rec, indent=2), encoding="utf-8")
        (staging_dir / "structural_analysis.json").write_text(
            json.dumps({"range_proposal": range_proposal, "structural_clusters": cluster_result}, indent=2),
            encoding="utf-8",
        )
        # Written last within staging so a mid-staging failure (e.g. disk
        # full while writing a Tier B file) never leaves a staged tree
        # whose manifest looks complete - `_validate_staged_tree` would
        # still catch a missing manifest, but this ordering makes the
        # common failure mode obviously incomplete rather than
        # misleadingly whole.
        (staging_dir / "manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    except OSError as exc:
        raise PublicationStagingError(
            "Could not write the staged Layer 0 tree at %s: %s" % (staging_dir, exc)
        ) from exc


def _validate_staged_tree(staging_dir, manifest):
    """Validates the staged tree BEFORE any publication is attempted -
    nothing in `output_dir` is touched by this function or by anything
    that runs before it succeeds."""
    errors = []
    manifest_path = staging_dir / "manifest.json"
    tier_a_path = staging_dir / "tier_a_summary.json"
    structural_path = staging_dir / "structural_analysis.json"
    tier_b_dir = staging_dir / _TIER_B_DIRNAME

    for label, path in (
        ("manifest", manifest_path),
        ("tier_a_summary", tier_a_path),
        ("structural_analysis", structural_path),
    ):
        if not path.exists():
            errors.append(f"{label} missing from staged tree: {path}")
            continue
        try:
            json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            errors.append(f"{label} is not valid JSON: {exc}")

    if not tier_b_dir.is_dir():
        errors.append(f"tier_b directory missing from staged tree: {tier_b_dir}")
    else:
        staged_tier_b_files = sorted(p.name for p in tier_b_dir.glob("slide_*.json"))
        expected = sorted(Path(p).name for p in manifest["tier_b_paths"])
        if staged_tier_b_files != expected:
            errors.append(
                "staged tier_b file set does not match manifest.tier_b_paths "
                f"(staged={staged_tier_b_files!r} expected={expected!r})"
            )
        for name in staged_tier_b_files:
            try:
                json.loads((tier_b_dir / name).read_text(encoding="utf-8"))
            except json.JSONDecodeError as exc:
                errors.append(f"tier_b/{name} is not valid JSON: {exc}")

    if not errors:
        try:
            tier_a_records = json.loads(tier_a_path.read_text(encoding="utf-8"))
            if len(tier_a_records) != manifest["slide_count"]:
                errors.append(
                    f"tier_a_summary has {len(tier_a_records)} records, "
                    f"manifest declares slide_count={manifest['slide_count']}"
                )
        except Exception as exc:
            errors.append(f"could not cross-check tier_a_summary against manifest: {exc}")

    if errors:
        raise StagingValidationError(
            "Layer 0 staged tree failed validation - nothing was published: " + "; ".join(errors)
        )


_MARKER_SCHEMA_VERSION = "1.0.0"


def _marker_path(output_dir):
    return Path(output_dir) / _COMMIT_MARKER_NAME


def _compute_owned_tree_digest(output_dir):
    """The whole-owned-tree commitment: a single SHA-256 over the sorted
    list of `(relative_path, file_sha256)` pairs for every owned
    canonical artifact - `manifest.json`, `tier_a_summary.json`,
    `structural_analysis.json`, and every `tier_b/slide_*.json` file
    *the manifest itself declares* via `tier_b_paths`. Deterministic
    regardless of filesystem traversal order (sorted explicitly).

    This binds the marker to the COMPLETE owned set, not just the
    manifest: mutating, removing, or adding any owned file - including a
    missing or an unexpected extra `tier_b/slide_*.json` - changes or
    invalidates this digest. The declared (`manifest["tier_b_paths"]`)
    and actual (glob of `tier_b/slide_*.json` on disk) Tier B file sets
    must match exactly, or this returns "not computable" - a stray extra
    file is exactly as disqualifying as a missing expected one.

    Returns `(digest_hex, file_count)`, or `(None, None)` if the digest
    is not computable for any reason (missing/unreadable/malformed
    manifest, a missing owned file, a Tier B set mismatch) - never
    raises. `(None, None)` is always treated by callers as "does not
    match," the safe direction.
    """
    output_dir = Path(output_dir)
    manifest_path = output_dir / "manifest.json"
    # Reading and parsing are two separate `try` blocks, each scoped to
    # exactly the errors that expected malformed input can produce:
    # invalid UTF-8 bytes (`UnicodeDecodeError`, NOT an `OSError`
    # subclass) or a missing file (`OSError`) on read; a JSON syntax
    # error or an oversized integer literal (Python's interpreter-level
    # int-to-str conversion limit raises a plain `ValueError`, not
    # `json.JSONDecodeError`, for that specific case) on parse. Both
    # fail closed - never escape as a raw decoding/parsing exception.
    try:
        manifest_text = manifest_path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError):
        return None, None
    try:
        manifest = json.loads(manifest_text)
    except ValueError:
        return None, None
    if not isinstance(manifest, dict):
        return None, None

    declared_tier_b = manifest.get("tier_b_paths")
    if not isinstance(declared_tier_b, list):
        return None, None
    try:
        expected_tier_b = sorted(Path(p).as_posix() for p in declared_tier_b)
    except TypeError:
        return None, None

    tier_b_dir = output_dir / _TIER_B_DIRNAME
    if not tier_b_dir.is_dir():
        return None, None
    actual_tier_b = sorted(f"{_TIER_B_DIRNAME}/{p.name}" for p in tier_b_dir.glob("slide_*.json"))
    if actual_tier_b != expected_tier_b:
        # A missing expected file or an extra stale owned file both land
        # here - either way the owned set does not match what this
        # publication is supposed to contain.
        return None, None

    expected_relative = sorted(list(_OWNED_TOP_LEVEL_FILES) + expected_tier_b)
    try:
        hasher = hashlib.sha256()
        for rel in expected_relative:
            file_bytes = (output_dir / rel).read_bytes()
            hasher.update(rel.encode("utf-8") + b"\x00")
            hasher.update(hashlib.sha256(file_bytes).digest() + b"\x00")
        return hasher.hexdigest(), len(expected_relative)
    except OSError:
        return None, None


def _write_commit_marker(output_dir):
    """Writes/replaces the commit-validity marker to record the current
    whole-owned-tree digest of `output_dir` - this is the single source
    of truth `is_tree_valid()` checks. Must only ever be called once the
    complete owned set is known to hold the intended, fully-committed
    content (either a fresh successful publish, or a fully-restored
    rollback). Raises `PublicationError` (never silently no-ops or
    writes a marker for content it could not verify) if the digest
    cannot be computed - the caller's own `try`/`except` machinery
    already treats that identically to any other commit-phase failure."""
    output_dir = Path(output_dir)
    digest, count = _compute_owned_tree_digest(output_dir)
    if digest is None:
        raise PublicationError(
            "Could not compute the owned-tree digest for %s while finalizing "
            "publication - the owned artifact set is not internally "
            "consistent (a missing/mismatched owned file)." % output_dir
        )
    marker = {
        "schema_version": _MARKER_SCHEMA_VERSION,
        "tree_sha256": digest,
        "owned_file_count": count,
    }
    _marker_path(output_dir).write_text(json.dumps(marker), encoding="utf-8")


def _invalidate_commit_marker_or_raise(output_dir):
    """HARD precondition for any canonical mutation (see the module's
    publication state machine): removes the existing commit-validity
    marker, if present, BEFORE any owned artifact is moved or replaced.
    A missing marker (first-ever publish, or already invalidated) is a
    no-op success, not an error.

    If the marker EXISTS and cannot be removed, this raises a typed
    `PublicationError` immediately - publication must abort before
    touching a single canonical artifact, so the previous tree is left
    completely untouched and remains valid. This is deliberately no
    longer a best-effort/silent-failure step (the prior version's
    silent-failure behavior is exactly what let a stale marker survive
    into a publication transition and later mismatch a mixed tree)."""
    marker_path = _marker_path(output_dir)
    try:
        if marker_path.exists():
            marker_path.unlink()
    except OSError as exc:
        raise PublicationError(
            "Could not invalidate the existing commit-validity marker at %s "
            "before publication - aborting before any canonical artifact was "
            "touched. The previous tree at %s remains valid and untouched: %s"
            % (marker_path, output_dir, exc)
        ) from exc


def is_tree_valid(output_dir):
    """Whether a published Layer 0 tree at `output_dir` can currently be
    trusted: `True` only if the commit-validity marker exists, has the
    expected shape, declares the exact currently-supported marker
    `schema_version` (`_MARKER_SCHEMA_VERSION` - an unknown/future/past
    version is never silently accepted), and its recorded whole-owned-
    tree digest matches a freshly recomputed digest of everything
    actually on disk right now (manifest, Tier A, structural analysis,
    and every declared Tier B file - see `_compute_owned_tree_digest`).
    This authenticates the COMPLETE owned set, not just the manifest -
    mutating, removing, or adding any single owned file (including a
    stray extra `tier_b/slide_*.json`) makes this `False`.

    Fails closed for every malformed/unexpected input this is a *safe
    predicate*, not a validating parser, so none of the following ever
    raise - they all simply return `False`: a missing marker; invalid
    UTF-8 bytes in the marker or the manifest; a JSON syntax error; an
    oversized integer JSON token (Python's interpreter-level int-to-str
    conversion limit raises a plain `ValueError` for this specific case,
    not `json.JSONDecodeError`); the wrong JSON shape (e.g. a bare
    list/scalar); a missing, `null`, or wrong-typed `schema_version`/
    `tree_sha256`/`owned_file_count` field; an unsupported
    `schema_version` value; a stale marker copied from a different tree;
    or any owned-file mismatch."""
    output_dir = Path(output_dir)
    marker_path = _marker_path(output_dir)
    if not marker_path.exists():
        return False
    try:
        marker_text = marker_path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError):
        return False
    try:
        marker = json.loads(marker_text)
    except ValueError:
        return False
    if not isinstance(marker, dict):
        return False
    if marker.get("schema_version") != _MARKER_SCHEMA_VERSION:
        return False
    recorded_digest = marker.get("tree_sha256")
    recorded_count = marker.get("owned_file_count")
    if not isinstance(recorded_digest, str) or not isinstance(recorded_count, int):
        return False
    actual_digest, actual_count = _compute_owned_tree_digest(output_dir)
    if actual_digest is None:
        return False
    return actual_digest == recorded_digest and actual_count == recorded_count


def _publish_owned_tree(staging_dir, output_dir):
    """Publishes the already-validated staged tree as one owned Layer 0
    tree (L0-03). Layer 0 owns exactly `_OWNED_TOP_LEVEL_FILES`, the
    entire `tier_b/` subdirectory, and the commit-validity marker
    (`_COMMIT_MARKER_NAME`) within `output_dir` - nothing else under or
    near `output_dir` is ever read, moved, or removed by this function,
    so unrelated files are never touched. Callers must hold
    `_acquire_publish_lock(output_dir)` around this call.

    Two-phase, rename-based publish: existing owned items are first
    moved aside (fast, atomic `os.replace`, not deleted yet), then the
    staged replacements are moved into place, then the commit-validity
    marker is written to match the now-published `manifest.json` - this
    marker write is the true "publish is complete" signal, checked via
    `is_tree_valid()`, not merely "the files happen to be present." Only
    once every new item AND the marker are successfully in place are the
    now-stale "moved aside" originals removed (this is what makes a
    shorter regeneration's stale extra Tier B files disappear - the
    entire old `tier_b/` is discarded as one unit, never diffed
    file-by-file).

    Publication-validity invariant (the specific gap this closes):
    consumers must never be able to mistake a mixed/inconsistent
    canonical tree for a valid one. Per the module's publication state
    machine, marker invalidation is now a HARD PRECONDITION performed
    BEFORE a single owned artifact is moved or replaced - not a
    best-effort corrective action taken after a failure is already
    detected. If the marker cannot be invalidated, this function raises
    immediately, before the prepare phase even begins: zero canonical
    mutation has happened, so no rollback is needed and the previous
    tree is left byte-identical and still valid (its own,
    already-invalidated-and-not-yet-replaced marker is simply gone,
    which is why invalidation happening first, not merely early, is what
    actually closes the gap - see `_invalidate_commit_marker_or_raise`).

    Once invalidation succeeds, there is no possible window - not during
    the prepare-phase move-asides, not during the commit-phase
    replacements, not for same-source regenerations with byte-identical
    content - in which `is_tree_valid()` can return `True`, because the
    marker simply does not exist again until the very last successful
    step. If a commit-phase failure occurs, rollback proceeds: each
    moved-aside item is restored independently (one item's restore
    failing does not abort restoring the rest). If every restore
    succeeds, the marker is re-established from a freshly recomputed
    whole-owned-tree digest of the now-fully-restored previous tree
    (State A: prior valid tree restored, `is_tree_valid()` -> `True`
    again), and `PublicationError` is raised describing the original
    commit failure. If any restore also fails, the marker is
    deliberately left absent (State B: canonical tree explicitly
    invalid/unavailable, `is_tree_valid()` -> `False`, even though some
    individual files may still be readable and even self-consistent-
    looking in isolation) and `PublicationRollbackError` is raised - the
    previous tree's content is NOT lost (each affected item's
    `.prev-<token>` backup is deliberately left on disk, named explicitly
    in the error), but must be manually restored before `output_dir`
    should be trusted again.
    """
    _mkdir_or_raise(output_dir, "output_dir")

    # HARD PRECONDITION: invalidate the existing marker (if any) before
    # touching a single owned artifact. A failure here aborts the whole
    # call with zero canonical mutation - see `_invalidate_commit_marker_or_raise`.
    _invalidate_commit_marker_or_raise(output_dir)

    token = f"{os.getpid()}-{uuid.uuid4().hex[:8]}"

    prepared = []  # (final_path, prev_path_or_None), in prepare order
    try:
        for name in _OWNED_TOP_LEVEL_FILES:
            final_path = output_dir / name
            if final_path.exists():
                prev_path = output_dir / f"{name}.prev-{token}"
                os.replace(final_path, prev_path)
                prepared.append((final_path, prev_path))
            else:
                prepared.append((final_path, None))

        tier_b_final = output_dir / _TIER_B_DIRNAME
        if tier_b_final.exists():
            tier_b_prev = output_dir / f"{_TIER_B_DIRNAME}.prev-{token}"
            os.replace(tier_b_final, tier_b_prev)
            prepared.append((tier_b_final, tier_b_prev))
        else:
            prepared.append((tier_b_final, None))

        for name in _OWNED_TOP_LEVEL_FILES:
            os.replace(staging_dir / name, output_dir / name)
        os.replace(staging_dir / _TIER_B_DIRNAME, tier_b_final)
        # The marker write is the true "publish complete" signal - it is
        # deliberately still inside this `try`, so a failure here (e.g.
        # disk full on this very small file) is treated exactly like any
        # other commit-phase failure and triggers the same rollback path
        # below, rather than leaving a fully-committed-but-unmarked tree.
        _write_commit_marker(output_dir)
    except Exception as commit_exc:
        # The marker was already invalidated as a hard precondition
        # before this `try` block even started (see above) - there is
        # nothing further to invalidate here; proceed straight to
        # rollback.
        rollback_errors = []
        for final_path, prev_path in prepared:
            try:
                if final_path.exists():
                    if final_path.is_dir():
                        shutil.rmtree(final_path)
                    else:
                        final_path.unlink()
                if prev_path is not None and prev_path.exists():
                    os.replace(prev_path, final_path)
            except OSError as rollback_exc:
                rollback_errors.append(
                    "%s (previous content retained, NOT deleted, at %s): %s"
                    % (final_path, prev_path, rollback_exc)
                )

        if rollback_errors:
            # State B: rollback could not fully restore the previous
            # tree. The marker stays absent (never re-established here)
            # so `is_tree_valid()` truthfully reports the canonical tree
            # as unavailable, even though some individual files may
            # still be present/readable/self-consistent-looking on their
            # own - never silently exposed as a valid publish.
            raise PublicationRollbackError(
                "Layer 0 publication to %s failed (%s) AND automatic rollback "
                "could not fully restore the previous tree. The previous, "
                "still-valid artifacts were retained under their "
                "'<name>.prev-%s' backup names and were NOT deleted - "
                "restore each one to its canonical name before trusting "
                "%s again. The canonical tree's commit-validity marker has "
                "been left absent, so is_tree_valid(%s) correctly reports "
                "False until that manual recovery happens. Rollback "
                "failures: %s"
                % (output_dir, commit_exc, token, output_dir, output_dir, "; ".join(rollback_errors))
            ) from commit_exc

        # State A: rollback fully succeeded - the previous owned tree is
        # genuinely back in place, so the marker is re-established from a
        # freshly recomputed whole-tree digest of it (best-effort; if
        # this particular write - or the digest computation itself -
        # fails, the marker simply stays absent, which is the safe
        # direction - `is_tree_valid()` would then under-report a tree
        # that is actually fine content-wise, never over-report one that
        # isn't). `_write_commit_marker` can raise either `OSError` (the
        # write itself) or `PublicationError` (digest not computable) -
        # both are equally non-fatal here.
        try:
            if (output_dir / "manifest.json").exists():
                _write_commit_marker(output_dir)
        except (OSError, PublicationError):
            pass

        raise PublicationError(
            "Layer 0 publication failed and was rolled back to the previous "
            "tree at %s: %s" % (output_dir, commit_exc)
        ) from commit_exc

    cleanup_errors = []
    for _, prev_path in prepared:
        if prev_path is not None and prev_path.exists():
            try:
                if prev_path.is_dir():
                    shutil.rmtree(prev_path)
                else:
                    prev_path.unlink()
            except OSError as exc:
                cleanup_errors.append(str(exc))
    if cleanup_errors:
        raise PublicationError(
            "Layer 0 publication succeeded but cleanup of stale previous "
            "artifacts failed (the new tree is live and correct): %s"
            % "; ".join(cleanup_errors)
        )


# --------------------------------------------------------------------------
# production entry point
# --------------------------------------------------------------------------

def generate_layer0_inventory(source_deck_path, output_dir):
    """Single production entry point. Read-only against
    `source_deck_path`.

    Source-integrity transaction (L0-01): the source is hashed BEFORE it
    is opened/parsed at all, then opened via the Safe PPT Engine's own
    typed `load_deck()` (reused unmodified), then re-hashed immediately
    before anything is written to disk. If the two hashes differ -
    including the case where the source was replaced between the first
    hash and `load_deck()`, which would otherwise let the manifest
    record a hash for bytes different from what was actually parsed -
    `SourceIntegrityError` is raised and nothing is staged or published.

    Atomic, owned-tree publication (L0-03): the complete output is built
    into a private staging directory, validated as a whole
    (`StagingValidationError` on failure - `output_dir` is never
    touched), and only then published as one atomically-swapped tree
    that Layer 0 exclusively owns within `output_dir`, under an
    exclusive single-writer lock (`PublicationLockError` if another
    generation currently holds it) so two concurrent generations for the
    same `output_dir` can never interleave (`PublicationError`/
    `PublicationRollbackError` on a publish failure - see
    `_publish_owned_tree` and `_acquire_publish_lock`). `output_dir`
    existing as a non-directory is rejected up front with a typed
    `PublicationError`, before the source is even hashed.
    """
    source_deck_path = str(source_deck_path)
    output_dir = Path(output_dir)

    if output_dir.exists() and not output_dir.is_dir():
        raise PublicationError(
            "output_dir %s already exists and is not a directory - refusing "
            "to publish a Layer 0 inventory there." % output_dir
        )

    canonical_source_hash = _hash_source_or_raise(source_deck_path)
    prs = load_deck(source_deck_path)

    deck_identity = extract_deck_identity(prs)
    tier_a_records = []
    tier_b_records = []
    for slide_index, slide in enumerate(prs.slides):
        layout_name = None
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass
        tier_a_records.append(
            extract_tier_a(slide, slide_index, prs.slide_height, layout_name)
        )
        tier_b_records.append(extract_tier_b(slide, slide_index))
    assign_topic_occurrence_ordinals(tier_a_records)

    range_proposal = propose_ranges(tier_a_records)
    cluster_result = cluster_structural_candidates(tier_a_records)

    final_source_hash = _hash_source_or_raise(source_deck_path)
    if final_source_hash != canonical_source_hash:
        raise SourceIntegrityError(
            "Source deck changed during Layer 0 extraction (initial=%s final=%s) "
            "- refusing to publish an inventory that may not match the bytes it "
            "was extracted from." % (canonical_source_hash, final_source_hash)
        )

    manifest = {
        "schema_version": SCHEMA_VERSION,
        "source_identifier": _normalize_source_identifier(source_deck_path),
        "source_sha256": canonical_source_hash,
        "slide_count": len(tier_a_records),
        "deck_identity": deck_identity,
        "tier_a_path": "tier_a_summary.json",
        "tier_b_paths": [f"{_TIER_B_DIRNAME}/slide_{r['slide_index']:04d}.json" for r in tier_a_records],
        "structural_analysis_path": "structural_analysis.json",
        "validation_status": "extracted_not_founder_reviewed",
    }

    with _acquire_publish_lock(output_dir):
        token = f"{os.getpid()}-{uuid.uuid4().hex[:8]}"
        staging_dir = output_dir.parent / f".{output_dir.name}.staging-{token}"
        try:
            staging_dir.mkdir(parents=True, exist_ok=False)
        except OSError as exc:
            raise PublicationStagingError(
                "Could not create Layer 0 staging directory %s: %s" % (staging_dir, exc)
            ) from exc
        try:
            _write_staged_tree(staging_dir, manifest, tier_a_records, tier_b_records, range_proposal, cluster_result)
            _validate_staged_tree(staging_dir, manifest)
            _publish_owned_tree(staging_dir, output_dir)
        finally:
            if staging_dir.exists():
                shutil.rmtree(staging_dir, ignore_errors=True)

    return manifest


if __name__ == "__main__":  # pragma: no cover - thin CLI wrapper
    import sys

    if len(sys.argv) != 3:
        print("usage: python -m documentation_intelligence.layer0 <source.pptx> <output_dir>", file=sys.stderr)
        raise SystemExit(2)
    result = generate_layer0_inventory(sys.argv[1], sys.argv[2])
    print(json.dumps(result, indent=2))
