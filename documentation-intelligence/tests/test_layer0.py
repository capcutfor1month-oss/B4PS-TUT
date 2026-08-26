"""Documentation Intelligence Layer 0 tests: Deterministic MasterSlide
Inventory. Fixtures are small synthetic `.pptx` files built with
python-pptx (raw-XML shape injection where python-pptx's own high-level
API has no builder, e.g. connectors/freeforms/groups), plus a real-deck
regression against the actual MasterSlide archive when present in the
working tree (guarded/skipped when it is not, matching
`tests/test_locate.py`'s own convention)."""

from __future__ import annotations

import fcntl
import hashlib
import json
import multiprocessing
import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches

from documentation_intelligence._safe_ppt_engine_import import DeckSourceError
from documentation_intelligence import layer0
from documentation_intelligence.layer0 import (
    SCHEMA_VERSION,
    PublicationError,
    PublicationLockError,
    PublicationRollbackError,
    PublicationStagingError,
    SourceIntegrityError,
    StagingValidationError,
    assign_topic_occurrence_ordinals,
    cluster_structural_candidates,
    compute_fingerprint,
    extract_deck_identity,
    extract_tier_a,
    extract_tier_b,
    extract_title,
    generate_layer0_inventory,
    is_tree_valid,
    normalize_title,
    propose_ranges,
)

_REPO_ROOT = Path(__file__).resolve().parents[2]
_REAL_DECK = (
    _REPO_ROOT
    / "documentation-artifacts"
    / "masterslide"
    / "old"
    / "MASTER Complete Bridge4PS Desktop-Browser Feature Tutorials.pptx"
)
_REAL_DECK_SHA256 = "e461613baae2874eaeede3268fff1aee081c33a04a6bcbb0c7769b94bab16834"


def _file_hash(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()


def _bound_connector_xml(shape_id, target_id, target_idx, with_end=False, end_id=None, end_idx="0"):
    if with_end:
        actual_end_id = end_id if end_id is not None else target_id + 1
        end = f'<a:endCxn id="{actual_end_id}" idx="{end_idx}"/>'
    else:
        end = ""
    return f"""<p:cxnSp {nsdecls("p", "a")}>
  <p:nvCxnSpPr>
    <p:cNvPr id="{shape_id}" name="Connector {shape_id}"/>
    <p:cNvCxnSpPr>
      <a:stCxn id="{target_id}" idx="{target_idx}"/>
      {end}
    </p:cNvCxnSpPr>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm><a:off x="0" y="0"/><a:ext cx="100000" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill></a:ln>
  </p:spPr>
</p:cxnSp>"""


def _custom_connector_xml(shape_id, st_cxn_attrs=None, end_cxn_attrs=None):
    """`st_cxn_attrs`/`end_cxn_attrs`: dict of XML attributes for the
    `<a:stCxn>`/`<a:endCxn>` element, or `None` to omit the element
    entirely (an absent endpoint)."""
    def _el(tag, attrs):
        if attrs is None:
            return ""
        attr_str = " ".join(f'{k}="{v}"' for k, v in attrs.items())
        return f"<a:{tag} {attr_str}/>"

    return f"""<p:cxnSp {nsdecls("p", "a")}>
  <p:nvCxnSpPr>
    <p:cNvPr id="{shape_id}" name="Connector {shape_id}"/>
    <p:cNvCxnSpPr>
      {_el("stCxn", st_cxn_attrs)}
      {_el("endCxn", end_cxn_attrs)}
    </p:cNvCxnSpPr>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm><a:off x="0" y="0"/><a:ext cx="100000" cy="0"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:cxnSp>"""


def _unbound_connector_xml(shape_id):
    return _custom_connector_xml(shape_id, None, None)


def _freeform_xml(shape_id, n_lnto=2, n_cubic=0):
    lnto = "".join('<a:lnTo><a:pt x="100" y="100"/></a:lnTo>' for _ in range(n_lnto))
    cubic = "".join(
        '<a:cubicBezTo><a:pt x="1" y="1"/><a:pt x="2" y="2"/><a:pt x="3" y="3"/></a:cubicBezTo>'
        for _ in range(n_cubic)
    )
    return f"""<p:sp {nsdecls("p", "a")}>
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="Freeform {shape_id}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="0" y="0"/><a:ext cx="200000" cy="200000"/></a:xfrm>
    <a:custGeom>
      <a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
      <a:rect l="0" t="0" r="0" b="0"/>
      <a:pathLst>
        <a:path w="200000" h="200000">
          <a:moveTo><a:pt x="0" y="0"/></a:moveTo>
          {lnto}
          {cubic}
        </a:path>
      </a:pathLst>
    </a:custGeom>
  </p:spPr>
</p:sp>"""


def _hidden_textbox_xml(shape_id, text="hidden shape"):
    return f"""<p:sp {nsdecls("p", "a")}>
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="Hidden {shape_id}" hidden="1"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="0" y="0"/><a:ext cx="500000" cy="200000"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/><a:lstStyle/>
    <a:p><a:r><a:t>{text}</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""


def _textbox_xml(shape_id, text, left=0, top=0, width=1000000, height=300000):
    return f"""<p:sp {nsdecls("p", "a")}>
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="TextBox {shape_id}"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/><a:lstStyle/>
    <a:p><a:r><a:t>{text}</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""


def _group_xml(
    group_id, child_xml_list, off=(0, 0), ext=(1000000, 1000000), ch_off=None, ch_ext=None,
    rot=None, flip_h=None, flip_v=None,
):
    ch_off = ch_off if ch_off is not None else off
    ch_ext = ch_ext if ch_ext is not None else ext
    children = "\n".join(child_xml_list)
    xfrm_attrs = ""
    if rot is not None:
        xfrm_attrs += f' rot="{rot}"'
    if flip_h is not None:
        xfrm_attrs += f' flipH="{flip_h}"'
    if flip_v is not None:
        xfrm_attrs += f' flipV="{flip_v}"'
    return f"""<p:grpSp {nsdecls("p", "a")}>
  <p:nvGrpSpPr>
    <p:cNvPr id="{group_id}" name="Group {group_id}"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm{xfrm_attrs}>
      <a:off x="{off[0]}" y="{off[1]}"/><a:ext cx="{ext[0]}" cy="{ext[1]}"/>
      <a:chOff x="{ch_off[0]}" y="{ch_off[1]}"/><a:chExt cx="{ch_ext[0]}" cy="{ch_ext[1]}"/>
    </a:xfrm>
  </p:grpSpPr>
  {children}
</p:grpSp>"""


def _blank_slide(prs=None):
    prs = prs or Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _append(slide, xml):
    el = parse_xml(xml)
    slide.shapes._spTree.append(el)
    return el


def _tier_a_for(slide, slide_index=0, prs=None):
    height = prs.slide_height if prs is not None else Emu(6858000)
    layout_name = slide.slide_layout.name
    return extract_tier_a(slide, slide_index, height, layout_name)


def _hold_lock_worker(lock_path, ready_event, release_event):
    """Module-level (picklable) worker for
    `TestSingleWriterLock.test_concurrent_process_holding_lock_blocks_a_second_writer`:
    acquires the same exclusive flock `_acquire_publish_lock` would, in a
    genuinely separate process, signals readiness, then holds it until
    told to release."""
    fd = os.open(lock_path, os.O_CREAT | os.O_RDWR)
    fcntl.flock(fd, fcntl.LOCK_EX)
    ready_event.set()
    release_event.wait(timeout=10)
    fcntl.flock(fd, fcntl.LOCK_UN)
    os.close(fd)


class TestConnectorBindingDescendantSearch:
    """`stCxn`/`endCxn` live under `p:nvCxnSpPr/p:cNvCxnSpPr`, a
    grandchild of `cxnSp` - a direct-children-only search silently
    reports every connector as unbound. This is exactly the failure the
    scratch prototype had."""

    def test_bound_connector_is_detected(self):
        prs, slide = _blank_slide()
        _append(slide, _bound_connector_xml(501, target_id=10, target_idx=2))
        rec = _tier_a_for(slide, prs=prs)
        assert rec["connector_count"] == 1
        assert rec["bound_connector_count"] == 1

    def test_bound_connector_tier_b_targets_are_correct(self):
        prs, slide = _blank_slide()
        _append(slide, _bound_connector_xml(501, target_id=10, target_idx=2, with_end=True))
        tb = extract_tier_b(slide, 0)
        cxn = [s for s in tb["shapes"] if s["xml_tag"] == "cxnSp"][0]
        assert cxn["connector"]["status"] == "VALID"
        assert cxn["connector"]["start_endpoint"]["status"] == "VALID"
        assert cxn["connector"]["start_endpoint"]["target_shape_id"] == "10"
        assert cxn["connector"]["start_endpoint"]["connection_idx"] == "2"
        assert cxn["connector"]["end_endpoint"]["status"] == "VALID"
        assert cxn["connector"]["end_endpoint"]["target_shape_id"] == "11"

    def test_unbound_connector_is_not_falsely_bound(self):
        prs, slide = _blank_slide()
        _append(slide, _unbound_connector_xml(502))
        rec = _tier_a_for(slide, prs=prs)
        assert rec["connector_count"] == 1
        assert rec["bound_connector_count"] == 0

    def test_mixed_bound_and_unbound_counted_separately(self):
        prs, slide = _blank_slide()
        _append(slide, _bound_connector_xml(501, target_id=10, target_idx=0))
        _append(slide, _unbound_connector_xml(502))
        _append(slide, _unbound_connector_xml(503))
        rec = _tier_a_for(slide, prs=prs)
        assert rec["connector_count"] == 3
        assert rec["bound_connector_count"] == 1


class TestMalformedConnectorBindings:
    """L0-06: `bound=True` must never be reported from mere element
    presence - each endpoint is ABSENT/VALID/INVALID, and the aggregate
    connector `status` surfaces INVALID explicitly rather than folding
    it into either a clean bind or a clean absence."""

    def test_missing_id_is_invalid(self):
        prs, slide = _blank_slide()
        _append(slide, _custom_connector_xml(1, st_cxn_attrs={"idx": "0"}, end_cxn_attrs=None))
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["start_endpoint"]["status"] == "INVALID"
        assert rec["connector"]["status"] == "INVALID"

    def test_missing_idx_is_invalid(self):
        prs, slide = _blank_slide()
        _append(slide, _custom_connector_xml(1, st_cxn_attrs={"id": "10"}, end_cxn_attrs=None))
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["start_endpoint"]["status"] == "INVALID"
        assert rec["connector"]["status"] == "INVALID"

    def test_non_numeric_id_is_invalid(self):
        prs, slide = _blank_slide()
        _append(slide, _custom_connector_xml(1, st_cxn_attrs={"id": "abc", "idx": "0"}, end_cxn_attrs=None))
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["start_endpoint"]["status"] == "INVALID"

    def test_negative_idx_is_invalid(self):
        prs, slide = _blank_slide()
        _append(slide, _custom_connector_xml(1, st_cxn_attrs={"id": "10", "idx": "-1"}, end_cxn_attrs=None))
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["start_endpoint"]["status"] == "INVALID"

    def test_valid_start_only(self):
        prs, slide = _blank_slide()
        _append(slide, _custom_connector_xml(1, st_cxn_attrs={"id": "10", "idx": "0"}, end_cxn_attrs=None))
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["start_endpoint"]["status"] == "VALID"
        assert rec["connector"]["end_endpoint"]["status"] == "ABSENT"
        assert rec["connector"]["status"] == "VALID"

    def test_valid_end_only(self):
        prs, slide = _blank_slide()
        _append(slide, _custom_connector_xml(1, st_cxn_attrs=None, end_cxn_attrs={"id": "11", "idx": "1"}))
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["start_endpoint"]["status"] == "ABSENT"
        assert rec["connector"]["end_endpoint"]["status"] == "VALID"
        assert rec["connector"]["status"] == "VALID"

    def test_both_valid(self):
        prs, slide = _blank_slide()
        _append(
            slide,
            _custom_connector_xml(
                1, st_cxn_attrs={"id": "10", "idx": "0"}, end_cxn_attrs={"id": "11", "idx": "1"}
            ),
        )
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["status"] == "VALID"

    def test_both_absent(self):
        prs, slide = _blank_slide()
        _append(slide, _custom_connector_xml(1, st_cxn_attrs=None, end_cxn_attrs=None))
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["status"] == "ABSENT"

    def test_one_invalid_one_valid_surfaces_invalid_not_valid(self):
        prs, slide = _blank_slide()
        _append(
            slide,
            _custom_connector_xml(
                1, st_cxn_attrs={"id": "10", "idx": "0"}, end_cxn_attrs={"idx": "0"}  # missing id
            ),
        )
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["connector"]["status"] == "INVALID"

    def test_invalid_connector_counted_separately_from_bound_in_tier_a(self):
        prs, slide = _blank_slide()
        _append(slide, _custom_connector_xml(1, st_cxn_attrs={"id": "abc", "idx": "0"}, end_cxn_attrs=None))
        _append(slide, _bound_connector_xml(2, target_id=20, target_idx=0))
        rec = _tier_a_for(slide, prs=prs)
        assert rec["connector_count"] == 2
        assert rec["bound_connector_count"] == 1
        assert rec["invalid_connector_count"] == 1


class TestGroupRecursion:
    def test_group_members_get_parent_and_depth(self):
        prs, slide = _blank_slide()
        child = _textbox_xml(2, "inside group")
        _append(slide, _group_xml(1, [child]))
        tb = extract_tier_b(slide, 0)
        group = [s for s in tb["shapes"] if s["xml_tag"] == "grpSp"][0]
        member = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        assert member["parent_group_id"] == 1
        assert member["group_depth"] == 1
        assert group["child_ids"] == [2]

    def test_child_ids_are_immediate_only_not_flattened_descendants(self):
        prs, slide = _blank_slide()
        inner_child = _textbox_xml(3, "deep")
        inner_group = _group_xml(2, [inner_child])
        outer_group = _group_xml(1, [inner_group])
        _append(slide, outer_group)
        tb = extract_tier_b(slide, 0)
        outer = [s for s in tb["shapes"] if s["shape_id"] == 1][0]
        inner = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        # the outer group's immediate child is the inner GROUP (id 2)
        # only - not the deeply nested textbox (id 3).
        assert outer["child_ids"] == [2]
        assert inner["child_ids"] == [3]

    def test_nested_group_depth_accumulates(self):
        prs, slide = _blank_slide()
        inner_child = _textbox_xml(3, "deep")
        inner_group = _group_xml(2, [inner_child])
        outer_group = _group_xml(1, [inner_group])
        _append(slide, outer_group)
        tb = extract_tier_b(slide, 0)
        deepest = [s for s in tb["shapes"] if s["shape_id"] == 3][0]
        assert deepest["group_depth"] == 2
        assert deepest["parent_group_id"] == 2
        rec = _tier_a_for(slide, prs=prs)
        assert rec["max_group_depth"] == 2
        assert rec["group_count"] == 2

    def test_top_level_vs_total_shape_count(self):
        prs, slide = _blank_slide()
        child = _textbox_xml(2, "inside")
        _append(slide, _group_xml(1, [child]))
        _append(slide, _textbox_xml(3, "top level sibling"))
        rec = _tier_a_for(slide, prs=prs)
        assert rec["shape_count_top_level"] == 2
        assert rec["shape_count_total"] == 3

    def test_local_z_order_is_per_parent_not_a_global_flattened_counter(self):
        """The previous version's `z_order` was a single counter threaded
        through the entire recursive traversal - a shape's position
        among its own siblings could not be recovered from it. Each
        parent's children must independently start their own local
        z-order at 0."""
        prs, slide = _blank_slide()
        child_a = _textbox_xml(2, "a")
        child_b = _textbox_xml(3, "b")
        _append(slide, _group_xml(1, [child_a, child_b]))
        _append(slide, _textbox_xml(4, "top-level sibling"))
        tb = extract_tier_b(slide, 0)
        by_id = {s["shape_id"]: s for s in tb["shapes"]}
        # top level: group (id 1) then sibling (id 4) -> local z 0, 1
        assert by_id[1]["local_z_order"] == 0
        assert by_id[4]["local_z_order"] == 1
        # inside the group: child_a then child_b -> local z 0, 1 (its
        # OWN sequence, not a continuation of the top-level counter)
        assert by_id[2]["local_z_order"] == 0
        assert by_id[3]["local_z_order"] == 0 or by_id[3]["local_z_order"] == 1
        assert by_id[3]["local_z_order"] == 1


class TestGroupCoordinateResolution:
    """L0-04: grouped-child geometry is honestly labelled `local_` (the
    immediate parent's coordinate space), and `resolved_geometry_in`
    gives the slide-absolute position when the group transform chain is
    reliably computable - never invented when it is not."""

    def test_top_level_shape_resolved_equals_local(self):
        prs, slide = _blank_slide()
        _append(slide, _textbox_xml(1, "top", left=100000, top=200000, width=50000, height=60000))
        tb = extract_tier_b(slide, 0)
        shp = tb["shapes"][0]
        assert shp["local_geometry_in"] == shp["resolved_geometry_in"]

    def test_grouped_child_local_geometry_is_the_raw_child_space_value(self):
        prs, slide = _blank_slide()
        child = _textbox_xml(2, "inside", left=9000000, top=9000000, width=100000, height=100000)
        group = _group_xml(
            1, [child],
            off=(0, 50000), ext=(2000000, 300000),
            ch_off=(9000000, 9000000), ch_ext=(2000000, 300000),
        )
        _append(slide, group)
        tb = extract_tier_b(slide, 0)
        member = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        assert member["local_geometry_in"]["l"] == round(9000000 / 914400, 3)
        assert member["local_geometry_in"]["t"] == round(9000000 / 914400, 3)

    def test_grouped_child_resolved_geometry_applies_group_transform(self):
        # off=(0, 50000), chOff=(9000000, 9000000), scale=1 (ext == chExt)
        # -> resolved_t = 50000 + (9000000 - 9000000) * 1 = 50000
        prs, slide = _blank_slide()
        child = _textbox_xml(2, "inside", left=9000000, top=9000000, width=100000, height=100000)
        group = _group_xml(
            1, [child],
            off=(0, 50000), ext=(2000000, 300000),
            ch_off=(9000000, 9000000), ch_ext=(2000000, 300000),
        )
        _append(slide, group)
        tb = extract_tier_b(slide, 0)
        member = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        assert member["resolved_geometry_in"]["t"] == round(50000 / 914400, 3)
        assert member["resolved_geometry_in"]["l"] == round(0 / 914400, 3)

    def test_nested_group_resolution_composes_both_levels(self):
        inner_child = _textbox_xml(3, "deep", left=5000000, top=5000000, width=10000, height=10000)
        inner_group = _group_xml(
            2, [inner_child],
            off=(1000000, 1000000), ext=(1000000, 1000000),
            ch_off=(5000000, 5000000), ch_ext=(1000000, 1000000),
        )
        outer_group = _group_xml(
            1, [inner_group],
            off=(0, 0), ext=(1000000, 1000000),
            ch_off=(1000000, 1000000), ch_ext=(1000000, 1000000),
        )
        prs, slide = _blank_slide()
        _append(slide, outer_group)
        tb = extract_tier_b(slide, 0)
        deepest = [s for s in tb["shapes"] if s["shape_id"] == 3][0]
        # inner resolve: 1000000 + (5000000-5000000)*1 = 1000000 (in outer's space)
        # outer resolve: 0 + (1000000-1000000)*1 = 0 (slide-absolute)
        assert deepest["resolved_geometry_in"]["l"] == 0.0
        assert deepest["resolved_geometry_in"]["t"] == 0.0

    def test_degenerate_group_transform_yields_none_not_a_guess(self):
        child = _textbox_xml(2, "inside", left=100, top=100, width=10, height=10)
        group = _group_xml(1, [child], off=(0, 0), ext=(1000000, 1000000), ch_off=(0, 0), ch_ext=(0, 0))
        prs, slide = _blank_slide()
        _append(slide, group)
        tb = extract_tier_b(slide, 0)
        member = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        assert member["resolved_geometry_in"] is None
        assert member["local_geometry_in"] is not None

    def test_rotated_group_yields_unresolved_not_approximated(self):
        """L0-04 (real-deck-confirmed): the canonical deck has real
        groups with `rot="10800000" flipH="1"` (slide index 32, groups
        1255/1259/1263). Rotation composition is out of scope - the
        transform must be treated as unresolvable, never approximated by
        silently ignoring the rotation."""
        child = _textbox_xml(2, "inside", left=100, top=100, width=10, height=10)
        group = _group_xml(1, [child], off=(0, 50000), ext=(2000000, 300000), rot="10800000")
        prs, slide = _blank_slide()
        _append(slide, group)
        tb = extract_tier_b(slide, 0)
        member = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        assert member["resolved_geometry_in"] is None
        assert member["local_geometry_in"] is not None
        # hierarchy/z-order/local geometry are entirely unaffected
        group_rec = [s for s in tb["shapes"] if s["shape_id"] == 1][0]
        assert group_rec["child_ids"] == [2]
        assert member["parent_group_id"] == 1
        assert member["group_depth"] == 1

    def test_flip_h_group_yields_unresolved(self):
        child = _textbox_xml(2, "inside", left=100, top=100, width=10, height=10)
        group = _group_xml(1, [child], off=(0, 50000), ext=(2000000, 300000), flip_h="1")
        prs, slide = _blank_slide()
        _append(slide, group)
        tb = extract_tier_b(slide, 0)
        member = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        assert member["resolved_geometry_in"] is None
        assert member["local_geometry_in"] is not None

    def test_flip_v_group_yields_unresolved(self):
        child = _textbox_xml(2, "inside", left=100, top=100, width=10, height=10)
        group = _group_xml(1, [child], off=(0, 50000), ext=(2000000, 300000), flip_v="1")
        prs, slide = _blank_slide()
        _append(slide, group)
        tb = extract_tier_b(slide, 0)
        member = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        assert member["resolved_geometry_in"] is None
        assert member["local_geometry_in"] is not None

    def test_ancestor_rotation_poisons_deeper_nested_descendant(self):
        """A clean (non-rotated) inner group nested inside a rotated
        outer group must still report `resolved_geometry_in is None` for
        its own descendants - the outer ancestor's unresolvable
        transform must propagate down the whole chain, not just to its
        immediate children."""
        deepest_child = _textbox_xml(3, "deep", left=100, top=100, width=10, height=10)
        inner_group = _group_xml(2, [deepest_child], off=(0, 0), ext=(1000000, 1000000))
        outer_group = _group_xml(1, [inner_group], off=(0, 50000), ext=(2000000, 300000), rot="10800000")
        prs, slide = _blank_slide()
        _append(slide, outer_group)
        tb = extract_tier_b(slide, 0)
        deepest = [s for s in tb["shapes"] if s["shape_id"] == 3][0]
        assert deepest["resolved_geometry_in"] is None
        assert deepest["local_geometry_in"] is not None
        assert deepest["group_depth"] == 2

    def test_rot_zero_is_a_normal_resolvable_group(self):
        """`rot="0"` is the explicit default and must resolve exactly
        like a group with no `rot` attribute at all - only a genuinely
        non-default rotation/flip should invalidate resolution."""
        child = _textbox_xml(2, "inside", left=9000000, top=9000000, width=100000, height=100000)
        group = _group_xml(
            1, [child], off=(0, 50000), ext=(2000000, 300000),
            ch_off=(9000000, 9000000), ch_ext=(2000000, 300000), rot="0",
        )
        prs, slide = _blank_slide()
        _append(slide, group)
        tb = extract_tier_b(slide, 0)
        member = [s for s in tb["shapes"] if s["shape_id"] == 2][0]
        assert member["resolved_geometry_in"] is not None
        assert member["resolved_geometry_in"]["t"] == round(50000 / 914400, 3)


class TestFreeformPaths:
    def test_freeform_segment_counts(self):
        prs, slide = _blank_slide()
        _append(slide, _freeform_xml(1, n_lnto=3, n_cubic=1))
        tb = extract_tier_b(slide, 0)
        shp = tb["shapes"][0]
        assert shp["freeform"] == {"moveTo": 1, "lnTo": 3, "cubicBezTo": 1}

    def test_non_freeform_shape_has_no_freeform_facts(self):
        prs, slide = _blank_slide()
        _append(slide, _textbox_xml(1, "plain"))
        tb = extract_tier_b(slide, 0)
        assert tb["shapes"][0]["freeform"] is None


class TestCroppedPicture:
    def test_crop_presence_and_values_detected(self, tmp_path):
        from PIL import Image

        img_path = tmp_path / "img.png"
        Image.new("RGB", (40, 20), color="red").save(img_path)
        prs, slide = _blank_slide()
        pic = slide.shapes.add_picture(str(img_path), Inches(1), Inches(1))
        pic.crop_left = 0.1
        pic.crop_top = 0.2
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["picture"]["has_crop"] is True
        assert rec["picture"]["source_width_px"] == 40
        assert rec["picture"]["source_height_px"] == 20

    def test_uncropped_picture_reports_no_crop(self, tmp_path):
        from PIL import Image

        img_path = tmp_path / "img.png"
        Image.new("RGB", (10, 10), color="blue").save(img_path)
        prs, slide = _blank_slide()
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(1))
        rec = extract_tier_b(slide, 0)["shapes"][0]
        assert rec["picture"]["has_crop"] is False


class TestHiddenShapes:
    def test_hidden_shape_is_flagged(self):
        prs, slide = _blank_slide()
        _append(slide, _hidden_textbox_xml(1))
        _append(slide, _textbox_xml(2, "visible"))
        rec = _tier_a_for(slide, prs=prs)
        assert rec["hidden_shape_count"] == 1


class TestUnsupportedStructuresSurfaced:
    def test_table_is_flagged_unsupported_not_silently_ignored(self):
        prs, slide = _blank_slide()
        slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))
        rec = _tier_a_for(slide, prs=prs)
        assert len(rec["unsupported_structures"]) == 1
        assert rec["shape_count_total"] == 1

    def test_ordinary_deck_has_zero_unsupported_structures(self):
        prs, slide = _blank_slide()
        _append(slide, _textbox_xml(1, "fine"))
        rec = _tier_a_for(slide, prs=prs)
        assert rec["unsupported_structures"] == []


class TestTitleExtraction:
    def test_title_placeholder_is_used_when_present(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Desktop: Placeholder Title"
        raw, basis = extract_title(slide, prs.slide_height)
        assert raw == "Desktop: Placeholder Title"
        assert basis == "title_placeholder"

    def test_header_region_title_without_placeholder(self):
        prs, slide = _blank_slide()
        _append(slide, _textbox_xml(1, "Desktop: Reactions", top=71425))
        raw, basis = extract_title(slide, prs.slide_height)
        assert raw == "Desktop: Reactions"
        assert basis == "header_region_desktop_convention"

    def test_misleading_earlier_body_text_is_not_mistaken_for_title(self):
        prs, slide = _blank_slide()
        _append(slide, _textbox_xml(1, "Copy link: copies a URL to clipboard.", top=2385232))
        _append(slide, _textbox_xml(2, "Desktop: Message Actions", top=71425))
        raw, basis = extract_title(slide, prs.slide_height)
        assert raw == "Desktop: Message Actions"
        assert basis == "header_region_desktop_convention"

    def test_desktop_colon_spacing_variants_normalize_identically(self):
        assert normalize_title("Desktop : Auto Join Channels") == normalize_title(
            "Desktop: Auto Join Channels"
        )
        assert normalize_title("Desktop:   Auto   Join Channels") == normalize_title(
            "Desktop: Auto Join Channels"
        )

    def test_no_text_on_slide_returns_no_title_found(self):
        prs, slide = _blank_slide()
        raw, basis = extract_title(slide, prs.slide_height)
        assert raw is None
        assert basis == "no_title_found"

    def test_header_region_without_desktop_convention_falls_back_to_topmost(self):
        prs, slide = _blank_slide()
        _append(slide, _textbox_xml(1, "Some Other Heading Style", top=50000))
        raw, basis = extract_title(slide, prs.slide_height)
        assert raw == "Some Other Heading Style"
        assert basis == "header_region_fallback"

    def test_grouped_title_beats_top_level_body_fallback(self):
        """L0-05: a title inside a group must be found via resolved
        (slide-absolute) coordinates, not ignored because
        `slide.shapes` only iterates top-level shapes. The group's raw
        local child coordinate (9,000,000 EMU) is far outside the
        header zone and even outside the slide entirely if treated as
        absolute - only correct transform resolution places it in the
        header region."""
        prs, slide = _blank_slide()
        grouped_title = _textbox_xml(2, "Desktop: Grouped Title", left=9000000, top=9000000, width=200000, height=200000)
        group = _group_xml(
            1, [grouped_title],
            off=(0, 50000), ext=(2000000, 300000),
            ch_off=(9000000, 9000000), ch_ext=(2000000, 300000),
        )
        _append(slide, group)
        _append(slide, _textbox_xml(3, "Some unrelated body copy", top=2000000))
        raw, basis = extract_title(slide, prs.slide_height)
        assert raw == "Desktop: Grouped Title"
        assert basis == "header_region_desktop_convention"

    def test_repeated_titles_get_increasing_occurrence_ordinal(self):
        prs = Presentation()
        records = []
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, "Desktop: Message Formatting", top=71425))
            records.append(extract_tier_a(slide, i, prs.slide_height, slide.slide_layout.name))
        assign_topic_occurrence_ordinals(records)
        assert [r["occurrence_ordinal"] for r in records] == [1, 2, 3]
        assert len({r["topic_key"] for r in records}) == 1

    def test_distinct_titles_each_start_at_ordinal_one(self):
        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        _append(s1, _textbox_xml(1, "Desktop: Reactions", top=71425))
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        _append(s2, _textbox_xml(1, "Desktop: Create New Channel", top=71425))
        records = [
            extract_tier_a(s1, 0, prs.slide_height, s1.slide_layout.name),
            extract_tier_a(s2, 1, prs.slide_height, s2.slide_layout.name),
        ]
        assign_topic_occurrence_ordinals(records)
        assert [r["occurrence_ordinal"] for r in records] == [1, 1]


class TestFingerprint:
    def test_fingerprint_is_bounded_and_explainable(self):
        prs, slide = _blank_slide()
        _append(slide, _bound_connector_xml(1, target_id=10, target_idx=0))
        rec = _tier_a_for(slide, prs=prs)
        fp = rec["fingerprint"]
        assert fp["bound_rate"] == 1.0
        assert 0.0 <= fp["cxn_rate"] <= 1.0
        assert fp["shape_count_norm"] <= 1.0

    def test_empty_slide_fingerprint_is_all_zero_not_guarded_denominator_leakage(self):
        """L0-07: the guarded (never-zero) denominator must be used only
        for ratio division, never as the value being normalized - an
        empty slide's `shape_count_norm` must be exactly 0.0, not
        `min(1/60, 1.0)`."""
        prs, slide = _blank_slide()
        rec = _tier_a_for(slide, prs=prs)
        assert rec["shape_count_total"] == 0
        fp = rec["fingerprint"]
        assert fp["shape_count_norm"] == 0.0
        assert fp == {
            "sp_rate": 0.0, "pic_rate": 0.0, "grp_rate": 0.0, "cxn_rate": 0.0,
            "bound_rate": 0.0, "crop_rate": 0.0, "freeform_rate": 0.0, "shape_count_norm": 0.0,
        }


class TestSourceIntegrityTransaction:
    """L0-01: the source is hashed BEFORE `load_deck()`, and re-hashed
    immediately before publication - a replacement occurring at any
    point in between must be caught, and nothing must be published."""

    def test_source_replaced_after_initial_hash_blocks_publication(self, tmp_path, monkeypatch):
        prs_a = Presentation()
        slide = prs_a.slides.add_slide(prs_a.slide_layouts[6])
        _append(slide, _textbox_xml(1, "Desktop: Original", top=71425))
        source = tmp_path / "source.pptx"
        prs_a.save(str(source))

        real_load_deck = layer0.load_deck

        def swapping_load_deck(path):
            # Simulate a concurrent replacement: the real bytes are
            # parsed first (this is what extraction will use), then the
            # file on disk is swapped for different content before the
            # function returns - exactly the boundary Codex described.
            prs = real_load_deck(path)
            prs_b = Presentation()
            slide_b = prs_b.slides.add_slide(prs_b.slide_layouts[6])
            _append(slide_b, _textbox_xml(1, "Desktop: Replaced", top=71425))
            slide_b2 = prs_b.slides.add_slide(prs_b.slide_layouts[6])
            prs_b.save(path)
            return prs

        monkeypatch.setattr(layer0, "load_deck", swapping_load_deck)

        out_dir = tmp_path / "out"
        with pytest.raises(SourceIntegrityError):
            generate_layer0_inventory(str(source), str(out_dir))

        assert not out_dir.exists() or not (out_dir / "manifest.json").exists()

    def test_missing_source_raises_deck_source_error(self, tmp_path):
        with pytest.raises(DeckSourceError):
            generate_layer0_inventory(str(tmp_path / "does-not-exist.pptx"), str(tmp_path / "out"))

    def test_source_integrity_error_is_also_a_deck_source_error(self):
        assert issubclass(SourceIntegrityError, DeckSourceError)


class TestNoMutationOfSource:
    def test_generate_inventory_never_writes_source(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _append(slide, _textbox_xml(1, "Desktop: Test Slide", top=71425))
        source = tmp_path / "source.pptx"
        prs.save(str(source))
        before = _file_hash(source)

        generate_layer0_inventory(str(source), str(tmp_path / "out"))

        assert _file_hash(source) == before


class TestDeterminism:
    def _build_source(self, tmp_path, name="source.pptx"):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _append(slide, _bound_connector_xml(1, target_id=10, target_idx=0))
        _append(slide, _textbox_xml(2, "Desktop: Repeatable Slide", top=71425))
        source = tmp_path / name
        prs.save(str(source))
        return source

    def _read_tree(self, out_dir):
        return {
            "manifest": json.loads((out_dir / "manifest.json").read_bytes()),
            "tier_a": (out_dir / "tier_a_summary.json").read_bytes(),
            "structural_analysis": (out_dir / "structural_analysis.json").read_bytes(),
            "tier_b": {
                p.name: p.read_bytes() for p in sorted((out_dir / "tier_b").glob("*.json"))
            },
        }

    def test_identical_input_produces_byte_identical_full_tree(self, tmp_path):
        source = self._build_source(tmp_path)
        generate_layer0_inventory(str(source), str(tmp_path / "out1"))
        generate_layer0_inventory(str(source), str(tmp_path / "out2"))

        tree1 = self._read_tree(tmp_path / "out1")
        tree2 = self._read_tree(tmp_path / "out2")
        assert tree1["manifest"] == tree2["manifest"]
        assert tree1["tier_a"] == tree2["tier_a"]
        assert tree1["structural_analysis"] == tree2["structural_analysis"]
        assert tree1["tier_b"] == tree2["tier_b"]

    def test_canonical_manifest_has_no_timestamp_field(self, tmp_path):
        source = self._build_source(tmp_path)
        manifest = generate_layer0_inventory(str(source), str(tmp_path / "out"))
        serialized = json.dumps(manifest)
        assert "generated_at" not in serialized
        assert "timestamp" not in serialized.lower()

    def test_relative_and_absolute_source_spelling_produce_identical_manifest(self, tmp_path, monkeypatch):
        source = self._build_source(tmp_path)
        abs_manifest = generate_layer0_inventory(str(source), str(tmp_path / "out_abs"))

        monkeypatch.chdir(tmp_path)
        rel_manifest = generate_layer0_inventory(source.name, str(tmp_path / "out_rel"))

        assert abs_manifest["source_identifier"] == rel_manifest["source_identifier"]
        assert abs_manifest["source_sha256"] == rel_manifest["source_sha256"]
        assert abs_manifest == rel_manifest

    def test_clustering_is_deterministic_given_fixed_seed(self):
        prs = Presentation()
        records = []
        for i in range(12):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if i % 2 == 0:
                _append(slide, _bound_connector_xml(1, target_id=10, target_idx=0))
            else:
                _append(slide, _textbox_xml(1, f"body {i}"))
            records.append(extract_tier_a(slide, i, prs.slide_height, slide.slide_layout.name))
        r1 = cluster_structural_candidates(records, k=3, seed=7)
        r2 = cluster_structural_candidates(records, k=3, seed=7)
        assert r1 == r2


class TestTierASeparation:
    def test_tier_a_record_has_no_shapes_key(self):
        prs, slide = _blank_slide()
        _append(slide, _textbox_xml(1, "some body copy", top=71425))
        rec = _tier_a_for(slide, prs=prs)
        assert "shapes" not in rec

    def test_tier_a_does_not_contain_full_run_text(self):
        prs, slide = _blank_slide()
        long_text = "This exact sentence must never leak into Tier A. " * 5
        _append(slide, _textbox_xml(1, long_text, top=71425))
        rec = _tier_a_for(slide, prs=prs)
        serialized = json.dumps(rec)
        assert long_text not in serialized

    def test_generated_tier_a_file_is_small_relative_to_tier_b(self, tmp_path):
        prs = Presentation()
        for i in range(5):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"Desktop: Slide {i}", top=71425))
        source = tmp_path / "source.pptx"
        prs.save(str(source))
        manifest = generate_layer0_inventory(str(source), str(tmp_path / "out"))
        tier_a_size = (tmp_path / "out" / manifest["tier_a_path"]).stat().st_size
        tier_b_total = sum(
            (tmp_path / "out" / p).stat().st_size for p in manifest["tier_b_paths"]
        )
        assert tier_a_size < tier_b_total


class TestAtomicPublication:
    """L0-03: staged, validated, all-or-nothing, owned-tree publication."""

    def _build_source(self, tmp_path, n_slides, name="source.pptx"):
        prs = Presentation()
        for i in range(n_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"Desktop: Slide {i}", top=71425))
        source = tmp_path / name
        prs.save(str(source))
        return source

    def test_mid_publish_failure_leaves_previous_tree_intact(self, tmp_path, monkeypatch):
        source = self._build_source(tmp_path, 3)
        out_dir = tmp_path / "out"
        manifest_1 = generate_layer0_inventory(str(source), str(out_dir))
        before_tree = {
            "manifest": (out_dir / "manifest.json").read_bytes(),
            "tier_a": (out_dir / "tier_a_summary.json").read_bytes(),
            "tier_b": {p.name: p.read_bytes() for p in (out_dir / "tier_b").glob("*.json")},
        }

        real_replace = os.replace
        call_count = {"n": 0}

        def flaky_replace(src, dst):
            call_count["n"] += 1
            if call_count["n"] == 5:
                raise OSError("simulated disk failure on the 5th replace")
            return real_replace(src, dst)

        monkeypatch.setattr(layer0.os, "replace", flaky_replace)

        with pytest.raises(PublicationError):
            generate_layer0_inventory(str(source), str(out_dir))

        after_tree = {
            "manifest": (out_dir / "manifest.json").read_bytes(),
            "tier_a": (out_dir / "tier_a_summary.json").read_bytes(),
            "tier_b": {p.name: p.read_bytes() for p in (out_dir / "tier_b").glob("*.json")},
        }
        assert before_tree == after_tree
        assert json.loads(before_tree["manifest"]) == manifest_1
        # no `.prev-*`/`.staging-*` leftovers in the parent directory
        leftovers = [p.name for p in tmp_path.iterdir() if ".prev-" in p.name or ".staging-" in p.name]
        assert leftovers == []

    def test_shortening_deck_removes_stale_owned_tier_b_files(self, tmp_path):
        source = self._build_source(tmp_path, 5)
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source), str(out_dir))
        assert sorted(p.name for p in (out_dir / "tier_b").glob("*.json")) == [
            f"slide_{i:04d}.json" for i in range(5)
        ]

        shorter_source = self._build_source(tmp_path, 2, name="shorter.pptx")
        generate_layer0_inventory(str(shorter_source), str(out_dir))
        remaining = sorted(p.name for p in (out_dir / "tier_b").glob("*.json"))
        assert remaining == ["slide_0000.json", "slide_0001.json"]
        assert "slide_0004.json" not in remaining

    def test_unrelated_files_in_output_dir_are_never_touched(self, tmp_path):
        source = self._build_source(tmp_path, 2)
        out_dir = tmp_path / "out"
        out_dir.mkdir()
        unrelated = out_dir / "README.md"
        unrelated.write_text("do not touch me")

        generate_layer0_inventory(str(source), str(out_dir))

        assert unrelated.exists()
        assert unrelated.read_text() == "do not touch me"
        # regenerate again to be sure a second publish still leaves it alone
        generate_layer0_inventory(str(source), str(out_dir))
        assert unrelated.exists()
        assert unrelated.read_text() == "do not touch me"

    def test_unrelated_file_adjacent_to_output_dir_is_never_touched(self, tmp_path):
        source = self._build_source(tmp_path, 2)
        out_dir = tmp_path / "out"
        sibling = tmp_path / "unrelated-sibling.txt"
        sibling.write_text("also not mine")

        generate_layer0_inventory(str(source), str(out_dir))

        assert sibling.exists()
        assert sibling.read_text() == "also not mine"

    def test_regeneration_over_existing_valid_tree_succeeds_and_is_consistent(self, tmp_path):
        source = self._build_source(tmp_path, 3)
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source), str(out_dir))
        manifest_2 = generate_layer0_inventory(str(source), str(out_dir))
        assert manifest_2["slide_count"] == 3
        assert sorted(p.name for p in (out_dir / "tier_b").glob("*.json")) == [
            "slide_0000.json", "slide_0001.json", "slide_0002.json"
        ]

    def test_staging_validation_failure_never_touches_output_dir(self, tmp_path, monkeypatch):
        source = self._build_source(tmp_path, 2)
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source), str(out_dir))
        before = (out_dir / "manifest.json").read_bytes()

        def broken_write(staging_dir, manifest, tier_a_records, tier_b_records, range_proposal, cluster_result):
            # write everything except tier_b, so validation fails
            (staging_dir / "manifest.json").write_text(json.dumps(manifest), encoding="utf-8")
            (staging_dir / "tier_a_summary.json").write_text(json.dumps(tier_a_records), encoding="utf-8")
            (staging_dir / "structural_analysis.json").write_text("{}", encoding="utf-8")

        monkeypatch.setattr(layer0, "_write_staged_tree", broken_write)

        with pytest.raises(StagingValidationError):
            generate_layer0_inventory(str(source), str(out_dir))

        assert (out_dir / "manifest.json").read_bytes() == before

    def test_output_dir_as_existing_file_raises_typed_error_not_file_exists_error(self, tmp_path):
        source = self._build_source(tmp_path, 2)
        out_dir = tmp_path / "out"
        out_dir.write_text("I am a plain file, not a directory")

        with pytest.raises(PublicationError):
            generate_layer0_inventory(str(source), str(out_dir))

        # the file must not have been mutated/replaced
        assert out_dir.is_file()
        assert out_dir.read_text() == "I am a plain file, not a directory"

    def test_rollback_failure_retains_recoverable_backup_and_raises_typed_error(self, tmp_path, monkeypatch):
        """Fault injection: publication commit fails AND the first
        rollback restore also fails. Must raise a typed
        `PublicationRollbackError` (never a raw OSError), must NOT lose
        the previous valid tree (its `.prev-<token>` backup is retained,
        not deleted), and must NOT silently present a mixed/broken tree
        as a successful publish."""
        source = self._build_source(tmp_path, 3)
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source), str(out_dir))
        original_manifest_bytes = (out_dir / "manifest.json").read_bytes()
        original_tier_a_bytes = (out_dir / "tier_a_summary.json").read_bytes()

        real_replace = os.replace
        call_count = {"n": 0}

        def flaky_replace(src, dst):
            call_count["n"] += 1
            # call 5 = the first phase-2 commit replace (manifest.json) -> fails
            # call 6 = the rollback's restore attempt for that same item -> also fails
            if call_count["n"] in (5, 6):
                raise OSError("simulated failure on call %d" % call_count["n"])
            return real_replace(src, dst)

        monkeypatch.setattr(layer0.os, "replace", flaky_replace)

        with pytest.raises(PublicationRollbackError):
            generate_layer0_inventory(str(source), str(out_dir))

        # the canonical manifest.json is missing (neither the new commit
        # nor the restore succeeded) - a truthful failure state, never a
        # silently-presented broken/mixed tree.
        assert not (out_dir / "manifest.json").exists()
        # its previous, still-valid content is retained, recoverable,
        # under an explicitly-named backup - never deleted.
        prev_candidates = list(out_dir.glob("manifest.json.prev-*"))
        assert len(prev_candidates) == 1
        assert prev_candidates[0].read_bytes() == original_manifest_bytes
        # the OTHER owned items, whose individual restores were not
        # injected with a failure, were successfully rolled back to
        # their canonical names.
        assert (out_dir / "tier_a_summary.json").read_bytes() == original_tier_a_bytes
        assert sorted(p.name for p in (out_dir / "tier_b").glob("*.json")) == [
            "slide_0000.json", "slide_0001.json", "slide_0002.json"
        ]
        # no consumer-facing validation may say this tree is valid - the
        # commit marker was left absent precisely because rollback did
        # not fully succeed, even though manifest.json is missing here
        # (a distinct signal `is_tree_valid` also independently catches).
        assert is_tree_valid(out_dir) is False

        # subsequent generation behavior is defined: manually completing
        # the recovery (renaming the backup back to its canonical name)
        # and regenerating again must succeed cleanly.
        prev_candidates[0].rename(out_dir / "manifest.json")
        manifest_after_recovery = generate_layer0_inventory(str(source), str(out_dir))
        assert manifest_after_recovery["slide_count"] == 3
        assert is_tree_valid(out_dir) is True

    def test_exact_mixed_tree_scenario_is_never_exposed_as_valid(self, tmp_path, monkeypatch):
        """Codex's exact reproduction: a later commit-phase failure
        (tier_b's commit) combined with a failure to roll back a
        newly-published canonical item (manifest.json's own removal
        step, during rollback) leaves the manifest advertising NEW
        (2-slide) content while Tier A/Tier B correctly show OLD
        (3-slide) content - individually each file still parses fine, so
        a naive "files exist and parse" check would wrongly call this
        tree valid. `is_tree_valid()` must report False regardless."""
        source_3 = self._build_source(tmp_path, 3, name="source3.pptx")
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source_3), str(out_dir))
        assert is_tree_valid(out_dir) is True
        original_manifest_bytes = (out_dir / "manifest.json").read_bytes()
        original_tier_a_bytes = (out_dir / "tier_a_summary.json").read_bytes()
        original_tier_b_names = sorted(p.name for p in (out_dir / "tier_b").glob("*.json"))

        source_2 = self._build_source(tmp_path, 2, name="source2.pptx")

        real_replace = os.replace
        real_unlink = Path.unlink

        def flaky_replace(src, dst):
            src_s, dst_s = str(src), str(dst)
            # fail exactly the tier_b COMMIT step (staging_dir/tier_b ->
            # output_dir/tier_b) - not the earlier prepare-phase
            # move-aside, and not any later rollback restore.
            if "staging" in src_s and dst_s.endswith(os.sep + "tier_b"):
                raise OSError("simulated failure committing tier_b")
            return real_replace(src, dst)

        def flaky_unlink(self_path, *a, **kw):
            # fail removing exactly the newly-committed canonical
            # manifest.json during rollback - not its `.prev-*` backup,
            # not any other file.
            if self_path.name == "manifest.json" and self_path.parent == out_dir:
                raise OSError("simulated failure removing newly-published manifest.json")
            return real_unlink(self_path, *a, **kw)

        monkeypatch.setattr(layer0.os, "replace", flaky_replace)
        monkeypatch.setattr(Path, "unlink", flaky_unlink)

        with pytest.raises(PublicationRollbackError):
            generate_layer0_inventory(str(source_2), str(out_dir))

        # Reproduces Codex's exact shape: manifest still shows the NEW
        # (2-slide) content...
        manifest_now = json.loads((out_dir / "manifest.json").read_text())
        assert manifest_now["slide_count"] == 2
        # ...while Tier A/Tier B correctly rolled back to the OLD
        # (3-slide) content - a genuinely mixed, self-inconsistent tree.
        assert (out_dir / "tier_a_summary.json").read_bytes() == original_tier_a_bytes
        assert sorted(p.name for p in (out_dir / "tier_b").glob("*.json")) == original_tier_b_names

        # The actual invariant under test: no consumer-facing validation
        # may say this tree is valid, even though every individual file
        # is present and independently parses.
        assert is_tree_valid(out_dir) is False

        # Old data remains recoverable.
        prev_manifest_candidates = list(out_dir.glob("manifest.json.prev-*"))
        assert len(prev_manifest_candidates) == 1
        assert prev_manifest_candidates[0].read_bytes() == original_manifest_bytes

        # Manually restoring the backup file makes the CONTENT
        # consistent again (all three tiers now agree at 3 slides), but
        # `is_tree_valid()` still correctly reports False - the marker,
        # not "do the files happen to agree," is the authoritative
        # signal, and it stays absent until a fresh, successful
        # `generate_layer0_inventory()` call re-establishes it (that
        # specific recovery path - a clean regeneration after manual
        # backup restoration - is exercised end-to-end by the companion
        # `test_rollback_failure_retains_recoverable_backup_and_raises_
        # typed_error` test above, without this test's extra fault
        # injections still active).
        prev_manifest_candidates[0].rename(out_dir / "manifest.json")
        assert json.loads((out_dir / "manifest.json").read_text())["slide_count"] == 3
        assert is_tree_valid(out_dir) is False


class TestPublicationValidityMarker:
    """L0-03 final closure: the commit-validity marker is a hard
    precondition (invalidated BEFORE any canonical mutation, not a
    best-effort corrective action after a failure), binds to the
    COMPLETE owned tree (not just the manifest), and `is_tree_valid()`
    fails closed for every malformed/unexpected input."""

    def _build_source(self, tmp_path, n_slides=2, name="source.pptx"):
        prs = Presentation()
        for i in range(n_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"Desktop: Slide {i}", top=71425))
        source = tmp_path / name
        prs.save(str(source))
        return source

    # ---- hard-precondition invalidation --------------------------------

    def test_marker_invalidation_failure_aborts_before_any_mutation(self, tmp_path, monkeypatch):
        """Codex's failure-mode #1: marker-unlink failure must abort
        publication BEFORE any owned artifact is touched - not merely be
        tolerated mid-rollback."""
        source = self._build_source(tmp_path, 3)
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source), str(out_dir))
        assert is_tree_valid(out_dir) is True
        owned_names = ["manifest.json", "tier_a_summary.json", "structural_analysis.json", "_committed.json"]
        before = {name: (out_dir / name).read_bytes() for name in owned_names}
        before_tier_b = sorted(p.name for p in (out_dir / "tier_b").glob("*.json"))

        real_unlink = Path.unlink

        def flaky_unlink(self_path, *a, **kw):
            if self_path.name == "_committed.json":
                raise OSError("simulated failure invalidating the commit marker")
            return real_unlink(self_path, *a, **kw)

        monkeypatch.setattr(Path, "unlink", flaky_unlink)

        with pytest.raises(PublicationError):
            generate_layer0_inventory(str(source), str(out_dir))

        # nothing was touched - not even a prepare-phase `.prev-*` backup
        # should exist, since publication must never begin moving or
        # replacing canonical artifacts if the marker cannot first be
        # invalidated.
        after = {name: (out_dir / name).read_bytes() for name in owned_names}
        assert after == before
        assert sorted(p.name for p in (out_dir / "tier_b").glob("*.json")) == before_tier_b
        assert list(out_dir.glob("*.prev-*")) == []
        assert is_tree_valid(out_dir) is True

    # ---- no mid-publication valid window --------------------------------

    def test_no_valid_window_exists_during_publication_transitions(self, tmp_path, monkeypatch):
        """Codex's failure-mode #2: same-source regeneration (a
        byte-identical manifest is entirely possible) must not create any
        interval - after moving old files aside, after each canonical
        commit, before the final marker write - where `is_tree_valid()`
        returns `True` for a tree that is not yet (or no longer) fully
        published."""
        source = self._build_source(tmp_path, 3)
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source), str(out_dir))
        assert is_tree_valid(out_dir) is True

        real_replace = os.replace
        validity_snapshots = []

        def probing_replace(src, dst):
            result = real_replace(src, dst)
            validity_snapshots.append(is_tree_valid(out_dir))
            return result

        monkeypatch.setattr(layer0.os, "replace", probing_replace)

        # regenerate from the SAME source - content may be byte-identical
        # to what's already published, which is exactly the case Codex
        # flagged as previously able to let a stale-but-matching marker
        # slip through.
        generate_layer0_inventory(str(source), str(out_dir))

        assert len(validity_snapshots) > 0
        assert all(v is False for v in validity_snapshots), (
            "is_tree_valid() must be False at every intermediate os.replace "
            "during publication - covering the prepare-phase move-asides "
            "and every commit-phase replacement (manifest, tier_a, "
            "structural_analysis, tier_b) - even for a same-source, "
            "content-identical regeneration"
        )
        assert is_tree_valid(out_dir) is True

    # ---- whole-tree commitment: healthy case -----------------------------

    def test_healthy_tree_validates_true(self, tmp_path):
        source = self._build_source(tmp_path, 3)
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source), str(out_dir))
        assert is_tree_valid(out_dir) is True

    # ---- fail-closed: marker itself malformed -----------------------------

    def _published_tree(self, tmp_path, n_slides=3, name="source.pptx"):
        source = self._build_source(tmp_path, n_slides, name)
        out_dir = tmp_path / f"out-{name}"
        generate_layer0_inventory(str(source), str(out_dir))
        assert is_tree_valid(out_dir) is True
        return out_dir

    def test_marker_missing_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "_committed.json").unlink()
        assert is_tree_valid(out_dir) is False

    def test_marker_invalid_json_is_false_not_raising(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "_committed.json").write_text("{not valid json", encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_marker_wrong_json_shape_array_is_false_not_raising(self, tmp_path):
        """Codex's failure-mode #4: `[]` must not raise AttributeError."""
        out_dir = self._published_tree(tmp_path)
        (out_dir / "_committed.json").write_text("[]", encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_marker_missing_required_field_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "_committed.json").write_text(json.dumps({"schema_version": "1.0.0"}), encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_marker_wrong_field_types_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "_committed.json").write_text(
            json.dumps({"tree_sha256": 12345, "owned_file_count": "not a number"}), encoding="utf-8"
        )
        assert is_tree_valid(out_dir) is False

    def test_marker_wrong_digest_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = json.loads((out_dir / "_committed.json").read_text())
        marker["tree_sha256"] = "0" * 64
        (out_dir / "_committed.json").write_text(json.dumps(marker), encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_stale_marker_copied_from_another_tree_is_false(self, tmp_path):
        """Codex's failure-mode: a marker that is valid JSON, has the
        right shape, and even matches ITS OWN tree's digest - but is the
        wrong tree's marker entirely."""
        out_dir_a = self._published_tree(tmp_path, n_slides=3, name="a.pptx")
        out_dir_b = self._published_tree(tmp_path, n_slides=2, name="b.pptx")
        marker_a_bytes = (out_dir_a / "_committed.json").read_bytes()
        (out_dir_b / "_committed.json").write_bytes(marker_a_bytes)
        assert is_tree_valid(out_dir_b) is False
        # tree A itself remains valid throughout - never touched.
        assert is_tree_valid(out_dir_a) is True

    # ---- fail-closed: owned-file mutation ---------------------------------

    def test_manifest_mutation_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        manifest = json.loads((out_dir / "manifest.json").read_text())
        manifest["slide_count"] = 999
        (out_dir / "manifest.json").write_text(json.dumps(manifest), encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_tier_a_mutation_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "tier_a_summary.json").write_text("[]", encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_structural_analysis_mutation_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "structural_analysis.json").write_text("{}", encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_one_tier_b_file_mutation_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        target = sorted((out_dir / "tier_b").glob("*.json"))[0]
        target.write_text(json.dumps({"slide_index": 0, "shapes": []}), encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_missing_tier_b_file_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        target = sorted((out_dir / "tier_b").glob("*.json"))[0]
        target.unlink()
        assert is_tree_valid(out_dir) is False

    def test_extra_stale_owned_tier_b_file_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        stray = out_dir / "tier_b" / "slide_9999.json"
        stray.write_text(json.dumps({"slide_index": 9999, "shapes": []}), encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    # ---- whole-tree digest determinism -------------------------------------

    def test_digest_deterministic_across_separate_generations(self, tmp_path):
        source = self._build_source(tmp_path, 3)
        out_dir_1 = tmp_path / "out1"
        out_dir_2 = tmp_path / "out2"
        generate_layer0_inventory(str(source), str(out_dir_1))
        generate_layer0_inventory(str(source), str(out_dir_2))
        marker_1 = json.loads((out_dir_1 / "_committed.json").read_text())
        marker_2 = json.loads((out_dir_2 / "_committed.json").read_text())
        assert marker_1 == marker_2
        assert is_tree_valid(out_dir_1) is True
        assert is_tree_valid(out_dir_2) is True


class TestParserFailClosed:
    """`is_tree_valid()` is a safe predicate over untrusted-shaped
    input, not a validating parser - it must return `False`, never
    raise, for every malformed marker/manifest state, including schema-
    version mismatches, invalid UTF-8, and JSON tokens (oversized
    integers) that make Python's own `json.loads` raise a plain
    `ValueError` rather than `json.JSONDecodeError`."""

    def _build_source(self, tmp_path, n_slides=2, name="source.pptx"):
        prs = Presentation()
        for i in range(n_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"Desktop: Slide {i}", top=71425))
        source = tmp_path / name
        prs.save(str(source))
        return source

    def _published_tree(self, tmp_path, n_slides=3, name="source.pptx"):
        source = self._build_source(tmp_path, n_slides, name)
        out_dir = tmp_path / f"out-{name}"
        generate_layer0_inventory(str(source), str(out_dir))
        assert is_tree_valid(out_dir) is True
        return out_dir

    def _set_marker(self, out_dir, marker_dict):
        (out_dir / "_committed.json").write_text(json.dumps(marker_dict), encoding="utf-8")

    def _valid_marker(self, out_dir):
        return json.loads((out_dir / "_committed.json").read_text())

    # ---- marker: schema_version enforcement --------------------------------

    def test_marker_missing_schema_version_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        del marker["schema_version"]
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_schema_version_null_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        marker["schema_version"] = None
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_schema_version_integer_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        marker["schema_version"] = 1
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_schema_version_empty_string_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        marker["schema_version"] = ""
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_schema_version_unknown_future_value_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        marker["schema_version"] = "999.0.0"
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_exact_supported_schema_version_validates_normally(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        assert marker["schema_version"] == layer0._MARKER_SCHEMA_VERSION
        assert is_tree_valid(out_dir) is True

    # ---- marker: missing/malformed fields ------------------------------------

    def test_marker_missing_tree_sha256_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        del marker["tree_sha256"]
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_tree_sha256_wrong_type_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        marker["tree_sha256"] = 12345
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_tree_sha256_malformed_string_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        marker["tree_sha256"] = "not-a-real-hex-digest"
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_missing_owned_file_count_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        del marker["owned_file_count"]
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_owned_file_count_wrong_type_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        marker = self._valid_marker(out_dir)
        marker["owned_file_count"] = "173"
        self._set_marker(out_dir, marker)
        assert is_tree_valid(out_dir) is False

    def test_marker_empty_file_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "_committed.json").write_text("", encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_marker_scalar_json_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "_committed.json").write_text("42", encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    # ---- marker: invalid UTF-8 -----------------------------------------------

    def test_marker_invalid_utf8_is_false_not_raising(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "_committed.json").write_bytes(b'{"schema_version": "1.0.0", "tree_sha256": "\xff\xfe invalid utf8')
        assert is_tree_valid(out_dir) is False

    # ---- marker: oversized integer JSON token ---------------------------------

    def test_marker_oversized_integer_is_false_not_raising(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        oversized = "9" * 5000
        (out_dir / "_committed.json").write_text(
            '{"schema_version": "1.0.0", "tree_sha256": "abc", "owned_file_count": %s}' % oversized,
            encoding="utf-8",
        )
        assert is_tree_valid(out_dir) is False

    # ---- manifest: invalid UTF-8 / oversized integer / malformed shape -------

    def test_manifest_invalid_utf8_is_false_not_raising(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "manifest.json").write_bytes(b'{"slide_count": 3, "tier_b_paths": ["\xff\xfe"')
        assert is_tree_valid(out_dir) is False

    def test_manifest_oversized_integer_is_false_not_raising(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        oversized = "9" * 5000
        (out_dir / "manifest.json").write_text(
            '{"slide_count": %s, "tier_b_paths": []}' % oversized, encoding="utf-8"
        )
        assert is_tree_valid(out_dir) is False

    def test_manifest_invalid_json_syntax_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "manifest.json").write_text("{not valid json", encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_manifest_wrong_top_level_shape_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "manifest.json").write_text("[]", encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_manifest_missing_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        (out_dir / "manifest.json").unlink()
        assert is_tree_valid(out_dir) is False

    def test_manifest_malformed_tier_b_declaration_is_false(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        manifest = json.loads((out_dir / "manifest.json").read_text())
        manifest["tier_b_paths"] = "not-a-list"
        (out_dir / "manifest.json").write_text(json.dumps(manifest), encoding="utf-8")
        assert is_tree_valid(out_dir) is False

    def test_manifest_valid_still_validates_true(self, tmp_path):
        out_dir = self._published_tree(tmp_path)
        assert is_tree_valid(out_dir) is True


class TestStagedWriteTypedErrors:
    """No raw `OSError` escapes an expected staging-time filesystem
    failure - `PublicationStagingError` instead, always with the
    original exception chained as `__cause__`, and any previously
    published tree left untouched."""

    def _build_source(self, tmp_path, n_slides=2, name="source.pptx"):
        prs = Presentation()
        for i in range(n_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"Desktop: Slide {i}", top=71425))
        source = tmp_path / name
        prs.save(str(source))
        return source

    def test_tier_b_file_write_failure_is_typed(self, tmp_path, monkeypatch):
        source = self._build_source(tmp_path, 3)
        real_write_text = Path.write_text

        def flaky_write_text(self_path, *a, **kw):
            if self_path.name == "slide_0001.json":
                raise OSError("simulated disk failure writing a staged Tier B file")
            return real_write_text(self_path, *a, **kw)

        monkeypatch.setattr(Path, "write_text", flaky_write_text)

        out_dir = tmp_path / "out"
        with pytest.raises(PublicationStagingError) as excinfo:
            generate_layer0_inventory(str(source), str(out_dir))
        assert isinstance(excinfo.value.__cause__, OSError)
        assert not out_dir.exists() or not (out_dir / "manifest.json").exists()

    def test_manifest_write_failure_is_typed(self, tmp_path, monkeypatch):
        source = self._build_source(tmp_path, 2)
        real_write_text = Path.write_text

        def flaky_write_text(self_path, *a, **kw):
            if self_path.name == "manifest.json" and "staging" in str(self_path):
                raise OSError("simulated disk failure writing the staged manifest")
            return real_write_text(self_path, *a, **kw)

        monkeypatch.setattr(Path, "write_text", flaky_write_text)

        out_dir = tmp_path / "out"
        with pytest.raises(PublicationStagingError):
            generate_layer0_inventory(str(source), str(out_dir))

    def test_staged_write_failure_leaves_prior_published_tree_untouched(self, tmp_path, monkeypatch):
        source_2 = self._build_source(tmp_path, 2, name="s2.pptx")
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source_2), str(out_dir))
        before = (out_dir / "manifest.json").read_bytes()
        assert is_tree_valid(out_dir) is True

        source_3 = self._build_source(tmp_path, 3, name="s3.pptx")
        real_write_text = Path.write_text

        def flaky_write_text(self_path, *a, **kw):
            if self_path.name == "structural_analysis.json" and "staging" in str(self_path):
                raise OSError("simulated failure writing staged structural_analysis.json")
            return real_write_text(self_path, *a, **kw)

        monkeypatch.setattr(Path, "write_text", flaky_write_text)

        with pytest.raises(PublicationStagingError):
            generate_layer0_inventory(str(source_3), str(out_dir))

        # the staging failure never reached publication at all - the
        # prior valid tree is untouched, not merely "rolled back".
        assert (out_dir / "manifest.json").read_bytes() == before
        assert is_tree_valid(out_dir) is True

    def test_staging_directory_creation_failure_is_typed(self, tmp_path, monkeypatch):
        source = self._build_source(tmp_path, 2)
        out_dir = tmp_path / "out"

        real_mkdir = Path.mkdir

        def flaky_mkdir(self_path, *a, **kw):
            if ".staging-" in self_path.name:
                raise OSError("simulated failure creating the staging directory")
            return real_mkdir(self_path, *a, **kw)

        monkeypatch.setattr(Path, "mkdir", flaky_mkdir)

        with pytest.raises(PublicationStagingError):
            generate_layer0_inventory(str(source), str(out_dir))


class TestSingleWriterLock:
    """L0-03: one Layer 0 output tree must have at most one active
    writer at a time."""

    def _build_source(self, tmp_path, n_slides=2, name="source.pptx"):
        prs = Presentation()
        for i in range(n_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"Desktop: Slide {i}", top=71425))
        source = tmp_path / name
        prs.save(str(source))
        return source

    def test_second_writer_fails_fast_while_lock_is_held(self, tmp_path):
        source = self._build_source(tmp_path)
        out_dir = tmp_path / "out"
        out_dir.mkdir()
        lock_path = tmp_path / f".{out_dir.name}.lock"

        fd = os.open(str(lock_path), os.O_CREAT | os.O_RDWR)
        try:
            fcntl.flock(fd, fcntl.LOCK_EX | fcntl.LOCK_NB)
            with pytest.raises(PublicationLockError):
                generate_layer0_inventory(str(source), str(out_dir))
        finally:
            fcntl.flock(fd, fcntl.LOCK_UN)
            os.close(fd)

        # once released, a normal generation succeeds
        manifest = generate_layer0_inventory(str(source), str(out_dir))
        assert manifest["slide_count"] == 2

    def test_lock_file_is_not_deleted_across_successful_generations(self, tmp_path):
        source = self._build_source(tmp_path)
        out_dir = tmp_path / "out"
        generate_layer0_inventory(str(source), str(out_dir))
        lock_path = tmp_path / f".{out_dir.name}.lock"
        assert lock_path.exists()
        generate_layer0_inventory(str(source), str(out_dir))
        assert lock_path.exists()

    def test_concurrent_process_holding_lock_blocks_a_second_writer(self, tmp_path):
        """Real two-process regression: a child process acquires and
        holds the publish lock; while it holds it, the parent process's
        own `generate_layer0_inventory` call must fail fast with
        `PublicationLockError` rather than hang or interleave."""
        source = self._build_source(tmp_path)
        out_dir = tmp_path / "out"
        out_dir.mkdir()

        ready = multiprocessing.Event()
        release = multiprocessing.Event()
        proc = multiprocessing.Process(
            target=_hold_lock_worker, args=(str(tmp_path / f".{out_dir.name}.lock"), ready, release)
        )
        proc.start()
        try:
            assert ready.wait(timeout=10), "child process never signalled it holds the lock"
            with pytest.raises(PublicationLockError):
                generate_layer0_inventory(str(source), str(out_dir))
        finally:
            release.set()
            proc.join(timeout=10)


class TestRangeProposal:
    def test_uniform_deck_produces_no_undersized_ranges(self):
        prs = Presentation()
        records = []
        for i in range(60):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"slide {i}"))
            records.append(extract_tier_a(slide, i, prs.slide_height, slide.slide_layout.name))
        result = propose_ranges(records, target_range_count=4, min_range_size=15)
        assert result["label"] == "WORKLOAD_DECOMPOSITION_ONLY"
        for r in result["ranges"]:
            assert not r["undersized"]

    def test_small_deck_below_min_range_size_yields_single_range(self):
        prs = Presentation()
        records = []
        for i in range(5):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"slide {i}"))
            records.append(extract_tier_a(slide, i, prs.slide_height, slide.slide_layout.name))
        result = propose_ranges(records, target_range_count=4, min_range_size=15)
        assert len(result["ranges"]) == 1
        assert result["ranges"][0]["slide_count"] == 5

    def test_seams_within_merge_window_collapse_to_one(self):
        prs = Presentation()
        records = []
        for i in range(60):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if i in (29, 30, 31):
                _append(slide, _bound_connector_xml(1, target_id=10, target_idx=0))
                _append(slide, _bound_connector_xml(2, target_id=11, target_idx=1))
            else:
                _append(slide, _textbox_xml(1, f"slide {i}"))
            records.append(extract_tier_a(slide, i, prs.slide_height, slide.slide_layout.name))
        result = propose_ranges(records, target_range_count=6, min_range_size=10, merge_window=3)
        seam_positions = [s["before_slide_index"] for s in result["accepted_seams"]]
        close_pairs = [
            (a, b) for a in seam_positions for b in seam_positions
            if a != b and abs(a - b) < 3
        ]
        assert close_pairs == []

    def test_uniform_input_may_have_zero_distance_seams_accepted_as_documented(self):
        """Accepted bounded limitation: workload-only decomposition, not
        a claim about editorial content - pathological uniform input may
        legitimately produce zero-distance adjacent slides."""
        prs = Presentation()
        records = []
        for i in range(30):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, "identical"))
            records.append(extract_tier_a(slide, i, prs.slide_height, slide.slide_layout.name))
        result = propose_ranges(records, target_range_count=4, min_range_size=5)
        assert result["label"] == "WORKLOAD_DECOMPOSITION_ONLY"


class TestStructuralClusteringLabel:
    def test_cluster_and_outlier_output_is_labelled_non_editorial(self):
        prs = Presentation()
        records = []
        for i in range(12):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _append(slide, _textbox_xml(1, f"slide {i}"))
            records.append(extract_tier_a(slide, i, prs.slide_height, slide.slide_layout.name))
        result = cluster_structural_candidates(records, k=3)
        assert result["label"] == "STRUCTURAL_CANDIDATE_CLUSTER"
        for c in result["clusters"]:
            assert c["label"] == "STRUCTURAL_CANDIDATE_CLUSTER"


class TestManifest:
    def test_manifest_records_schema_version_and_hash_and_status(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _append(slide, _textbox_xml(1, "Desktop: One Slide", top=71425))
        source = tmp_path / "source.pptx"
        prs.save(str(source))
        manifest = generate_layer0_inventory(str(source), str(tmp_path / "out"))
        assert manifest["schema_version"] == SCHEMA_VERSION
        assert manifest["source_sha256"] == _file_hash(source)
        assert manifest["slide_count"] == 1
        assert manifest["validation_status"] == "extracted_not_founder_reviewed"
        assert "source_identifier" in manifest
        assert "source_deck_path" not in manifest


@pytest.mark.skipif(not _REAL_DECK.exists(), reason="real MasterSlide deck not present in working tree")
class TestRealDeckValidation:
    def test_source_hash_matches_known_canonical_value_before_and_after(self, tmp_path):
        assert _file_hash(_REAL_DECK) == _REAL_DECK_SHA256
        generate_layer0_inventory(str(_REAL_DECK), str(tmp_path / "out"))
        assert _file_hash(_REAL_DECK) == _REAL_DECK_SHA256

    def test_full_deck_extraction_reports_170_slides_no_unsupported_structures(self, tmp_path):
        manifest = generate_layer0_inventory(str(_REAL_DECK), str(tmp_path / "out"))
        assert manifest["slide_count"] == 170
        tier_a = json.loads((tmp_path / "out" / manifest["tier_a_path"]).read_text())
        assert all(r["unsupported_structures"] == [] for r in tier_a)
        assert is_tree_valid(tmp_path / "out") is True

    def test_known_complex_slides_match_cross_checked_structure(self, tmp_path):
        manifest = generate_layer0_inventory(str(_REAL_DECK), str(tmp_path / "out"))
        tier_a = json.loads((tmp_path / "out" / manifest["tier_a_path"]).read_text())

        reactions = tier_a[78]
        assert reactions["topic_key"] == "desktop: reactions"
        assert reactions["connector_count"] == 2
        assert reactions["bound_connector_count"] == 1

        message_actions = tier_a[74]
        assert message_actions["topic_key"] == "desktop: message actions"
        assert message_actions["connector_count"] == 8
        assert message_actions["bound_connector_count"] == 0
        assert message_actions["freeform_count"] == 8

        auto_join = tier_a[8]
        assert auto_join["topic_key"] == "desktop: auto join channels"
        assert auto_join["group_count"] == 3

        create_channel = tier_a[15]
        assert create_channel["topic_key"] == "desktop: create new channel"

    def test_max_group_depth_reaches_real_nested_groups(self, tmp_path):
        manifest = generate_layer0_inventory(str(_REAL_DECK), str(tmp_path / "out"))
        tier_a = json.loads((tmp_path / "out" / manifest["tier_a_path"]).read_text())
        assert max(r["max_group_depth"] for r in tier_a) >= 2

    def test_deck_identity_matches_known_slide_geometry(self):
        prs = Presentation(str(_REAL_DECK))
        identity = extract_deck_identity(prs)
        assert identity["slide_count"] == 170

    def test_real_rotated_flipped_groups_yield_unresolved_geometry(self):
        """L0-04 real-deck regression. Cross-checked directly against
        the raw OOXML this session (not assumed from prior chat output):
        slide index 32 (`slide33.xml`) contains three top-level groups
        with `rot="10800000" flipH="1"` on their own `p:grpSpPr/a:xfrm` -
        confirmed via `p:grpSpPr/a:xfrm/@rot`/`@flipH` inspection, not
        hardcoded IDs (the IDs are asserted only as evidence of what was
        actually found, per the audit's own "verify against raw XML,
        don't rely on exact IDs" instruction - if production indexing
        ever changes, `rotated_or_flipped_group_ids` below is what this
        test actually discovers and checks, not a fixed expectation)."""
        prs = Presentation(str(_REAL_DECK))
        slide = prs.slides[32]

        def find_rotated_or_flipped_groups(shapes):
            found = []
            for shape in shapes:
                if shape._element.tag.split("}")[-1] == "grpSp":
                    grpSpPr = shape._element.find(qn("p:grpSpPr"))
                    xfrm = grpSpPr.find(qn("a:xfrm")) if grpSpPr is not None else None
                    if xfrm is not None:
                        rot = xfrm.get("rot")
                        flip_h = xfrm.get("flipH")
                        flip_v = xfrm.get("flipV")
                        if (rot not in (None, "0")) or flip_h == "1" or flip_v == "1":
                            found.append(shape.shape_id)
                    found.extend(find_rotated_or_flipped_groups(shape.shapes))
            return found

        rotated_or_flipped_group_ids = find_rotated_or_flipped_groups(slide.shapes)
        assert rotated_or_flipped_group_ids, (
            "expected at least one rotated/flipped group on real slide index 32 - "
            "if this now fails, the canonical deck's structure has changed and "
            "this regression's premise needs re-verification against raw XML, "
            "not silently loosened"
        )
        # matches what this session independently confirmed via raw XML
        assert set(rotated_or_flipped_group_ids) >= {1255, 1259, 1263}

        tier_b = extract_tier_b(slide, 32)
        by_id = {s["shape_id"]: s for s in tier_b["shapes"]}
        for group_id in rotated_or_flipped_group_ids:
            group_rec = by_id[group_id]
            # The group SHAPE's own resolved position is unaffected: per
            # OOXML, `off`/`ext` describe the un-rotated/un-flipped
            # bounding box, so the group's own placement within ITS
            # parent is still reliably resolvable regardless of its own
            # rot/flip - only resolving its CHILDREN through that
            # rotated/flipped child-coordinate space is what must be
            # refused.
            assert group_rec["local_geometry_in"] is not None
            for child_id in group_rec["child_ids"]:
                child_rec = by_id[child_id]
                assert child_rec["resolved_geometry_in"] is None, (
                    f"child {child_id} of rotated/flipped group {group_id} must "
                    "have unresolved geometry, not an approximated one"
                )
                assert child_rec["local_geometry_in"] is not None
