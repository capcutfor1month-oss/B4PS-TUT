"""Documentation Intelligence Slice 2 tests: Existing Documentation Text
Locator. Fixtures are small synthetic .pptx files built with python-pptx,
plus one real-deck regression against the actual MasterSlide archive when
present in the working tree (guarded/skipped when it is not, so this
suite stays portable)."""

from __future__ import annotations

import hashlib
import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from documentation_intelligence.locate import locate_documentation_text
from documentation_intelligence._safe_ppt_engine_import import DeckSourceError

_REPO_ROOT = Path(__file__).resolve().parents[2]
_REAL_DECK = (
    _REPO_ROOT
    / "documentation-artifacts"
    / "masterslide"
    / "old"
    / "MASTER Complete Bridge4PS Desktop-Browser Feature Tutorials.pptx"
)


def _build_deck(path, slide_texts):
    """`slide_texts`: list of lists of shape text (one inner list per
    slide; `None` entries become a shape with no text frame)."""
    prs = Presentation()
    for shapes in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        top = 0.5
        for text in shapes:
            box = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(4), Inches(1))
            if text is not None:
                box.text_frame.text = text
            top += 1.2
    prs.save(str(path))
    return str(path)


def _file_hash(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()


class TestExactMatch:
    def test_exact_text_match_returns_matched(self, tmp_path):
        deck = _build_deck(tmp_path / "d.pptx", [["Pin messages for everyone."]])
        result = locate_documentation_text(deck, "Pin messages for everyone.")
        assert result.status == "matched"
        assert result.match_basis == "exact"
        assert result.slide_index == 0
        assert result.matched_text == "Pin messages for everyone."


class TestWhitespaceNormalizedMatch:
    def test_whitespace_difference_still_matches(self, tmp_path):
        deck = _build_deck(tmp_path / "d.pptx", [["Pin:\nMessages   bookmarked for everyone."]])
        result = locate_documentation_text(deck, "Pin: Messages bookmarked for everyone.")
        assert result.status == "matched"
        assert result.match_basis == "whitespace_normalized"

    def test_exact_takes_priority_over_whitespace_normalized(self, tmp_path):
        # A shape whose text is byte-identical to the query must be
        # reported "exact", never demoted to a weaker basis.
        deck = _build_deck(tmp_path / "d.pptx", [["Pin messages."]])
        result = locate_documentation_text(deck, "Pin messages.")
        assert result.match_basis == "exact"


class TestBoundedSubstringMatch:
    def test_query_is_substring_of_shape_text_unique(self, tmp_path):
        deck = _build_deck(
            tmp_path / "d.pptx",
            [["Header\nPin messages for everyone.\nFooter note."]],
        )
        result = locate_documentation_text(deck, "Pin messages for everyone.")
        assert result.status == "matched"
        assert result.match_basis == "bounded_substring"

    def test_shape_text_is_substring_of_query_unique(self, tmp_path):
        deck = _build_deck(tmp_path / "d.pptx", [["Pin messages."]])
        result = locate_documentation_text(
            deck, "Context: Pin messages. End of context."
        )
        assert result.status == "matched"
        assert result.match_basis == "bounded_substring"

    def test_empty_normalized_text_never_matches_via_substring(self, tmp_path):
        # An all-whitespace query normalizes to "" - it must not
        # spuriously "substring match" every non-whitespace shape in the
        # deck (an empty string is technically a substring of anything).
        deck = _build_deck(tmp_path / "d.pptx", [["Some real text.", "Other text."]])
        result = locate_documentation_text(deck, "   ")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "no_match"


class TestBlankQueryNeverMatches:
    """DI-S2-01: a whitespace-normalized-empty query must never produce
    `status="matched"`, even against an equally blank shape."""

    def test_empty_string_against_empty_text_box(self, tmp_path):
        deck = _build_deck(tmp_path / "d.pptx", [[""]])
        result = locate_documentation_text(deck, "")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "no_match"

    def test_whitespace_only_query_against_empty_text_box(self, tmp_path):
        deck = _build_deck(tmp_path / "d.pptx", [[""]])
        result = locate_documentation_text(deck, "   \n\t ")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "no_match"

    def test_whitespace_only_query_against_whitespace_only_shape_text(self, tmp_path):
        deck = _build_deck(tmp_path / "d.pptx", [["   "]])
        result = locate_documentation_text(deck, "\t\t")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "no_match"

    def test_whitespace_only_query_against_textless_shapes(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))  # no text frame content set
        path = str(tmp_path / "d.pptx")
        prs.save(path)
        result = locate_documentation_text(path, "  ")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "no_match"

    def test_blank_query_in_deck_with_empty_and_nonempty_shapes(self, tmp_path):
        deck = _build_deck(
            tmp_path / "d.pptx", [["", "Pin messages for everyone.", "   "]]
        )
        result = locate_documentation_text(deck, "")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "no_match"
        assert result.slide_index is None
        assert result.shape_index is None


class TestNoMatch:
    def test_no_shape_matches_anything(self, tmp_path):
        deck = _build_deck(tmp_path / "d.pptx", [["Completely unrelated text."]])
        result = locate_documentation_text(deck, "Pin messages for everyone.")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "no_match"
        assert result.slide_index is None
        assert result.shape_index is None

    def test_shape_with_no_text_frame_is_skipped_not_a_match(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))  # empty text frame
        path = str(tmp_path / "d.pptx")
        prs.save(path)
        result = locate_documentation_text(path, "Pin messages for everyone.")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "no_match"


class TestAmbiguousMultipleMatches:
    def test_two_shapes_same_slide_match_is_unresolved(self, tmp_path):
        deck = _build_deck(
            tmp_path / "d.pptx",
            [["Pin messages for everyone.", "Pin messages for everyone."]],
        )
        result = locate_documentation_text(deck, "Pin messages for everyone.")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "ambiguous_multiple_matches"

    def test_duplicate_wording_across_different_slides_is_unresolved(self, tmp_path):
        deck = _build_deck(
            tmp_path / "d.pptx",
            [["Pin messages for everyone."], ["Pin messages for everyone."]],
        )
        result = locate_documentation_text(deck, "Pin messages for everyone.")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "ambiguous_multiple_matches"
        assert result.slide_index is None

    def test_mixed_basis_matches_still_count_as_ambiguous(self, tmp_path):
        # One exact match and one substring match - both are candidates,
        # so this is still ambiguous, never resolved by preferring the
        # "better" basis.
        deck = _build_deck(
            tmp_path / "d.pptx",
            [["Pin messages for everyone.", "Context: Pin messages for everyone. More."]],
        )
        result = locate_documentation_text(deck, "Pin messages for everyone.")
        assert result.status == "unresolved"
        assert result.unresolved_reason == "ambiguous_multiple_matches"


class TestReadOnlyBehavior:
    def test_deck_file_is_byte_identical_before_and_after(self, tmp_path):
        deck = _build_deck(tmp_path / "d.pptx", [["Pin messages for everyone."]])
        before = _file_hash(deck)
        locate_documentation_text(deck, "Pin messages for everyone.")
        locate_documentation_text(deck, "no such text anywhere")
        assert _file_hash(deck) == before

    def test_module_never_references_any_mutation_primitive(self):
        import documentation_intelligence.locate as locate_module

        for mutating_name in (
            "set_shape_text",
            "move_shape",
            "resize_shape",
            "set_shape_geometry",
            "replace_picture",
        ):
            assert not hasattr(locate_module, mutating_name)

    def test_locate_module_imports_only_inspect_deck_from_the_bridge(self):
        # The shared Safe PPT Engine bridge also serves Slice 3
        # (`mutate.py`), which legitimately needs `set_shape_text` - so
        # the bridge's own `__all__` is not the right place to assert
        # Slice 2's read-only guarantee. What matters for Slice 2 is that
        # `locate.py` itself never imports or references anything beyond
        # `inspect_deck` from that bridge - checked directly here.
        import ast
        import inspect

        import documentation_intelligence.locate as locate_module

        source = inspect.getsource(locate_module)
        tree = ast.parse(source)
        imported_names = {
            alias.name
            for node in ast.walk(tree)
            if isinstance(node, ast.ImportFrom) and node.module == "_safe_ppt_engine_import"
            for alias in node.names
        }
        assert imported_names == {"inspect_deck"}


class TestProductKnowledgeNeverParticipates:
    def test_function_signature_takes_no_editorial_memory_or_product_input(self):
        import inspect

        sig = inspect.signature(locate_documentation_text)
        assert list(sig.parameters) == ["deck_path", "documentation_text"]


class TestTypedErrorsPropagate:
    def test_missing_deck_path_raises_deck_source_error(self, tmp_path):
        missing = str(tmp_path / "does-not-exist.pptx")
        with pytest.raises(DeckSourceError):
            locate_documentation_text(missing, "anything")


@pytest.mark.skipif(not _REAL_DECK.exists(), reason="real MasterSlide deck not present in working tree")
class TestRealPinnedMessagesFixture:
    def test_locates_real_pinned_messages_documentation_claim(self):
        # The exact existing-documentation claim already recorded in the
        # real Editorial Memory store's `pinned-messages-masterslide-claim`
        # KnowledgeState, whose own Evidence `source_ref` already asserts
        # "OLD deck slide 75" in this exact file.
        documentation_text = (
            "Pin: Messages bookmarked for all members of the channel to "
            "view in the Pinned Messages list."
        )
        before = _file_hash(_REAL_DECK)
        result = locate_documentation_text(str(_REAL_DECK), documentation_text)
        assert _file_hash(_REAL_DECK) == before, "real deck must never be modified"

        assert result.status == "matched"
        # 0-indexed slide_index 74 == the Evidence source_ref's "slide 75".
        assert result.slide_index == 74
        assert result.match_basis == "whitespace_normalized"

        # Approved PRODUCT Knowledge content for the same subject is a
        # different string entirely and must never be what this test
        # searches for - locating documentation wording is independent of
        # product truth, per the locked Slice 2 boundary.
        product_content = (
            "Pinning a message (via a channel message's 'Pin' action) "
            "bookmarks it for all members of that channel"
        )
        assert product_content != documentation_text
