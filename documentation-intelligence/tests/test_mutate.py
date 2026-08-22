"""Documentation Intelligence Slice 3 tests: Controlled PPT Mutation
Handoff. Fixtures are small synthetic .pptx files built with
python-pptx, plus one real-deck-derived regression using a disposable
`tmp_path` COPY of the actual MasterSlide archive (the real file itself
is never a mutation target in any test)."""

from __future__ import annotations

import hashlib
import os
import shutil
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from documentation_intelligence._safe_ppt_engine_import import (
    DeckSourceError,
    MutationError,
    ValidationError,
)
from documentation_intelligence.mutate import (
    MutationResult,
    OutputCleanupError,
    ShapeIdMismatchError,
    _verify_mutation,
    apply_approved_replacement,
)

_REPO_ROOT = Path(__file__).resolve().parents[2]
_REAL_DECK = (
    _REPO_ROOT
    / "documentation-artifacts"
    / "masterslide"
    / "old"
    / "MASTER Complete Bridge4PS Desktop-Browser Feature Tutorials.pptx"
)


def _build_deck(path, slide_texts):
    """`slide_texts`: list of lists of shape text (one inner list per
    slide)."""
    prs = Presentation()
    for shapes in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        top = 0.5
        for text in shapes:
            box = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(4), Inches(1))
            box.text_frame.text = text
            top += 1.2
    prs.save(str(path))
    return str(path)


def _file_hash(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()


def _shape_id_at(deck_path, slide_index, shape_index):
    prs = Presentation(deck_path)
    return list(prs.slides[slide_index].shapes)[shape_index].shape_id


class TestSuccessfulReplacement:
    def test_target_text_is_replaced_and_result_is_deterministic(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Old text.", "Unrelated."]])
        output = str(tmp_path / "output.pptx")

        result = apply_approved_replacement(source, output, 0, 0, "New approved text.")

        assert isinstance(result, MutationResult)
        assert result.status == "applied"
        assert result.output_path == output
        assert result.slide_index == 0
        assert result.shape_index == 0
        assert result.previous_text == "Old text."
        assert result.new_text == "New approved text."
        assert result.shape_id == _shape_id_at(source, 0, 0)

    def test_shape_id_guard_matching_succeeds(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")
        real_shape_id = _shape_id_at(source, 0, 0)

        result = apply_approved_replacement(
            source, output, 0, 0, "New text.", shape_id=real_shape_id
        )
        assert result.status == "applied"


class TestOutputPublicationAndVerification:
    def test_output_reopens_as_valid_deck_with_target_change_persisted(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")

        apply_approved_replacement(source, output, 0, 0, "New text.")

        prs = Presentation(output)
        assert prs.slides[0].shapes[0].text_frame.text == "New text."

    def test_correct_target_changed(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["A", "B", "C"]])
        output = str(tmp_path / "output.pptx")

        apply_approved_replacement(source, output, 0, 1, "B-updated")

        prs = Presentation(output)
        shapes = list(prs.slides[0].shapes)
        assert shapes[1].text_frame.text == "B-updated"

    def test_unrelated_target_text_unchanged(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["A", "B", "C"], ["D", "E"]])
        output = str(tmp_path / "output.pptx")

        apply_approved_replacement(source, output, 0, 1, "B-updated")

        prs = Presentation(output)
        slide0_shapes = list(prs.slides[0].shapes)
        slide1_shapes = list(prs.slides[1].shapes)
        assert slide0_shapes[0].text_frame.text == "A"
        assert slide0_shapes[2].text_frame.text == "C"
        assert slide1_shapes[0].text_frame.text == "D"
        assert slide1_shapes[1].text_frame.text == "E"

    def test_original_source_bytes_unchanged(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")
        before = _file_hash(source)

        apply_approved_replacement(source, output, 0, 0, "New text.")

        assert _file_hash(source) == before


class TestTargetMissing:
    def test_target_slide_missing_raises_mutation_error(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Only slide."]])
        output = str(tmp_path / "output.pptx")

        with pytest.raises(MutationError):
            apply_approved_replacement(source, output, 5, 0, "New text.")

    def test_target_shape_missing_raises_mutation_error(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Only shape."]])
        output = str(tmp_path / "output.pptx")

        with pytest.raises(MutationError):
            apply_approved_replacement(source, output, 0, 9, "New text.")

    def test_no_output_file_created_on_target_missing(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Only slide."]])
        output = str(tmp_path / "output.pptx")

        with pytest.raises(MutationError):
            apply_approved_replacement(source, output, 5, 0, "New text.")

        assert not Path(output).exists()


class TestShapeIdMismatch:
    def test_shape_id_mismatch_raises_and_is_a_mutation_error(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")
        real_shape_id = _shape_id_at(source, 0, 0)

        with pytest.raises(ShapeIdMismatchError):
            apply_approved_replacement(
                source, output, 0, 0, "New text.", shape_id=real_shape_id + 999
            )

    def test_shape_id_mismatch_is_a_mutation_error_subclass(self):
        assert issubclass(ShapeIdMismatchError, MutationError)

    def test_shape_id_mismatch_never_writes_output(self, tmp_path):
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")
        real_shape_id = _shape_id_at(source, 0, 0)

        with pytest.raises(ShapeIdMismatchError):
            apply_approved_replacement(
                source, output, 0, 0, "New text.", shape_id=real_shape_id + 1
            )

        assert not Path(output).exists()

    def test_shape_id_none_skips_guard_entirely(self, tmp_path):
        # No shape_id supplied -> no guard check at all, just the normal
        # slide/shape-index resolution.
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")

        result = apply_approved_replacement(source, output, 0, 0, "New text.", shape_id=None)
        assert result.status == "applied"


class TestSourceDeckMissing:
    def test_missing_source_deck_raises_deck_source_error(self, tmp_path):
        missing = str(tmp_path / "does-not-exist.pptx")
        output = str(tmp_path / "output.pptx")

        with pytest.raises(DeckSourceError):
            apply_approved_replacement(missing, output, 0, 0, "New text.")


class TestSafePptEngineTypedErrorsPreserved:
    def test_mutation_error_is_the_real_safe_ppt_engine_type(self):
        import documentation_intelligence._safe_ppt_engine_import as bridge

        assert MutationError is bridge.MutationError

    def test_deck_source_error_is_the_real_safe_ppt_engine_type(self):
        import documentation_intelligence._safe_ppt_engine_import as bridge

        assert DeckSourceError is bridge.DeckSourceError


class TestNoEditorialMemoryAccess:
    def test_function_signature_takes_no_editorial_memory_input(self):
        import inspect

        sig = inspect.signature(apply_approved_replacement)
        assert list(sig.parameters) == [
            "source_deck_path",
            "output_deck_path",
            "slide_index",
            "shape_index",
            "replacement_text",
            "shape_id",
        ]

    def test_module_never_imports_editorial_memory(self):
        import ast
        import inspect

        import documentation_intelligence.mutate as mutate_module

        tree = ast.parse(inspect.getsource(mutate_module))
        for node in ast.walk(tree):
            if isinstance(node, (ast.Import, ast.ImportFrom)):
                names = " ".join(
                    [getattr(node, "module", "") or ""]
                    + [alias.name for alias in node.names]
                )
                assert "editorial_memory" not in names
                assert "EditorialMemory" not in names


class TestNoSlice2Involvement:
    def test_mutate_module_never_imports_locate(self):
        import ast
        import inspect

        import documentation_intelligence.mutate as mutate_module

        tree = ast.parse(inspect.getsource(mutate_module))
        for node in ast.walk(tree):
            if isinstance(node, ast.ImportFrom):
                assert node.module != "locate"
            if isinstance(node, ast.Import):
                assert all(alias.name != "documentation_intelligence.locate" for alias in node.names)

    def test_mismatched_target_never_falls_back_to_a_search(self, tmp_path):
        # A deliberately wrong shape_index must fail, never silently
        # locate "the closest text" anywhere else in the deck.
        source = _build_deck(tmp_path / "source.pptx", [["Old text.", "Something else."]])
        output = str(tmp_path / "output.pptx")

        with pytest.raises(MutationError):
            apply_approved_replacement(source, output, 0, 99, "New text.")


@pytest.mark.skipif(not _REAL_DECK.exists(), reason="real MasterSlide deck not present in working tree")
class TestRealPinnedMessagesFixtureCopy:
    def test_applies_replacement_against_a_copy_of_the_real_deck(self, tmp_path):
        # The exact target Slice 2 already proved against the real file.
        # The REAL file is only ever copied here, never itself the
        # mutation source or target.
        source = str(tmp_path / "source.pptx")
        shutil.copyfile(str(_REAL_DECK), source)
        output = str(tmp_path / "output.pptx")

        before_real_hash = _file_hash(_REAL_DECK)

        result = apply_approved_replacement(
            source,
            output,
            slide_index=74,
            shape_index=12,
            replacement_text=(
                "Pin: Messages bookmarked for all members of the channel to "
                "view in the Pinned Messages list. (approved test update)"
            ),
            shape_id=2428,
        )

        assert result.status == "applied"
        assert result.slide_index == 74
        assert result.shape_index == 12
        assert result.shape_id == 2428

        # The real file on disk must be byte-identical to before - only
        # the tmp_path copy was ever touched.
        assert _file_hash(_REAL_DECK) == before_real_hash

        prs = Presentation(output)
        assert "(approved test update)" in list(prs.slides[74].shapes)[12].text_frame.text


def _structure(slides):
    """Builds a minimal `inspect_deck()`-shaped structure for direct
    `_verify_mutation()` unit tests. `slides`: list of lists of
    `(shape_id, text)` tuples, one inner list per slide."""
    return {
        "slide_count": len(slides),
        "slides": [
            {
                "slide_index": s_idx,
                "shape_count": len(shapes),
                "shapes": [
                    {"shape_index": sh_idx, "shape_id": shape_id, "text": text}
                    for sh_idx, (shape_id, text) in enumerate(shapes)
                ],
            }
            for s_idx, shapes in enumerate(slides)
        ],
    }


class TestVerifyMutationAdversarial:
    """DI-S3-01: `_verify_mutation()` must catch every discrepancy a
    naive `zip()`-based comparison could silently miss."""

    def test_target_text_remains_old_after_mutation(self):
        before = _structure([[(1, "old text")]])
        after = _structure([[(1, "old text")]])  # mutation never actually applied
        with pytest.raises(ValidationError):
            _verify_mutation(before, after, 0, 0, "new text")

    def test_trailing_shape_disappears(self):
        before = _structure([[(1, "target"), (2, "trailing shape")]])
        after = _structure([[(1, "new target")]])  # shape 2 vanished
        with pytest.raises(ValidationError):
            _verify_mutation(before, after, 0, 0, "new target")

    def test_trailing_slide_disappears(self):
        before = _structure([[(1, "target")], [(2, "trailing slide content")]])
        after = _structure([[(1, "new target")]])  # slide 1 vanished entirely
        with pytest.raises(ValidationError):
            _verify_mutation(before, after, 0, 0, "new target")

    def test_shape_correspondence_identity_differs(self):
        # Same counts, but the shape_id sequence changed - a different
        # shape now occupies position 1, even though text superficially
        # "lines up" by position.
        before = _structure([[(1, "target"), (2, "second")]])
        after = _structure([[(1, "new target"), (99, "second")]])
        with pytest.raises(ValidationError):
            _verify_mutation(before, after, 0, 0, "new target")

    def test_unrelated_text_changes(self):
        before = _structure([[(1, "target"), (2, "unrelated")]])
        after = _structure([[(1, "new target"), (2, "unrelated but changed")]])
        with pytest.raises(ValidationError):
            _verify_mutation(before, after, 0, 0, "new target")

    def test_correct_mutation_passes_cleanly(self):
        # Positive control: a genuinely correct mutation must not raise.
        before = _structure([[(1, "old target"), (2, "unrelated")]])
        after = _structure([[(1, "new target"), (2, "unrelated")]])
        _verify_mutation(before, after, 0, 0, "new target")  # no raise


class TestOutputCleanupOnVerificationFailure:
    """DI-S3-02: a verification failure after `set_shape_text()` has
    already published must remove only the newly-created output."""

    def test_verification_failure_leaves_no_final_output(self, tmp_path, monkeypatch):
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")

        import documentation_intelligence.mutate as mutate_module

        def _always_fail(*args, **kwargs):
            raise ValidationError("simulated post-mutation verification failure")

        monkeypatch.setattr(mutate_module, "_verify_mutation", _always_fail)

        with pytest.raises(ValidationError):
            apply_approved_replacement(source, output, 0, 0, "New text.")

        assert not Path(output).exists()

    def test_source_remains_byte_identical_on_verification_failure(self, tmp_path, monkeypatch):
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")
        before_hash = _file_hash(source)

        import documentation_intelligence.mutate as mutate_module

        monkeypatch.setattr(
            mutate_module,
            "_verify_mutation",
            lambda *a, **k: (_ for _ in ()).throw(ValidationError("simulated failure")),
        )

        with pytest.raises(ValidationError):
            apply_approved_replacement(source, output, 0, 0, "New text.")

        assert _file_hash(source) == before_hash

    def test_cleanup_never_removes_a_preexisting_output(self, tmp_path):
        # overwrite=False means set_shape_text() itself refuses to
        # publish over an existing output path - this must fail before
        # ever reaching Slice 3's own verification/cleanup logic, and the
        # pre-existing file must be left completely untouched.
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output_path = tmp_path / "output.pptx"
        output_path.write_bytes(b"pre-existing unrelated content")
        before_hash = hashlib.sha256(output_path.read_bytes()).hexdigest()

        from documentation_intelligence._safe_ppt_engine_import import SafeDeckError

        with pytest.raises(SafeDeckError):
            apply_approved_replacement(source, str(output_path), 0, 0, "New text.")

        assert hashlib.sha256(output_path.read_bytes()).hexdigest() == before_hash

    def test_cleanup_failure_is_typed_and_explicit_with_both_contexts(self, tmp_path, monkeypatch):
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")

        import documentation_intelligence.mutate as mutate_module

        original_verification_error = ValidationError("simulated verification failure")

        def _always_fail(*args, **kwargs):
            raise original_verification_error

        # Only the cleanup attempt on this call's own output path must
        # fail - the Safe PPT Engine's own internal use of `os.remove`
        # for its unrelated temp-file cleanup during publication must be
        # left completely unaffected, or this test would spuriously fail
        # inside `set_shape_text()` itself before ever reaching Slice
        # 3's own verification/cleanup logic.
        real_remove = os.remove

        def _remove_only_fails_for_target(path):
            if path == output:
                raise OSError("simulated cleanup failure")
            return real_remove(path)

        monkeypatch.setattr(mutate_module, "_verify_mutation", _always_fail)
        monkeypatch.setattr(mutate_module.os, "remove", _remove_only_fails_for_target)

        with pytest.raises(OutputCleanupError) as excinfo:
            apply_approved_replacement(source, output, 0, 0, "New text.")

        assert excinfo.value.__cause__ is original_verification_error
        assert "simulated cleanup failure" in str(excinfo.value)

    def test_ownership_race_replacement_artifact_is_never_deleted(self, tmp_path, monkeypatch):
        # Codex Re-audit 1 (DI-S3-02, ownership race): between
        # set_shape_text() publishing this call's own output and this
        # function's own verification running, some OTHER process
        # atomically replaces whatever is at output_deck_path with an
        # unrelated artifact. Verification then fails. Cleanup must
        # refuse to delete the replacement - it is no longer the file
        # this call itself created.
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")
        before_source_hash = _file_hash(source)

        import documentation_intelligence.mutate as mutate_module

        replacement_content = b"REPLACEMENT ARTIFACT - NOT SLICE 3'S OWN OUTPUT"

        def _replace_output_pathname_then_fail(*args, **kwargs):
            # Simulate an external, unrelated process atomically
            # replacing the pathname (a rename onto the same path is
            # exactly how a real concurrent writer would do this too).
            replacement_tmp = output + ".external-replacement-tmp"
            Path(replacement_tmp).write_bytes(replacement_content)
            os.replace(replacement_tmp, output)
            raise ValidationError("simulated post-mutation verification failure")

        monkeypatch.setattr(mutate_module, "_verify_mutation", _replace_output_pathname_then_fail)

        with pytest.raises(ValidationError):
            apply_approved_replacement(source, output, 0, 0, "New text.")

        # The replacement artifact must survive completely untouched.
        assert Path(output).exists()
        assert Path(output).read_bytes() == replacement_content

        # The source is never touched by any part of this path.
        assert _file_hash(source) == before_source_hash

    def test_ownership_still_matches_when_nothing_replaces_the_output(self, tmp_path, monkeypatch):
        # Contrast case: when no replacement happens, ordinary cleanup
        # must still remove Slice 3's own output on a verification
        # failure - the ownership check must not become overly cautious
        # and start refusing legitimate cleanups.
        source = _build_deck(tmp_path / "source.pptx", [["Old text."]])
        output = str(tmp_path / "output.pptx")

        import documentation_intelligence.mutate as mutate_module

        monkeypatch.setattr(
            mutate_module,
            "_verify_mutation",
            lambda *a, **k: (_ for _ in ()).throw(ValidationError("simulated failure")),
        )

        with pytest.raises(ValidationError):
            apply_approved_replacement(source, output, 0, 0, "New text.")

        assert not Path(output).exists()
