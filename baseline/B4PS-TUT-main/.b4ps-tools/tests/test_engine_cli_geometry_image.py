"""Safe PPT Engine Build 2 CLI smoke tests - proves the new commands are
reachable end-to-end through b4ps.py."""

import io
import json
import os
import subprocess
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

TOOLS_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
B4PS_PY = os.path.join(TOOLS_DIR, "b4ps.py")


def _run(*args):
    return subprocess.run(
        [sys.executable, B4PS_PY, *args],
        capture_output=True, text=True, cwd=TOOLS_DIR)


def _build_fixture_with_picture(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    buf = io.BytesIO()
    Image.new("RGB", (2, 2), color=(255, 0, 0)).save(buf, format="PNG")
    slide.shapes.add_picture(io.BytesIO(buf.getvalue()), Inches(3), Inches(3), Inches(1), Inches(1))
    prs.save(path)


def test_engine_move_shape_cli(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    _build_fixture_with_picture(str(fixture))
    output = tmp_path / "out.pptx"

    result = _run("engine-move-shape", "--input", str(fixture), "--output", str(output),
                  "--slide", "0", "--shape", "0", "--left", "914400", "--top", "914400")
    assert result.returncode == 0, result.stderr
    assert output.exists()


def test_engine_resize_shape_cli(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    _build_fixture_with_picture(str(fixture))
    output = tmp_path / "out.pptx"

    result = _run("engine-resize-shape", "--input", str(fixture), "--output", str(output),
                  "--slide", "0", "--shape", "0", "--width", "1828800", "--height", "914400")
    assert result.returncode == 0, result.stderr

    verify = _run("engine-inspect", "--input", str(output), "--json")
    payload = json.loads(verify.stdout)
    shape = payload["slides"][0]["shapes"][0]
    assert shape["width"] == 1828800
    assert shape["height"] == 914400


def test_engine_set_geometry_cli(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    _build_fixture_with_picture(str(fixture))
    output = tmp_path / "out.pptx"

    result = _run("engine-set-geometry", "--input", str(fixture), "--output", str(output),
                  "--slide", "0", "--shape", "0",
                  "--left", "457200", "--top", "457200",
                  "--width", "914400", "--height", "457200")
    assert result.returncode == 0, result.stderr

    verify = _run("engine-inspect", "--input", str(output), "--json")
    payload = json.loads(verify.stdout)
    shape = payload["slides"][0]["shapes"][0]
    assert (shape["left"], shape["top"], shape["width"], shape["height"]) == (
        457200, 457200, 914400, 457200)


def test_engine_set_geometry_rejects_zero_width_cli(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    _build_fixture_with_picture(str(fixture))
    output = tmp_path / "out.pptx"

    result = _run("engine-set-geometry", "--input", str(fixture), "--output", str(output),
                  "--slide", "0", "--shape", "0",
                  "--left", "0", "--top", "0", "--width", "0", "--height", "914400")
    assert result.returncode == 1
    assert "error:" in result.stderr
    assert not output.exists()


def test_engine_replace_image_cli(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    _build_fixture_with_picture(str(fixture))
    output = tmp_path / "out.pptx"

    replacement = tmp_path / "new.png"
    buf = io.BytesIO()
    Image.new("RGB", (3, 3), color=(0, 255, 0)).save(buf, format="PNG")
    replacement.write_bytes(buf.getvalue())

    # shape 1 is the picture (textbox added first).
    result = _run("engine-replace-image", "--input", str(fixture), "--output", str(output),
                  "--slide", "0", "--shape", "1", "--image", str(replacement))
    assert result.returncode == 0, result.stderr

    verify = _run("engine-inspect", "--input", str(output), "--json")
    payload = json.loads(verify.stdout)
    assert payload["slides"][0]["shapes"][1]["has_image"] is True


def test_engine_replace_image_rejects_non_picture_target_cli(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    _build_fixture_with_picture(str(fixture))
    output = tmp_path / "out.pptx"

    replacement = tmp_path / "new.png"
    buf = io.BytesIO()
    Image.new("RGB", (3, 3), color=(0, 255, 0)).save(buf, format="PNG")
    replacement.write_bytes(buf.getvalue())

    # shape 0 is the textbox, not a picture.
    result = _run("engine-replace-image", "--input", str(fixture), "--output", str(output),
                  "--slide", "0", "--shape", "0", "--image", str(replacement))
    assert result.returncode == 1
    assert "not a picture" in result.stderr
    assert not output.exists()
