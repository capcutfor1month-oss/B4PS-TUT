"""Safe PPT Engine Audit Repair 1: adversarial regression tests for
independent Codex audit findings F-01 through F-08. Synthetic fixtures
only - no production PPTX content is used or fabricated.
"""

import io
import os
import tempfile
import zipfile
from unittest import mock

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from lib import ppt_engine


def _build_fixture(path, n_slides=2):
    prs = Presentation()
    for i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        box.text_frame.text = "slide %d text" % i
    prs.save(path)
    return path


def _png_bytes(color, size=(4, 4)):
    buf = io.BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def fixture_pptx(tmp_path):
    return _build_fixture(str(tmp_path / "fixture.pptx"))


@pytest.fixture
def fixture_pptx_with_picture(tmp_path):
    def build(path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(io.BytesIO(_png_bytes((255, 0, 0))),
                                 Inches(1), Inches(1), Inches(2), Inches(2))
        prs.save(path)
        return path
    return build(str(tmp_path / "fixture_pic.pptx"))


ALL_MUTATION_CALLS = {
    "set_shape_text": lambda i, o, ow=False: ppt_engine.set_shape_text(
        i, o, 0, 0, "x", overwrite=ow),
    "move_shape": lambda i, o, ow=False: ppt_engine.move_shape(
        i, o, 0, 0, Inches(1), Inches(1), overwrite=ow),
    "resize_shape": lambda i, o, ow=False: ppt_engine.resize_shape(
        i, o, 0, 0, Inches(1), Inches(1), overwrite=ow),
    "set_shape_geometry": lambda i, o, ow=False: ppt_engine.set_shape_geometry(
        i, o, 0, 0, Inches(1), Inches(1), Inches(1), Inches(1), overwrite=ow),
}


# --------------------------------------------------------------------------
# F-01 - source aliasing (symlink / alias) protection, all primitives
# --------------------------------------------------------------------------

@pytest.mark.parametrize("name", list(ALL_MUTATION_CALLS))
def test_f01_input_symlink_to_output_referent_rejected(name, fixture_pptx, tmp_path):
    """A symlink standing in for the input, whose real target is the same
    file the output would be, must still be rejected."""
    real_output = tmp_path / "real_target.pptx"
    _build_fixture(str(real_output))  # distinct content from fixture_pptx
    symlink_input = tmp_path / "input_link.pptx"
    symlink_input.symlink_to(real_output)

    before_hash = ppt_engine.file_sha256(str(real_output))
    with pytest.raises(ppt_engine.OutputPathError):
        ALL_MUTATION_CALLS[name](str(symlink_input), str(real_output))
    assert ppt_engine.file_sha256(str(real_output)) == before_hash


@pytest.mark.parametrize("name", list(ALL_MUTATION_CALLS))
def test_f01_output_symlink_to_input_rejected(name, fixture_pptx, tmp_path):
    """A symlink standing in for the output, whose real target is the
    input file itself, must still be rejected."""
    symlink_output = tmp_path / "output_link.pptx"
    symlink_output.symlink_to(fixture_pptx)

    before_hash = ppt_engine.file_sha256(fixture_pptx)
    with pytest.raises(ppt_engine.OutputPathError):
        ALL_MUTATION_CALLS[name](fixture_pptx, str(symlink_output))
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash


def test_f01_replace_picture_symlink_alias_rejected(fixture_pptx_with_picture, tmp_path):
    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    symlink_output = tmp_path / "output_link.pptx"
    symlink_output.symlink_to(fixture_pptx_with_picture)

    before_hash = ppt_engine.file_sha256(fixture_pptx_with_picture)
    with pytest.raises(ppt_engine.OutputPathError):
        ppt_engine.replace_picture(fixture_pptx_with_picture, str(symlink_output),
                                   0, 0, str(replacement))
    assert ppt_engine.file_sha256(fixture_pptx_with_picture) == before_hash


@pytest.mark.parametrize("name", list(ALL_MUTATION_CALLS))
def test_f01_equivalent_resolved_paths_rejected(name, fixture_pptx, tmp_path):
    """Two different-looking path strings (one with a redundant './' /
    '..' segment) that resolve to the identical file must be rejected."""
    weird_path = os.path.join(os.path.dirname(fixture_pptx), ".",
                              os.path.basename(fixture_pptx))
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    with pytest.raises(ppt_engine.OutputPathError):
        ALL_MUTATION_CALLS[name](fixture_pptx, weird_path)
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash


@pytest.mark.parametrize("overwrite", [True, False])
def test_f01_distinct_files_still_permitted(overwrite, fixture_pptx, tmp_path):
    """The alias protection must not reject ordinary, genuinely distinct
    input/output paths - only aliases of the same file."""
    output = tmp_path / "distinct_output.pptx"
    if overwrite:
        output.write_bytes(b"placeholder")
    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "ok", overwrite=overwrite)
    assert output.exists()


# --------------------------------------------------------------------------
# F-02 - overwrite=False publication race
# --------------------------------------------------------------------------

def test_f02_destination_created_after_preflight_fails_cleanly(fixture_pptx, tmp_path):
    """Simulate the exact race: nothing at the destination during preflight,
    but something appears there before the atomic create happens. The
    real O_CREAT|O_EXCL publish step (not a separate exists-check) must
    catch this and fail cleanly."""
    output = tmp_path / "race_output.pptx"
    real_publish = ppt_engine._publish_atomically

    def racing_publish(staged_path, output_path, overwrite):
        # Winning writer: create the destination right before the real
        # publish call would run.
        with open(output_path, "wb") as fh:
            fh.write(b"raced content")
        return real_publish(staged_path, output_path, overwrite)

    before_hash = ppt_engine.file_sha256(fixture_pptx)
    with mock.patch.object(ppt_engine, "_publish_atomically", side_effect=racing_publish):
        with pytest.raises(ppt_engine.OutputPathError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")

    # The racing writer's content is untouched by our failed publish.
    assert output.read_bytes() == b"raced content"
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash
    assert _no_leaked_temp_files(tmp_path)


def test_f02_overwrite_false_uses_exclusive_create_not_check_then_replace(fixture_pptx, tmp_path):
    """Direct proof the no-overwrite path is O_CREAT|O_EXCL, not a
    preflight exists-check followed by os.replace: create the destination
    file for real, then call the publish primitive directly."""
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"already here")
    staged = ppt_engine.create_working_copy(fixture_pptx)
    try:
        with pytest.raises(ppt_engine.OutputPathError):
            ppt_engine._publish_atomically(staged, str(output), overwrite=False)
    finally:
        if os.path.exists(staged):
            os.remove(staged)
    assert output.read_bytes() == b"already here"


def test_f02_overwrite_true_still_replaces_atomically(fixture_pptx, tmp_path):
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"stale")
    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "new", overwrite=True)
    result = ppt_engine.inspect_deck(str(output))
    assert result["slides"][0]["shapes"][0]["text"] == "new"


def _no_leaked_temp_files(directory):
    names = os.listdir(str(directory)) if os.path.isdir(str(directory)) else []
    leaked_dir = [n for n in names if n.startswith(".b4ps_publish_")]
    leaked_tmp = [n for n in os.listdir(tempfile.gettempdir())
                 if n.startswith("b4ps_engine_")]
    return not leaked_dir and not leaked_tmp


# --------------------------------------------------------------------------
# F-03 - typed errors at every mutation entry point (engine + CLI)
# --------------------------------------------------------------------------

@pytest.mark.parametrize("name", list(ALL_MUTATION_CALLS))
def test_f03_missing_input_raises_typed_error_not_bare_oserror(name, tmp_path):
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.DeckSourceError, match="not found"):
        ALL_MUTATION_CALLS[name](str(tmp_path / "does-not-exist.pptx"), str(output))
    assert not output.exists()


@pytest.mark.parametrize("name", list(ALL_MUTATION_CALLS))
def test_f03_empty_input_raises_typed_error(name, tmp_path):
    empty = tmp_path / "empty.pptx"
    empty.write_bytes(b"")
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.DeckSourceError, match="empty"):
        ALL_MUTATION_CALLS[name](str(empty), str(output))
    assert not output.exists()


@pytest.mark.parametrize("name", list(ALL_MUTATION_CALLS))
def test_f03_corrupt_input_raises_typed_error(name, tmp_path):
    corrupt = tmp_path / "corrupt.pptx"
    corrupt.write_bytes(b"not a real pptx file")
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.DeckSourceError):
        ALL_MUTATION_CALLS[name](str(corrupt), str(output))
    assert not output.exists()


def test_f03_replace_picture_missing_input_typed_error(tmp_path):
    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.DeckSourceError, match="not found"):
        ppt_engine.replace_picture(str(tmp_path / "missing.pptx"), str(output),
                                   0, 0, str(replacement))
    assert not output.exists()


@pytest.mark.parametrize("cli_args,fixture_needs_picture", [
    (["engine-set-text", "--slide", "0", "--shape", "0", "--text", "x"], False),
    (["engine-move-shape", "--slide", "0", "--shape", "0", "--left", "0", "--top", "0"], False),
    (["engine-resize-shape", "--slide", "0", "--shape", "0", "--width", "914400", "--height", "914400"], False),
    (["engine-set-geometry", "--slide", "0", "--shape", "0", "--left", "0", "--top", "0",
      "--width", "914400", "--height", "914400"], False),
])
def test_f03_cli_missing_input_no_traceback_all_mutation_commands(
        cli_args, fixture_needs_picture, tmp_path):
    import subprocess
    import sys
    tools_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output = tmp_path / "out.pptx"
    result = subprocess.run(
        [sys.executable, os.path.join(tools_dir, "b4ps.py"), *cli_args,
         "--input", str(tmp_path / "missing.pptx"), "--output", str(output)],
        capture_output=True, text=True, cwd=tools_dir)
    assert result.returncode == 1
    assert "error:" in result.stderr
    assert "Traceback" not in result.stderr
    assert not output.exists()


def test_f03_cli_replace_image_missing_input_no_traceback(tmp_path):
    import subprocess
    import sys
    tools_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    result = subprocess.run(
        [sys.executable, os.path.join(tools_dir, "b4ps.py"), "engine-replace-image",
         "--input", str(tmp_path / "missing.pptx"), "--output", str(output),
         "--slide", "0", "--shape", "0", "--image", str(replacement)],
        capture_output=True, text=True, cwd=tools_dir)
    assert result.returncode == 1
    assert "error:" in result.stderr
    assert "Traceback" not in result.stderr
    assert not output.exists()


# --------------------------------------------------------------------------
# F-04 - stale media cleanup on picture replacement
# --------------------------------------------------------------------------

def _media_names(pptx_path):
    with zipfile.ZipFile(pptx_path) as zf:
        return {n for n in zf.namelist() if n.startswith("ppt/media/")}


def test_f04_single_picture_replacement_old_media_disappears(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(_png_bytes((255, 0, 0))),
                             Inches(1), Inches(1), Inches(2), Inches(2))
    prs.save(str(fixture))
    before_media = _media_names(str(fixture))
    assert len(before_media) == 1

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))

    after_media = _media_names(str(output))
    assert len(after_media) == 1
    assert after_media.isdisjoint(before_media)


def test_f04_shared_image_replacement_does_not_break_other_picture(tmp_path):
    """Two pictures on the same slide sharing the same underlying image
    (python-pptx assigns them the same relationship id automatically).
    Replacing one must leave the other's image completely intact."""
    fixture = tmp_path / "fixture.pptx"
    shared_bytes = _png_bytes((255, 0, 0))
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(shared_bytes), Inches(1), Inches(1), Inches(1), Inches(1))
    slide.shapes.add_picture(io.BytesIO(shared_bytes), Inches(3), Inches(1), Inches(1), Inches(1))
    prs.save(str(fixture))

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))

    result_prs = Presentation(str(output))
    shape0, shape1 = result_prs.slides[0].shapes[0], result_prs.slides[0].shapes[1]
    with open(replacement, "rb") as fh:
        expected_new = fh.read()
    assert shape0.image.blob == expected_new
    # The untouched sibling picture keeps the original shared image exactly.
    assert shape1.image.blob == shared_bytes


def test_f04_replacement_media_and_geometry_and_reopen(fixture_pptx_with_picture, tmp_path):
    before = ppt_engine.inspect_deck(fixture_pptx_with_picture)["slides"][0]["shapes"][0]
    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(fixture_pptx_with_picture, str(output), 0, 0, str(replacement))

    result = ppt_engine.inspect_deck(str(output))
    assert result["slide_count"] == 1
    after = result["slides"][0]["shapes"][0]
    assert (after["left"], after["top"], after["width"], after["height"]) == (
        before["left"], before["top"], before["width"], before["height"])


# --------------------------------------------------------------------------
# F-05 - exception-safe staging / no leaked temp files on failure
# --------------------------------------------------------------------------

def test_f05_staging_copy_failure_leaves_no_temp_file(fixture_pptx):
    # Superseded by Audit Repair 2 (F-05): staging failures are now typed
    # as TransactionIOError specifically, not the broader DeckSourceError.
    with mock.patch("shutil.copy2", side_effect=OSError("simulated disk failure")):
        with pytest.raises(ppt_engine.TransactionIOError):
            ppt_engine.create_working_copy(fixture_pptx)
    leaked = [n for n in os.listdir(tempfile.gettempdir()) if n.startswith("b4ps_engine_")]
    assert leaked == []


def test_f05_save_failure_cleans_up_staged_file_and_leaves_source_untouched(
        fixture_pptx, tmp_path):
    # Superseded by Audit Repair 2 (F-05): save failures are now typed as
    # TransactionIOError instead of leaking the raw OSError.
    output = tmp_path / "out.pptx"
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    with mock.patch("pptx.presentation.Presentation.save",
                    side_effect=OSError("simulated save failure")):
        with pytest.raises(ppt_engine.TransactionIOError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert not output.exists()
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash
    assert _no_leaked_temp_files(tmp_path)


def test_f05_validation_failure_cleans_up_staged_file(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    with mock.patch("lib.ppt_engine.validate_mutated_output",
                    return_value=(False, "simulated validation failure")):
        with pytest.raises(ppt_engine.ValidationError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert not output.exists()
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash
    assert _no_leaked_temp_files(tmp_path)


def test_f05_publication_failure_cleans_up_temp_and_leaves_source_and_output_untouched(
        fixture_pptx, tmp_path):
    """A failure inside the atomic-publish step itself (after a
    successful, validated mutation) must still leave no partial output and
    no leaked staged/publish temp files, and must not touch the source."""
    output = tmp_path / "out.pptx"
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    with mock.patch.object(ppt_engine, "_publish_atomically",
                           side_effect=OSError("simulated publish failure")):
        with pytest.raises(OSError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert not output.exists()
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash
    assert _no_leaked_temp_files(tmp_path)


# --------------------------------------------------------------------------
# F-06 - required verification records (checked structurally here; content
# reviewed manually as part of the repair - see openspec/changes/
# safe-ppt-engine/verification-report.md)
# --------------------------------------------------------------------------

def test_f06_verification_report_exists_and_is_nonempty():
    tools_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    repo_root = os.path.abspath(os.path.join(tools_dir, "..", "..", ".."))
    report_path = os.path.join(
        repo_root, "openspec", "changes", "safe-ppt-engine", "verification-report.md")
    assert os.path.isfile(report_path), "missing %s" % report_path
    with open(report_path, encoding="utf-8") as fh:
        text = fh.read()
    assert "F-01" in text and "F-08" in text
    assert "re-audit" in text.lower()


# --------------------------------------------------------------------------
# F-07 - non-picture target rejected before staging
# --------------------------------------------------------------------------

def test_f07_non_picture_target_rejected_before_any_staging(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    prs.save(str(fixture))

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"

    before_temp_files = set(os.listdir(tempfile.gettempdir()))
    with pytest.raises(ppt_engine.MutationError, match="not a picture"):
        ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))
    after_temp_files = set(os.listdir(tempfile.gettempdir()))

    new_engine_temp_files = {
        n for n in (after_temp_files - before_temp_files)
        if n.startswith("b4ps_engine_")
    }
    assert new_engine_temp_files == set()
    assert not output.exists()


# --------------------------------------------------------------------------
# F-08 - deterministic canonical inspection across identical copies
# --------------------------------------------------------------------------

def test_f08_identical_pptx_at_different_paths_inspect_identically(fixture_pptx, tmp_path):
    copy_path = tmp_path / "nested" / "copy_of_fixture.pptx"
    copy_path.parent.mkdir()
    import shutil as _shutil
    _shutil.copy2(fixture_pptx, str(copy_path))

    result_a = ppt_engine.inspect_deck(fixture_pptx)
    result_b = ppt_engine.inspect_deck(str(copy_path))
    assert result_a == result_b


def test_f08_canonical_inspection_has_no_path_key(fixture_pptx):
    result = ppt_engine.inspect_deck(fixture_pptx)
    assert "path" not in result
    assert "source_path" not in result


def test_f08_describe_deck_wraps_structure_and_adds_source_path(fixture_pptx):
    result = ppt_engine.describe_deck(fixture_pptx)
    assert result["source_path"] == os.path.abspath(fixture_pptx)
    assert result["structure"] == ppt_engine.inspect_deck(fixture_pptx)
