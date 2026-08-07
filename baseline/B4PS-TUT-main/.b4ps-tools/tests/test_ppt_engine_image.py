"""Safe PPT Engine Build 2: picture-replacement primitive. Synthetic PPTX
and image fixtures only - the 8 unavailable Git LFS production assets are
never touched, fabricated, or stood in for."""

import io
import os

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from lib import ppt_engine


def _png_bytes(color, size=(4, 4)):
    buf = io.BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="PNG")
    return buf.getvalue()


def _build_fixture_with_pictures(path):
    """Slide 0: one picture + one textbox. Slide 1: one picture (unrelated,
    for the "unrelated shapes untouched" check)."""
    prs = Presentation()

    slide0 = prs.slides.add_slide(prs.slide_layouts[6])
    slide0.shapes.add_picture(io.BytesIO(_png_bytes((255, 0, 0))),
                              Inches(1), Inches(1), Inches(2), Inches(2))
    box = slide0.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(1), Inches(0.5))
    box.text_frame.text = "caption"

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.shapes.add_picture(io.BytesIO(_png_bytes((0, 0, 255))),
                              Inches(0.5), Inches(0.5), Inches(1), Inches(1))

    prs.save(path)
    return path


@pytest.fixture
def fixture_pptx(tmp_path):
    return _build_fixture_with_pictures(str(tmp_path / "fixture.pptx"))


@pytest.fixture
def replacement_image(tmp_path):
    path = tmp_path / "replacement.png"
    # Deliberately a different aspect ratio / size than the original 4x4,
    # to exercise "preserve the existing frame regardless of new image size".
    path.write_bytes(_png_bytes((0, 255, 0), size=(8, 2)))
    return str(path)


def test_replace_picture_success(fixture_pptx, replacement_image, tmp_path):
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, replacement_image)

    result = ppt_engine.inspect_deck(str(output))
    shape = result["slides"][0]["shapes"][0]
    assert shape["has_image"] is True


def test_replace_picture_image_actually_differs(fixture_pptx, replacement_image, tmp_path):
    prs_before = Presentation(fixture_pptx)
    before_blob = prs_before.slides[0].shapes[0].image.blob

    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, replacement_image)

    prs_after = Presentation(str(output))
    after_blob = prs_after.slides[0].shapes[0].image.blob
    assert after_blob != before_blob
    with open(replacement_image, "rb") as fh:
        assert after_blob == fh.read()


def test_replace_picture_preserves_existing_frame_geometry(
        fixture_pptx, replacement_image, tmp_path):
    before = ppt_engine.inspect_deck(fixture_pptx)["slides"][0]["shapes"][0]
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, replacement_image)

    after = ppt_engine.inspect_deck(str(output))["slides"][0]["shapes"][0]
    assert (after["left"], after["top"], after["width"], after["height"]) == (
        before["left"], before["top"], before["width"], before["height"])


def test_replace_picture_rejects_non_picture_target(fixture_pptx, replacement_image, tmp_path):
    output = tmp_path / "out.pptx"
    # shape 1 on slide 0 is the textbox, not a picture.
    with pytest.raises(ppt_engine.MutationError, match="not a picture"):
        ppt_engine.replace_picture(fixture_pptx, str(output), 0, 1, replacement_image)
    assert not output.exists()


def test_replace_picture_rejects_missing_image(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.MutationError, match="not found"):
        ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0,
                                   str(tmp_path / "does-not-exist.png"))
    assert not output.exists()


def test_replace_picture_rejects_empty_image(fixture_pptx, tmp_path):
    empty = tmp_path / "empty.png"
    empty.write_bytes(b"")
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.MutationError, match="empty"):
        ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, str(empty))
    assert not output.exists()


def test_replace_picture_rejects_corrupt_image(fixture_pptx, tmp_path):
    corrupt = tmp_path / "corrupt.png"
    corrupt.write_bytes(b"this is not a real image file")
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.MutationError, match="not a valid"):
        ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, str(corrupt))
    assert not output.exists()


def test_replace_picture_source_unchanged(fixture_pptx, replacement_image, tmp_path):
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, replacement_image)
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash


def test_replace_picture_output_reopens(fixture_pptx, replacement_image, tmp_path):
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, replacement_image)
    result = ppt_engine.inspect_deck(str(output))
    assert result["slide_count"] == 2


def test_replace_picture_unrelated_shapes_and_slides_preserved(
        fixture_pptx, replacement_image, tmp_path):
    before = ppt_engine.inspect_deck(fixture_pptx)
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, replacement_image)
    after = ppt_engine.inspect_deck(str(output))

    # The caption textbox on slide 0 is untouched.
    assert after["slides"][0]["shapes"][1]["text"] == before["slides"][0]["shapes"][1]["text"]
    # Slide 1's unrelated picture (its own image + geometry) is untouched.
    before_pic1 = Presentation(fixture_pptx).slides[1].shapes[0].image.blob
    after_pic1 = Presentation(str(output)).slides[1].shapes[0].image.blob
    assert after_pic1 == before_pic1
    assert after["slides"][1]["shapes"][0]["left"] == before["slides"][1]["shapes"][0]["left"]


def test_replace_picture_output_path_collision_rejected(fixture_pptx, replacement_image):
    with pytest.raises(ppt_engine.OutputPathError):
        ppt_engine.replace_picture(fixture_pptx, fixture_pptx, 0, 0, replacement_image)
