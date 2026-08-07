"""Safe PPT Engine Audit Repair 2: adversarial regression tests for the
remaining independent Codex re-audit findings (F-02, F-04, F-05, R-01,
R-02). Synthetic fixtures only - no production PPTX content is used or
fabricated.
"""

import io
import os
import tempfile
import zipfile
from unittest import mock

import pytest
from PIL import Image
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches

from lib import ppt_engine


def _build_fixture(path, n_slides=2):
    prs = Presentation()
    for i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        box.text_frame.text = "slide %d text" % i
    prs.save(path)
    return path


def _png_bytes(color, size=(4, 4)):
    buf = io.BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="PNG")
    return buf.getvalue()


def _media_names(pptx_path):
    with zipfile.ZipFile(pptx_path) as zf:
        return {n for n in zf.namelist() if n.startswith("ppt/media/")}


@pytest.fixture
def fixture_pptx(tmp_path):
    return _build_fixture(str(tmp_path / "fixture.pptx"))


def _no_leaked_temp_files(directory):
    dir_leaks = ([n for n in os.listdir(str(directory)) if n.startswith(".b4ps_publish_")]
                if os.path.isdir(str(directory)) else [])
    tmp_leaks = [n for n in os.listdir(tempfile.gettempdir())
                if n.startswith("b4ps_engine_") or n.startswith(".b4ps_publish_")]
    return not dir_leaks and not tmp_leaks


ALL_MUTATION_CALLS = {
    "set_shape_text": lambda i, o, ow=False: ppt_engine.set_shape_text(
        i, o, 0, 0, "x", overwrite=ow),
    "move_shape": lambda i, o, ow=False: ppt_engine.move_shape(
        i, o, 0, 0, Inches(1), Inches(1), overwrite=ow),
    "resize_shape": lambda i, o, ow=False: ppt_engine.resize_shape(
        i, o, 0, 0, Inches(1), Inches(1), overwrite=ow),
    "set_shape_geometry": lambda i, o, ow=False: ppt_engine.set_shape_geometry(
        i, o, 0, 0, Inches(1), Inches(1), Inches(1), Inches(1), overwrite=ow),
}


# --------------------------------------------------------------------------
# F-02 - true atomic no-overwrite publication
# --------------------------------------------------------------------------

def test_f02_no_overwrite_uses_hardlink_not_exclusive_stream(fixture_pptx, tmp_path):
    """Direct mechanism proof: the no-overwrite path must not open/write
    the final path directly (which would expose a partial file window).
    It must instead write a complete same-directory temp file and use
    os.link to publish - so os.open/os.fdopen is never called with the
    final path itself. A pass-through wrapper is used (not a bare
    MagicMock) so the real filesystem calls made in between still work."""
    output = tmp_path / "out.pptx"
    staged = ppt_engine.create_working_copy(fixture_pptx)
    real_open = os.open
    opened_paths = []

    def spy_open(path, *a, **kw):
        opened_paths.append(path)
        return real_open(path, *a, **kw)

    try:
        with mock.patch("os.open", side_effect=spy_open):
            ppt_engine._publish_atomically(staged, str(output), overwrite=False)
    finally:
        if os.path.exists(staged):
            os.remove(staged)
    assert str(output) not in opened_paths, (
        "final output path must never be opened directly for no-overwrite "
        "publication")
    assert output.exists()


def test_f02_final_path_never_exposes_partial_bytes(fixture_pptx, tmp_path):
    """Instrument os.link itself: by the time the destination becomes
    visible (the link call), the temp file it links from must already
    contain the *complete* staged content - proving no partial-write
    window can exist at the final path, since the temp file that becomes
    the final path is fully written and fsynced before any operation
    touches the destination name."""
    output = tmp_path / "out.pptx"
    staged = ppt_engine.create_working_copy(fixture_pptx)
    expected_size = os.path.getsize(staged)
    real_link = os.link

    def instrumented_link(src, dst):
        assert os.path.getsize(src) == expected_size, (
            "temp file linked to the final path was not fully written first")
        return real_link(src, dst)

    try:
        with mock.patch("os.link", side_effect=instrumented_link):
            ppt_engine._publish_atomically(staged, str(output), overwrite=False)
    finally:
        if os.path.exists(staged):
            os.remove(staged)
    assert output.stat().st_size == expected_size


def test_f02_destination_created_after_preflight_preserved(fixture_pptx, tmp_path):
    """The exact race the finding describes: nothing at the destination
    during preflight, but something appears there before the real publish
    call. The atomic hard-link publish must fail without touching it."""
    output = tmp_path / "race_output.pptx"
    real_publish = ppt_engine._publish_atomically

    def racing_publish(staged_path, output_path, overwrite):
        with open(output_path, "wb") as fh:
            fh.write(b"raced content")
        return real_publish(staged_path, output_path, overwrite)

    before_hash = ppt_engine.file_sha256(fixture_pptx)
    with mock.patch.object(ppt_engine, "_publish_atomically", side_effect=racing_publish):
        with pytest.raises(ppt_engine.OutputPathError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")

    assert output.read_bytes() == b"raced content"
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash
    assert _no_leaked_temp_files(tmp_path)


def test_f02_injected_failure_before_publication_leaves_no_final_output(
        fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    with mock.patch.object(ppt_engine, "_publish_atomically",
                           side_effect=ppt_engine.TransactionIOError("simulated")):
        with pytest.raises(ppt_engine.TransactionIOError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert not output.exists()
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash
    assert _no_leaked_temp_files(tmp_path)


def test_f02_handled_publication_failure_cleans_temp_artifacts(fixture_pptx, tmp_path):
    """Fail the publish-side copy (writing the same-directory temp file)
    and confirm the allocated .b4ps_publish_* temp is removed."""
    output = tmp_path / "out.pptx"
    staged = ppt_engine.create_working_copy(fixture_pptx)
    try:
        with mock.patch("shutil.copyfileobj", side_effect=OSError("simulated")):
            with pytest.raises(ppt_engine.TransactionIOError):
                ppt_engine._publish_atomically(staged, str(output), overwrite=False)
    finally:
        if os.path.exists(staged):
            os.remove(staged)
    assert not output.exists()
    assert _no_leaked_temp_files(tmp_path)


def test_f02_overwrite_true_still_atomic_replace(fixture_pptx, tmp_path):
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"stale")
    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "new", overwrite=True)
    assert ppt_engine.inspect_deck(str(output))["slides"][0]["shapes"][0]["text"] == "new"
    assert _no_leaked_temp_files(tmp_path)


# --------------------------------------------------------------------------
# R-01 - source integrity verified before publication, not after
# --------------------------------------------------------------------------

def test_r01_source_changed_before_publish_blocks_publication_entirely(
        fixture_pptx, tmp_path):
    """Simulate the source changing between hashing and the finalize step:
    the integrity check must be the gate that runs BEFORE _publish_atomically
    is ever called, not a check performed only after output already exists."""
    output = tmp_path / "out.pptx"
    real_publish = ppt_engine._publish_atomically
    publish_was_called = []

    def mutate_source_then_publish(staged_path, output_path, overwrite):
        publish_was_called.append(True)
        return real_publish(staged_path, output_path, overwrite)

    original_finalize = ppt_engine._finalize_transaction

    def spy_finalize(input_path, output_path, staged, overwrite, before_hash):
        # Change the source right where the real integrity gate must catch
        # it - immediately before finalize would otherwise proceed to
        # publish.
        with open(input_path, "ab") as fh:
            fh.write(b"tampered")
        return original_finalize(input_path, output_path, staged, overwrite, before_hash)

    with mock.patch.object(ppt_engine, "_publish_atomically", side_effect=mutate_source_then_publish):
        with mock.patch.object(ppt_engine, "_finalize_transaction", side_effect=spy_finalize):
            with pytest.raises(ppt_engine.SafeDeckError):
                ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")

    assert publish_was_called == [], (
        "publication must never be attempted once the source-integrity "
        "check has failed")
    assert not output.exists()


def test_r01_pre_existing_destination_untouched_on_integrity_failure(
        fixture_pptx, tmp_path):
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"pre-existing")
    original_finalize = ppt_engine._finalize_transaction

    def spy_finalize(input_path, output_path, staged, overwrite, before_hash):
        with open(input_path, "ab") as fh:
            fh.write(b"tampered")
        return original_finalize(input_path, output_path, staged, overwrite, before_hash)

    with mock.patch.object(ppt_engine, "_finalize_transaction", side_effect=spy_finalize):
        with pytest.raises(ppt_engine.SafeDeckError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x", overwrite=True)
    assert output.read_bytes() == b"pre-existing"


def test_r01_integrity_gate_runs_before_publish_for_all_mutation_families(tmp_path):
    """Direct proof, for every mutation family, that _finalize_transaction
    (which performs the pre-publish integrity check) is called before
    _publish_atomically ever runs - by asserting the call sequence."""
    call_order = []
    real_finalize = ppt_engine._finalize_transaction
    real_publish = ppt_engine._publish_atomically

    def tracking_finalize(*a, **kw):
        call_order.append("finalize")
        return real_finalize(*a, **kw)

    def tracking_publish(*a, **kw):
        call_order.append("publish")
        return real_publish(*a, **kw)

    for name, call in ALL_MUTATION_CALLS.items():
        call_order.clear()
        with tempfile.TemporaryDirectory() as d:
            fixture = _build_fixture(os.path.join(d, "fixture.pptx"))
            output = os.path.join(d, "out.pptx")
            with mock.patch.object(ppt_engine, "_finalize_transaction", side_effect=tracking_finalize):
                with mock.patch.object(ppt_engine, "_publish_atomically", side_effect=tracking_publish):
                    call(fixture, output)
        assert call_order == ["finalize", "publish"], (
            "%s: finalize (integrity gate) must run strictly before publish" % name)


def test_r01_cleanup_of_temp_files_still_happens_on_integrity_failure(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    original_finalize = ppt_engine._finalize_transaction

    def spy_finalize(input_path, output_path, staged, overwrite, before_hash):
        with open(input_path, "ab") as fh:
            fh.write(b"tampered")
        return original_finalize(input_path, output_path, staged, overwrite, before_hash)

    with mock.patch.object(ppt_engine, "_finalize_transaction", side_effect=spy_finalize):
        with pytest.raises(ppt_engine.SafeDeckError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert _no_leaked_temp_files(tmp_path)


# --------------------------------------------------------------------------
# F-05 - typed transaction I/O errors (expanded matrix)
# --------------------------------------------------------------------------

@pytest.mark.parametrize("overwrite", [True, False])
def test_f05_publish_copy_failure_typed_both_overwrite_modes(fixture_pptx, tmp_path, overwrite):
    output = tmp_path / "out.pptx"
    if overwrite:
        output.write_bytes(b"stale")
    with mock.patch("shutil.copyfileobj", side_effect=OSError("simulated")):
        with pytest.raises(ppt_engine.TransactionIOError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x", overwrite=overwrite)
    assert _no_leaked_temp_files(tmp_path)


def test_f05_replace_and_publish_failure_typed_overwrite_true(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    output.write_bytes(b"stale")
    with mock.patch("os.replace", side_effect=OSError("simulated replace failure")):
        with pytest.raises(ppt_engine.TransactionIOError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x", overwrite=True)
    assert output.read_bytes() == b"stale"


def test_f05_link_failure_typed_overwrite_false(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    with mock.patch("os.link", side_effect=OSError("simulated link failure")):
        with pytest.raises(ppt_engine.TransactionIOError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert not output.exists()


def test_f05_cli_save_failure_no_traceback(fixture_pptx, tmp_path):
    """CLI-level proof: an I/O failure inside the transaction (not a bad
    argument) must still exit cleanly with no traceback."""
    import subprocess
    import sys
    tools_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output = tmp_path / "out.pptx"
    # Run the CLI against a read-only *directory* as the output location,
    # so publication genuinely fails at the filesystem level without any
    # mocking (a real, reproducible I/O failure).
    readonly_dir = tmp_path / "readonly"
    readonly_dir.mkdir()
    os.chmod(str(readonly_dir), 0o500)
    blocked_output = readonly_dir / "out.pptx"
    try:
        result = subprocess.run(
            [sys.executable, os.path.join(tools_dir, "b4ps.py"), "engine-set-text",
             "--input", fixture_pptx, "--output", str(blocked_output),
             "--slide", "0", "--shape", "0", "--text", "x"],
            capture_output=True, text=True, cwd=tools_dir)
    finally:
        os.chmod(str(readonly_dir), 0o700)
    assert result.returncode == 1
    assert "error:" in result.stderr
    assert "Traceback" not in result.stderr


# --------------------------------------------------------------------------
# F-04 - relationship cleanup inspects the whole slide XML tree
# --------------------------------------------------------------------------

def _add_background_picture_fill(slide, image_path):
    """Manually build a p:bg/p:bgPr/a:blipFill referencing image_path - a
    slide element that carries a blip but is not inside p:spTree, since
    python-pptx has no public API for a picture-fill slide background."""
    _bg_part, bg_rid = slide.part.get_or_add_image_part(image_path)
    bgPr = slide._element.cSld.get_or_add_bgPr()
    no_fill = bgPr.find(qn("a:noFill"))
    if no_fill is not None:
        bgPr.remove(no_fill)
    blip_fill_xml = (
        '<a:blipFill %s><a:blip r:embed="%s"/><a:stretch><a:fillRect/></a:stretch></a:blipFill>'
        % (nsdecls("a", "r"), bg_rid))
    bgPr.insert(0, parse_xml(blip_fill_xml))
    return bg_rid


def test_f04_single_picture_replacement_old_media_disappears(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(_png_bytes((255, 0, 0))),
                             Inches(1), Inches(1), Inches(2), Inches(2))
    prs.save(str(fixture))
    assert len(_media_names(str(fixture))) == 1

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))

    after_media = _media_names(str(output))
    assert len(after_media) == 1


def test_f04_shared_picture_relationship_preserved(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    shared_bytes = _png_bytes((255, 0, 0))
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(shared_bytes), Inches(1), Inches(1), Inches(1), Inches(1))
    slide.shapes.add_picture(io.BytesIO(shared_bytes), Inches(3), Inches(1), Inches(1), Inches(1))
    prs.save(str(fixture))

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))

    result_prs = Presentation(str(output))
    shape1 = result_prs.slides[0].shapes[1]
    assert shape1.image.blob == shared_bytes


def test_f04_background_sharing_relationship_preserved_old_bug_would_break_this(tmp_path):
    """The exact scenario the old spTree-only scan missed: the foreground
    picture and the slide background share one relationship id (python-pptx
    assigns the same rId when the same image content is related to the
    same slide part twice). Replacing the foreground picture must not drop
    the relationship the background still needs - the old implementation,
    scanning only spTree, would have incorrectly concluded the id was
    unused (since the background blip lives outside spTree) and removed it,
    breaking the background."""
    fixture = tmp_path / "fixture.pptx"
    shared_bytes = _png_bytes((123, 45, 67))
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fg_pic = slide.shapes.add_picture(io.BytesIO(shared_bytes), Inches(1), Inches(1), Inches(1), Inches(1))

    fd, tmp_img = tempfile.mkstemp(suffix=".png")
    os.write(fd, shared_bytes)
    os.close(fd)
    try:
        bg_rid = _add_background_picture_fill(slide, tmp_img)
    finally:
        os.remove(tmp_img)
    fg_rid = fg_pic._pic.blipFill.blip.rEmbed
    assert bg_rid == fg_rid, "test setup assumption: fg picture and bg share one relationship id"
    prs.save(str(fixture))

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))

    # The background's blip (untouched by the mutation) must still resolve
    # to a valid, present relationship and image part.
    result_prs = Presentation(str(output))
    result_slide = result_prs.slides[0]._element
    bg_blip = None
    for blip in result_slide.iter(qn("a:blip")):
        # The foreground picture now points at the new image; the
        # remaining blip belongs to the background.
        rid = blip.get(qn("r:embed"))
        target = result_prs.slides[0].part.related_part(rid)
        if target.blob == shared_bytes:
            bg_blip = blip
    assert bg_blip is not None, "background relationship/media was incorrectly dropped"

    # Foreground picture itself was correctly replaced.
    assert result_prs.slides[0].shapes[0].image.blob != shared_bytes


def test_f04_crop_preserved_across_replacement(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(io.BytesIO(_png_bytes((255, 0, 0))),
                                   Inches(1), Inches(1), Inches(2), Inches(2))
    pic.crop_left = 0.1
    pic.crop_top = 0.2
    pic.crop_right = 0.05
    pic.crop_bottom = 0.15
    prs.save(str(fixture))

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))

    result_pic = Presentation(str(output)).slides[0].shapes[0]
    assert result_pic.crop_left == pytest.approx(0.1)
    assert result_pic.crop_top == pytest.approx(0.2)
    assert result_pic.crop_right == pytest.approx(0.05)
    assert result_pic.crop_bottom == pytest.approx(0.15)


def test_f04_result_reopens_successfully_background_case(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    shared_bytes = _png_bytes((10, 20, 30))
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(shared_bytes), Inches(1), Inches(1), Inches(1), Inches(1))
    fd, tmp_img = tempfile.mkstemp(suffix=".png")
    os.write(fd, shared_bytes)
    os.close(fd)
    try:
        _add_background_picture_fill(slide, tmp_img)
    finally:
        os.remove(tmp_img)
    prs.save(str(fixture))

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((99, 88, 77)))
    output = tmp_path / "out.pptx"
    ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))

    result = ppt_engine.inspect_deck(str(output))
    assert result["slide_count"] == 1


# --------------------------------------------------------------------------
# R-02 - tests that actually prove the mechanism, not just end state
# --------------------------------------------------------------------------

def test_r02_f07_non_picture_target_zero_staging_calls_via_spy(tmp_path):
    """Direct spy on create_working_copy - a non-picture target must
    result in literally zero calls to it, not merely "no leftover temp
    file" (which could pass even if staging happened and cleaned up)."""
    fixture = tmp_path / "fixture.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    prs.save(str(fixture))

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"

    with mock.patch.object(ppt_engine, "create_working_copy") as spy:
        with pytest.raises(ppt_engine.MutationError, match="not a picture"):
            ppt_engine.replace_picture(str(fixture), str(output), 0, 0, str(replacement))
        spy.assert_not_called()


@pytest.mark.parametrize("name", list(ALL_MUTATION_CALLS))
@pytest.mark.parametrize("defect", ["missing", "empty", "corrupt", "valid_zip_not_pptx"])
def test_r02_f03_matrix_all_mutation_families(name, defect, tmp_path):
    bad_path = tmp_path / "bad.pptx"
    if defect == "missing":
        pass  # do not create it
    elif defect == "empty":
        bad_path.write_bytes(b"")
    elif defect == "corrupt":
        bad_path.write_bytes(b"not a pptx or zip at all")
    elif defect == "valid_zip_not_pptx":
        with zipfile.ZipFile(str(bad_path), "w") as zf:
            zf.writestr("hello.txt", "not a pptx")

    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.DeckSourceError):
        ALL_MUTATION_CALLS[name](str(bad_path), str(output))
    assert not output.exists()


@pytest.mark.parametrize("defect", ["missing", "empty", "corrupt", "valid_zip_not_pptx"])
def test_r02_f03_matrix_replace_picture(defect, tmp_path):
    bad_path = tmp_path / "bad.pptx"
    if defect == "missing":
        pass
    elif defect == "empty":
        bad_path.write_bytes(b"")
    elif defect == "corrupt":
        bad_path.write_bytes(b"not a pptx or zip at all")
    elif defect == "valid_zip_not_pptx":
        with zipfile.ZipFile(str(bad_path), "w") as zf:
            zf.writestr("hello.txt", "not a pptx")

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.DeckSourceError):
        ppt_engine.replace_picture(str(bad_path), str(output), 0, 0, str(replacement))
    assert not output.exists()


@pytest.mark.parametrize("cli_args", [
    ["engine-set-text", "--slide", "0", "--shape", "0", "--text", "x"],
    ["engine-move-shape", "--slide", "0", "--shape", "0", "--left", "0", "--top", "0"],
    ["engine-resize-shape", "--slide", "0", "--shape", "0", "--width", "914400", "--height", "914400"],
    ["engine-set-geometry", "--slide", "0", "--shape", "0", "--left", "0", "--top", "0",
     "--width", "914400", "--height", "914400"],
])
@pytest.mark.parametrize("defect", ["missing", "empty", "corrupt", "valid_zip_not_pptx"])
def test_r02_f03_matrix_cli(cli_args, defect, tmp_path):
    import subprocess
    import sys
    tools_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    bad_path = tmp_path / "bad.pptx"
    if defect == "missing":
        pass
    elif defect == "empty":
        bad_path.write_bytes(b"")
    elif defect == "corrupt":
        bad_path.write_bytes(b"not a pptx or zip at all")
    elif defect == "valid_zip_not_pptx":
        with zipfile.ZipFile(str(bad_path), "w") as zf:
            zf.writestr("hello.txt", "not a pptx")

    output = tmp_path / "out.pptx"
    result = subprocess.run(
        [sys.executable, os.path.join(tools_dir, "b4ps.py"), *cli_args,
         "--input", str(bad_path), "--output", str(output)],
        capture_output=True, text=True, cwd=tools_dir)
    assert result.returncode == 1
    assert "error:" in result.stderr
    assert "Traceback" not in result.stderr
    assert not output.exists()


@pytest.mark.parametrize("defect", ["missing", "empty", "corrupt", "valid_zip_not_pptx"])
def test_r02_f03_matrix_cli_replace_image(defect, tmp_path):
    import subprocess
    import sys
    tools_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    bad_path = tmp_path / "bad.pptx"
    if defect == "missing":
        pass
    elif defect == "empty":
        bad_path.write_bytes(b"")
    elif defect == "corrupt":
        bad_path.write_bytes(b"not a pptx or zip at all")
    elif defect == "valid_zip_not_pptx":
        with zipfile.ZipFile(str(bad_path), "w") as zf:
            zf.writestr("hello.txt", "not a pptx")

    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    result = subprocess.run(
        [sys.executable, os.path.join(tools_dir, "b4ps.py"), "engine-replace-image",
         "--input", str(bad_path), "--output", str(output),
         "--slide", "0", "--shape", "0", "--image", str(replacement)],
        capture_output=True, text=True, cwd=tools_dir)
    assert result.returncode == 1
    assert "error:" in result.stderr
    assert "Traceback" not in result.stderr
    assert not output.exists()


# --------------------------------------------------------------------------
# Regression guard - already-closed findings must remain closed
# --------------------------------------------------------------------------

def test_regression_f01_alias_protection_still_active(fixture_pptx, tmp_path):
    symlink_output = tmp_path / "output_link.pptx"
    symlink_output.symlink_to(fixture_pptx)
    with pytest.raises(ppt_engine.OutputPathError):
        ppt_engine.set_shape_text(fixture_pptx, str(symlink_output), 0, 0, "x")


def test_regression_f03_typed_error_on_missing_input_still_active(tmp_path):
    with pytest.raises(ppt_engine.DeckSourceError):
        ppt_engine.set_shape_text(str(tmp_path / "missing.pptx"), str(tmp_path / "out.pptx"),
                                  0, 0, "x")


def test_regression_f07_preflight_still_active(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    with pytest.raises(ppt_engine.MutationError, match="not a picture"):
        ppt_engine.replace_picture(fixture_pptx, str(output), 0, 0, str(replacement))


def test_regression_f08_deterministic_inspection_still_active(fixture_pptx, tmp_path):
    import shutil as _shutil
    copy_path = tmp_path / "nested" / "copy.pptx"
    copy_path.parent.mkdir()
    _shutil.copy2(fixture_pptx, str(copy_path))
    assert ppt_engine.inspect_deck(fixture_pptx) == ppt_engine.inspect_deck(str(copy_path))
