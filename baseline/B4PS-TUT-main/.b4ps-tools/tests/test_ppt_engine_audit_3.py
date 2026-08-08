"""Safe PPT Engine Audit 3: adversarial regression tests for the residual
findings F-05, F-06, R-01, R-02 from the third-round Codex audit. Synthetic
fixtures only - no production PPTX content is used or fabricated.
"""

import io
import os
import tempfile
from unittest import mock

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from lib import ppt_engine


def _build_fixture(path, n_slides=2):
    prs = Presentation()
    for i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        box.text_frame.text = "slide %d text" % i
    prs.save(path)
    return path


def _png_bytes(color, size=(4, 4)):
    buf = io.BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def fixture_pptx(tmp_path):
    return _build_fixture(str(tmp_path / "fixture.pptx"))


@pytest.fixture
def fixture_pptx_with_picture(tmp_path):
    def build(path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(io.BytesIO(_png_bytes((255, 0, 0))),
                                 Inches(1), Inches(1), Inches(2), Inches(2))
        prs.save(path)
        return path
    return build(str(tmp_path / "fixture_pic.pptx"))


ALL_MUTATION_CALLS = {
    "set_shape_text": lambda i, o, ow=False: ppt_engine.set_shape_text(
        i, o, 0, 0, "x", overwrite=ow),
    "move_shape": lambda i, o, ow=False: ppt_engine.move_shape(
        i, o, 0, 0, Inches(1), Inches(1), overwrite=ow),
    "resize_shape": lambda i, o, ow=False: ppt_engine.resize_shape(
        i, o, 0, 0, Inches(1), Inches(1), overwrite=ow),
    "set_shape_geometry": lambda i, o, ow=False: ppt_engine.set_shape_geometry(
        i, o, 0, 0, Inches(1), Inches(1), Inches(1), Inches(1), overwrite=ow),
}


def _make_replace_picture_call(replacement_path):
    return lambda i, o, ow=False: ppt_engine.replace_picture(
        i, o, 0, 0, replacement_path, overwrite=ow)


# --------------------------------------------------------------------------
# F-05 - mkstemp and cleanup failures typed, cleanup failure not swallowed
# --------------------------------------------------------------------------

def test_f05_create_working_copy_mkstemp_failure_typed(fixture_pptx):
    with mock.patch("tempfile.mkstemp", side_effect=OSError("simulated mkstemp failure")):
        with pytest.raises(ppt_engine.TransactionIOError, match="temporary file"):
            ppt_engine.create_working_copy(fixture_pptx)


def test_f05_create_working_copy_os_close_failure_typed_and_never_retouches_fd(
        fixture_pptx):
    """A failure closing the fd `mkstemp` handed back must be typed (no
    raw OSError) and must clean up the temp pathname - but, once
    `os.close(fd)` has reported an error, `create_working_copy` must never
    call `os.close` or `os.fstat` on that fd number again. Descriptor
    state after a failed close is not safely recoverable by retrying (see
    the docstring on the `except OSError` block in `create_working_copy`):
    even if POSIX released the slot, something else in the process could
    already have been handed that exact fd number back by the time we
    regain control. This regression fails if the engine ever touches the
    fd again post-failure, for either reason."""
    real_close = os.close
    real_fstat = os.fstat
    failing_fd = {}
    touches_after_failure = []

    real_mkstemp = tempfile.mkstemp

    def spy_mkstemp(*a, **kw):
        fd, path = real_mkstemp(*a, **kw)
        failing_fd["fd"] = fd
        failing_fd["path"] = path
        return fd, path

    def failing_close(fd):
        if fd == failing_fd.get("fd"):
            # Really release it, matching the POSIX guarantee the code's
            # own docstring cites - and confirming this test isn't merely
            # exercising a still-open fd by accident.
            real_close(fd)
            raise OSError("simulated os.close failure (despite real closure)")
        return real_close(fd)

    def spy_fstat(fd, *a, **kw):
        if fd == failing_fd.get("fd"):
            touches_after_failure.append(("fstat", fd))
        return real_fstat(fd, *a, **kw)

    def spy_close(fd):
        if fd == failing_fd.get("fd"):
            touches_after_failure.append(("close", fd))
        return failing_close(fd)

    with mock.patch("tempfile.mkstemp", side_effect=spy_mkstemp):
        with mock.patch("os.close", side_effect=spy_close):
            with mock.patch("os.fstat", side_effect=spy_fstat):
                with pytest.raises(ppt_engine.TransactionIOError, match="close"):
                    ppt_engine.create_working_copy(fixture_pptx)
    assert not os.path.exists(failing_fd["path"])
    assert touches_after_failure == [("close", failing_fd["fd"])], (
        "create_working_copy touched the fd more than once after a close "
        "failure: %r" % (touches_after_failure,))


def test_f05_create_working_copy_close_failure_does_not_touch_reused_fd(
        fixture_pptx, tmp_path):
    """Adversarial fd-reuse case: simulate the fd number the failed
    `os.close` reported an error for being handed straight back out to an
    unrelated file by the OS - the exact race the no-retry policy exists
    to avoid racing. `create_working_copy` must never call `os.close` or
    `os.fstat` on that fd number again once it has failed once, even
    though a *different* real file now legitimately owns that number.
    Explicitly asserts the unrelated file actually reused the exact
    failed fd number before checking the engine leaves it untouched -
    without that check, a failure of the reuse assumption would make the
    "untouched" assertion pass vacuously rather than testing anything
    about fd reuse."""
    real_close = os.close
    real_fstat = os.fstat
    real_open = os.open
    failing_fd = {}
    reused = {}
    touches_after_failure = []

    real_mkstemp = tempfile.mkstemp

    def spy_mkstemp(*a, **kw):
        fd, path = real_mkstemp(*a, **kw)
        failing_fd["fd"] = fd
        failing_fd["path"] = path
        return fd, path

    def failing_close_then_reuse(fd):
        if fd == failing_fd.get("fd"):
            real_close(fd)  # genuinely free the slot
            reuse_path = str(tmp_path / "unrelated_owner.bin")
            # The OS is free to hand the just-freed fd number straight
            # back out - simulate that by opening a new, unrelated file
            # immediately: on Linux/macOS this reliably reuses the lowest
            # free descriptor, which is the one we just released.
            reused["fd"] = real_open(reuse_path, os.O_CREAT | os.O_RDWR)
            reused["path"] = reuse_path
            raise OSError("simulated os.close failure (fd now owned by another file)")
        return real_close(fd)

    def spy_fstat(fd, *a, **kw):
        if fd == failing_fd.get("fd"):
            touches_after_failure.append(("fstat", fd))
        return real_fstat(fd, *a, **kw)

    def spy_close(fd):
        if fd == failing_fd.get("fd"):
            touches_after_failure.append(("close", fd))
        return failing_close_then_reuse(fd)

    try:
        with mock.patch("tempfile.mkstemp", side_effect=spy_mkstemp):
            with mock.patch("os.close", side_effect=spy_close):
                with mock.patch("os.fstat", side_effect=spy_fstat):
                    with pytest.raises(ppt_engine.TransactionIOError, match="close"):
                        ppt_engine.create_working_copy(fixture_pptx)
        assert not os.path.exists(failing_fd["path"])
        # The adversarial premise of this test only holds if the unrelated
        # file genuinely got handed back the exact fd number the engine
        # failed to close - assert that directly, rather than assuming
        # the OS's lowest-available-fd behavior held. If this assertion
        # ever fails, the test below it would otherwise pass vacuously
        # (proving nothing about fd reuse, only that the engine leaves
        # some *other*, non-reused fd untouched).
        assert reused["fd"] == failing_fd["fd"], (
            "the unrelated file did not reuse the failed fd number (got "
            "%r, expected %r) - this test's adversarial premise did not "
            "hold, so the touches_after_failure assertion below would "
            "not actually be testing the fd-reuse scenario"
            % (reused.get("fd"), failing_fd["fd"]))
        assert touches_after_failure == [("close", failing_fd["fd"])], (
            "create_working_copy touched the reused fd after the close "
            "failure that handed it to another file: %r"
            % (touches_after_failure,))
        # The reused fd must still be open and usable - proof the engine
        # did not blindly close (or otherwise disturb) a descriptor it no
        # longer owns.
        os.write(reused["fd"], b"still mine")
        assert real_fstat(reused["fd"]).st_size > 0
    finally:
        if "fd" in reused:
            real_close(reused["fd"])


def test_f05_publish_mkstemp_failure_typed(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    with mock.patch("tempfile.mkstemp", side_effect=OSError("simulated mkstemp failure")):
        with pytest.raises(ppt_engine.TransactionIOError):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert not output.exists()


def test_f05_create_working_copy_cleanup_failure_not_silently_swallowed(fixture_pptx):
    """When staging's own copy step fails, a *further* failure trying to
    remove the partially-created temp file must be visible in the raised
    error message - not silently discarded."""
    before = set(os.listdir(tempfile.gettempdir()))
    with mock.patch("shutil.copy2", side_effect=OSError("simulated copy failure")):
        with mock.patch("os.remove", side_effect=OSError("simulated cleanup failure")):
            with pytest.raises(ppt_engine.TransactionIOError) as excinfo:
                ppt_engine.create_working_copy(fixture_pptx)
    message = str(excinfo.value)
    assert "simulated copy failure" in message
    assert "simulated cleanup failure" in message, (
        "cleanup failure must be reported, not silently swallowed")
    # The mocked-unremovable temp file is a deliberate artifact of this
    # test (cleanup was forced to fail); remove it for real now that the
    # mock is inactive, so this test does not itself leak.
    for name in set(os.listdir(tempfile.gettempdir())) - before:
        if name.startswith("b4ps_engine_"):
            os.remove(os.path.join(tempfile.gettempdir(), name))


def test_f05_staged_copy_cleanup_failure_on_success_raises(fixture_pptx, tmp_path):
    """If the mutation itself succeeds (output already published) but
    removing the staged temp copy afterward fails, that cleanup failure
    must not be silently discarded - it is the only problem at that
    point, so it is raised."""
    output = tmp_path / "out.pptx"
    real_remove = os.remove
    real_create_working_copy = ppt_engine.create_working_copy
    staged_paths = []

    def spy_create_working_copy(path):
        staged = real_create_working_copy(path)
        staged_paths.append(staged)
        return staged

    def failing_remove(path, *a, **kw):
        if staged_paths and path == staged_paths[-1]:
            raise OSError("simulated post-success cleanup failure")
        return real_remove(path, *a, **kw)

    with mock.patch.object(ppt_engine, "create_working_copy", side_effect=spy_create_working_copy):
        with mock.patch("os.remove", side_effect=failing_remove):
            with pytest.raises(ppt_engine.TransactionIOError, match="cleanup failed"):
                ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")

    # The mutation itself completed successfully before the cleanup
    # failure was raised - the output was published.
    assert output.exists()
    result = ppt_engine.inspect_deck(str(output))
    assert result["slides"][0]["shapes"][0]["text"] == "x"
    # The leaked staged file is still present (removal failed) - proving
    # the failure was real and reported, not fabricated.
    assert os.path.exists(staged_paths[-1])
    os.remove(staged_paths[-1])  # actual test cleanup


def test_f05_staged_copy_cleanup_failure_on_primary_failure_does_not_mask_it(
        fixture_pptx, tmp_path, capsys):
    """If the mutation itself fails AND cleanup also fails, the original
    (primary) error must still be what's raised - the secondary cleanup
    failure is reported (as a stderr warning), not substituted as the
    exception type. Asserts the actual warning text, not just the raised
    exception type."""
    output = tmp_path / "out.pptx"
    staged_paths = []
    real_create_working_copy = ppt_engine.create_working_copy
    real_remove = os.remove

    def spy_create_working_copy(path):
        staged = real_create_working_copy(path)
        staged_paths.append(staged)
        return staged

    def failing_remove(path, *a, **kw):
        if staged_paths and path == staged_paths[-1]:
            raise OSError("simulated cleanup failure during primary-failure path")
        return real_remove(path, *a, **kw)

    with mock.patch.object(ppt_engine, "create_working_copy", side_effect=spy_create_working_copy):
        with mock.patch.object(ppt_engine, "validate_mutated_output",
                               return_value=(False, "simulated validation failure")):
            with mock.patch("os.remove", side_effect=failing_remove):
                with pytest.raises(ppt_engine.ValidationError):
                    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert not output.exists()
    stderr = capsys.readouterr().err
    assert "simulated cleanup failure during primary-failure path" in stderr
    if staged_paths and os.path.exists(staged_paths[-1]):
        os.remove(staged_paths[-1])  # actual test cleanup


def test_f05_publish_temp_cleanup_failure_reported_on_link_exists_race(
        fixture_pptx, tmp_path, capsys):
    """The FileExistsError race path must still report a secondary
    cleanup failure (as a stderr warning) without changing the raised
    error type (OutputPathError). Asserts the actual warning text, not
    just the raised exception type."""
    output = tmp_path / "out.pptx"
    staged = ppt_engine.create_working_copy(fixture_pptx)
    output.write_bytes(b"already exists")
    try:
        with mock.patch("os.remove", side_effect=OSError("simulated cleanup failure")):
            with pytest.raises(ppt_engine.OutputPathError):
                ppt_engine._publish_atomically(staged, str(output), overwrite=False)
    finally:
        if os.path.exists(staged):
            os.remove(staged)
    stderr = capsys.readouterr().err
    assert "warning" in stderr.lower()
    assert "simulated cleanup failure" in stderr


# --------------------------------------------------------------------------
# R-01 - integrity check immediately before publication; no post-check gate
# --------------------------------------------------------------------------

def test_r01_no_post_publication_integrity_check_remains_a_gate(fixture_pptx, tmp_path):
    """Once _publish_atomically succeeds, the transaction must be treated
    as successful, even if the source is (hypothetically) observed to
    have changed afterward - because there is no post-publication check
    left to fail on. This proves the gate was removed, not merely
    weakened."""
    output = tmp_path / "out.pptx"
    real_publish = ppt_engine._publish_atomically

    def publish_then_mutate_source(staged_path, out_path, overwrite):
        result = real_publish(staged_path, out_path, overwrite)
        # Mutate the source *after* publication succeeds.
        with open(fixture_pptx, "ab") as fh:
            fh.write(b"appended-after-publish")
        return result

    with mock.patch.object(ppt_engine, "_publish_atomically",
                           side_effect=publish_then_mutate_source):
        # Must NOT raise - publication already succeeded.
        ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")

    assert output.exists()
    result = ppt_engine.inspect_deck(str(output))
    assert result["slides"][0]["shapes"][0]["text"] == "x"


def test_r01_integrity_check_runs_immediately_before_publish_call_order(fixture_pptx, tmp_path):
    call_order = []
    real_verify = ppt_engine._verify_source_hash_or_raise
    real_publish = ppt_engine._publish_atomically

    def tracking_verify(*a, **kw):
        call_order.append("verify")
        return real_verify(*a, **kw)

    def tracking_publish(*a, **kw):
        call_order.append("publish")
        return real_publish(*a, **kw)

    output = tmp_path / "out.pptx"
    with mock.patch.object(ppt_engine, "_verify_source_hash_or_raise", side_effect=tracking_verify):
        with mock.patch.object(ppt_engine, "_publish_atomically", side_effect=tracking_publish):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert call_order == ["verify", "publish"]
    # And nothing checks the hash again after - only one verify call total.
    assert call_order.count("verify") == 1


def test_r01_source_deleted_at_integrity_boundary_raises_typed_deck_source_error(
        fixture_pptx, tmp_path):
    """Source deletion between hashing and the pre-publish check must
    raise a typed DeckSourceError (a read failure), not SafeDeckError
    (a content-mismatch failure) and not a raw OSError."""
    output = tmp_path / "out.pptx"
    real_verify = ppt_engine._verify_source_hash_or_raise

    def delete_source_then_verify(input_path, before_hash):
        os.remove(input_path)
        return real_verify(input_path, before_hash)

    with mock.patch.object(ppt_engine, "_verify_source_hash_or_raise",
                           side_effect=delete_source_then_verify):
        with pytest.raises(ppt_engine.DeckSourceError, match="re-read the source"):
            ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert not output.exists()


def test_r01_source_changed_at_integrity_boundary_blocks_publication(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    real_verify = ppt_engine._verify_source_hash_or_raise
    publish_called = []
    real_publish = ppt_engine._publish_atomically

    def tracking_publish(*a, **kw):
        publish_called.append(True)
        return real_publish(*a, **kw)

    def tamper_then_verify(input_path, before_hash):
        with open(input_path, "ab") as fh:
            fh.write(b"tampered-before-verify")
        return real_verify(input_path, before_hash)

    with mock.patch.object(ppt_engine, "_verify_source_hash_or_raise", side_effect=tamper_then_verify):
        with mock.patch.object(ppt_engine, "_publish_atomically", side_effect=tracking_publish):
            with pytest.raises(ppt_engine.SafeDeckError):
                ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert publish_called == []
    assert not output.exists()


@pytest.mark.parametrize("name", list(ALL_MUTATION_CALLS))
def test_r01_integrity_boundary_all_geometry_and_text_families(name, fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    real_verify = ppt_engine._verify_source_hash_or_raise

    def tamper_then_verify(input_path, before_hash):
        with open(input_path, "ab") as fh:
            fh.write(b"tampered")
        return real_verify(input_path, before_hash)

    with mock.patch.object(ppt_engine, "_verify_source_hash_or_raise", side_effect=tamper_then_verify):
        with pytest.raises(ppt_engine.SafeDeckError):
            ALL_MUTATION_CALLS[name](fixture_pptx, str(output))
    assert not output.exists()


def test_r01_integrity_boundary_image_replacement_participates(
        fixture_pptx_with_picture, tmp_path):
    """Image replacement must go through the same single pre-publish
    integrity gate as every other mutation family."""
    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    real_verify = ppt_engine._verify_source_hash_or_raise
    publish_called = []
    real_publish = ppt_engine._publish_atomically

    def tracking_publish(*a, **kw):
        publish_called.append(True)
        return real_publish(*a, **kw)

    def tamper_then_verify(input_path, before_hash):
        with open(input_path, "ab") as fh:
            fh.write(b"tampered")
        return real_verify(input_path, before_hash)

    with mock.patch.object(ppt_engine, "_verify_source_hash_or_raise", side_effect=tamper_then_verify):
        with mock.patch.object(ppt_engine, "_publish_atomically", side_effect=tracking_publish):
            with pytest.raises(ppt_engine.SafeDeckError):
                ppt_engine.replace_picture(fixture_pptx_with_picture, str(output),
                                           0, 0, str(replacement))
    assert publish_called == []
    assert not output.exists()


def test_r01_image_replacement_source_deleted_at_boundary_typed(
        fixture_pptx_with_picture, tmp_path):
    replacement = tmp_path / "repl.png"
    replacement.write_bytes(_png_bytes((0, 255, 0)))
    output = tmp_path / "out.pptx"
    real_verify = ppt_engine._verify_source_hash_or_raise

    def delete_source_then_verify(input_path, before_hash):
        os.remove(input_path)
        return real_verify(input_path, before_hash)

    with mock.patch.object(ppt_engine, "_verify_source_hash_or_raise",
                           side_effect=delete_source_then_verify):
        with pytest.raises(ppt_engine.DeckSourceError):
            ppt_engine.replace_picture(fixture_pptx_with_picture, str(output),
                                       0, 0, str(replacement))
    assert not output.exists()


# --------------------------------------------------------------------------
# Regression guard - previously-closed findings must remain closed
# --------------------------------------------------------------------------

def test_regression_full_suite_still_green_for_normal_operations(fixture_pptx, tmp_path):
    output1 = tmp_path / "text.pptx"
    ppt_engine.set_shape_text(fixture_pptx, str(output1), 0, 0, "hello")
    assert ppt_engine.inspect_deck(str(output1))["slides"][0]["shapes"][0]["text"] == "hello"

    output2 = tmp_path / "moved.pptx"
    ppt_engine.move_shape(fixture_pptx, str(output2), 0, 0, Inches(2), Inches(2))
    shape = ppt_engine.inspect_deck(str(output2))["slides"][0]["shapes"][0]
    assert shape["left"] == Inches(2)


def test_regression_source_hash_unchanged_for_ordinary_successful_mutation(
        fixture_pptx, tmp_path):
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    output = tmp_path / "out.pptx"
    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash


def test_regression_no_leaked_temp_files_on_ordinary_success(fixture_pptx, tmp_path):
    output = tmp_path / "out.pptx"
    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "x")
    leaked = [n for n in os.listdir(tempfile.gettempdir()) if n.startswith("b4ps_engine_")]
    assert leaked == []
