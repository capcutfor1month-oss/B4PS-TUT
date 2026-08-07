"""Safe PPT Engine CLI smoke tests - proves the engine is reachable end-to-end
through b4ps.py, not just as a library. Uses a small synthetic fixture."""

import json
import os
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches

TOOLS_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
B4PS_PY = os.path.join(TOOLS_DIR, "b4ps.py")


def _run(*args):
    return subprocess.run(
        [sys.executable, B4PS_PY, *args],
        capture_output=True, text=True, cwd=TOOLS_DIR)


def _build_fixture(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text_frame.text = "hello"
    prs.save(path)


def test_engine_inspect_json(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    _build_fixture(str(fixture))

    result = _run("engine-inspect", "--input", str(fixture), "--json")
    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    # F-08: canonical structure is nested under "structure", separate from
    # the non-canonical "source_path" provenance metadata.
    assert payload["source_path"] == os.path.abspath(str(fixture))
    structure = payload["structure"]
    assert structure["slide_count"] == 1
    assert structure["slides"][0]["shapes"][0]["text"] == "hello"


def test_engine_inspect_missing_file_fails_clearly():
    result = _run("engine-inspect", "--input", "/nonexistent/path.pptx")
    assert result.returncode == 1
    assert "error:" in result.stderr
    assert "Traceback" not in result.stderr


def test_engine_set_text_end_to_end(tmp_path):
    fixture = tmp_path / "fixture.pptx"
    _build_fixture(str(fixture))
    output = tmp_path / "out.pptx"

    result = _run("engine-set-text", "--input", str(fixture), "--output", str(output),
                  "--slide", "0", "--shape", "0", "--text", "updated")
    assert result.returncode == 0, result.stderr
    assert output.exists()

    verify = _run("engine-inspect", "--input", str(output), "--json")
    payload = json.loads(verify.stdout)
    assert payload["structure"]["slides"][0]["shapes"][0]["text"] == "updated"
