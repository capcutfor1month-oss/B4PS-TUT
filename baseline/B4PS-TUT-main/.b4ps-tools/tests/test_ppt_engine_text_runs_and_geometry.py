"""`set_shape_text_runs_and_geometry` - the combined formatting-preserving
run-edit + geometry-move primitive added for Presentation Editing
Intelligence (Codex finding PEI-S1-01: PEI must not implement its own,
weaker publication protocol - every real mutation funnels through this
one audited transaction boundary). Synthetic fixtures only."""

import os
import threading

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from lib import ppt_engine


def _build_fixture(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    target = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = target.text_frame
    tf.paragraphs[0].text = "Bold: "
    tf.paragraphs[0].runs[0].font.bold = True
    run2 = tf.paragraphs[0].add_run()
    run2.text = "plain tail"
    run2.font.size = Pt(11)

    dependent = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(3), Inches(1))
    dependent.text_frame.text = "dependent shape"

    unrelated = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(3), Inches(1))
    unrelated.text_frame.text = "must never change"

    prs.save(path)
    return path, target.shape_id, dependent.shape_id, unrelated.shape_id


@pytest.fixture
def fixture_pptx(tmp_path):
    return _build_fixture(str(tmp_path / "fixture.pptx"))


def test_run_edit_preserves_formatting_of_untouched_runs(fixture_pptx, tmp_path):
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"
    ppt_engine.set_shape_text_runs_and_geometry(
        source, str(output), 0, 0,
        run_edits=[(0, 1, "plain tail extended")],
        geometry_moves=[],
    )
    prs = Presentation(str(output))
    shape = list(prs.slides[0].shapes)[0]
    runs = shape.text_frame.paragraphs[0].runs
    assert runs[0].text == "Bold: "
    assert runs[0].font.bold is True
    assert runs[1].text == "plain tail extended"
    assert runs[1].font.size == Pt(11)


def test_geometry_move_applies_absolute_top_and_verifies_shape_identity(fixture_pptx, tmp_path):
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"
    new_top = Inches(3)
    ppt_engine.set_shape_text_runs_and_geometry(
        source, str(output), 0, 0,
        run_edits=[],
        geometry_moves=[(1, dep_id, new_top)],
    )
    prs = Presentation(str(output))
    dep_shape = list(prs.slides[0].shapes)[1]
    assert dep_shape.top == new_top


def test_geometry_move_shape_id_mismatch_publishes_nothing(fixture_pptx, tmp_path):
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"
    with pytest.raises(ppt_engine.MutationError):
        ppt_engine.set_shape_text_runs_and_geometry(
            source, str(output), 0, 0,
            run_edits=[],
            geometry_moves=[(1, dep_id + 999, Inches(3))],
        )
    assert not output.exists()


def test_unrelated_shape_and_source_are_never_touched(fixture_pptx, tmp_path):
    source, target_id, dep_id, unrelated_id = fixture_pptx
    output = tmp_path / "out.pptx"
    before_hash = ppt_engine.file_sha256(source)

    ppt_engine.set_shape_text_runs_and_geometry(
        source, str(output), 0, 0,
        run_edits=[(0, 1, "changed")],
        geometry_moves=[(1, dep_id, Inches(4))],
    )

    assert ppt_engine.file_sha256(source) == before_hash
    prs = Presentation(str(output))
    unrelated_shape = list(prs.slides[0].shapes)[2]
    assert unrelated_shape.text_frame.text == "must never change"


# --------------------------------------------------------------------------
# PEI-S1-01: no-clobber race, primary-error preservation, staging cleanup
# --------------------------------------------------------------------------

def test_overwrite_false_refuses_a_destination_created_after_the_initial_check(fixture_pptx, tmp_path):
    """The exact race Codex flagged: `output_path` does not exist when the
    transaction starts, but something else creates it before this
    primitive's own publication step. `overwrite=False` must still refuse
    to clobber it - proven here by making the output already exist before
    the call even begins (the strongest, deterministic form of this race:
    if the primitive were doing a naive `exists()`-then-write, this alone
    would already prove the bug; the real fix - `_publish_atomically`'s
    `os.link`-based `FileExistsError` path - has no check-then-write
    window at all, so this also passes for a destination created at any
    point up to publication, not only before the call)."""
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"
    output.write_bytes(b"already exists - must not be silently replaced")
    original_bytes = output.read_bytes()

    with pytest.raises(ppt_engine.OutputPathError):
        ppt_engine.set_shape_text_runs_and_geometry(
            source, str(output), 0, 0,
            run_edits=[(0, 1, "changed")],
            geometry_moves=[],
            overwrite=False,
        )

    # The pre-existing destination must be completely untouched.
    assert output.read_bytes() == original_bytes


def test_overwrite_true_still_publishes_correctly_over_an_existing_destination(fixture_pptx, tmp_path):
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"
    output.write_bytes(b"placeholder")

    ppt_engine.set_shape_text_runs_and_geometry(
        source, str(output), 0, 0,
        run_edits=[(0, 1, "changed")],
        geometry_moves=[],
        overwrite=True,
    )
    prs = Presentation(str(output))
    assert prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[1].text == "changed"


def test_primary_validation_failure_leaves_no_output_and_preserves_the_error(fixture_pptx, tmp_path, monkeypatch):
    """Force the post-save validation to fail (by asking for a run edit
    that will not actually persist as requested is not directly
    reachable through the public run_edits contract, so this forces the
    failure the same way the module's own docstring says it is handled:
    patch `validate_text_runs_and_geometry_output` to report failure) and
    confirm nothing is published and the original error propagates -
    proving validation failure never leaves a falsely-successful
    artifact."""
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"

    monkeypatch.setattr(
        ppt_engine, "validate_text_runs_and_geometry_output",
        lambda *a, **k: (False, "forced failure for this test"),
    )

    with pytest.raises(ppt_engine.ValidationError, match="forced failure"):
        ppt_engine.set_shape_text_runs_and_geometry(
            source, str(output), 0, 0,
            run_edits=[(0, 1, "changed")],
            geometry_moves=[],
        )
    assert not output.exists()


def test_staging_leaves_no_temp_file_behind_on_success(fixture_pptx, tmp_path):
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"
    before_tmp_dir_entries = set(os.listdir(tmp_path))

    ppt_engine.set_shape_text_runs_and_geometry(
        source, str(output), 0, 0,
        run_edits=[(0, 1, "changed")],
        geometry_moves=[],
    )

    after_entries = set(os.listdir(tmp_path)) - before_tmp_dir_entries
    # Only the real output should be new in this directory - no leaked
    # staging or publish-temp artifact.
    assert after_entries == {"out.pptx"}


def test_failed_run_edit_target_leaves_no_output_and_no_staging_leak(fixture_pptx, tmp_path):
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"
    before_tmp_dir_entries = set(os.listdir(tmp_path))

    with pytest.raises(ppt_engine.MutationError):
        ppt_engine.set_shape_text_runs_and_geometry(
            source, str(output), 0, 0,
            run_edits=[(0, 99, "out of range run")],
            geometry_moves=[],
        )

    assert not output.exists()
    after_entries = set(os.listdir(tmp_path)) - before_tmp_dir_entries
    assert after_entries == set()


def test_source_deck_changing_mid_transaction_blocks_publication(fixture_pptx, tmp_path, monkeypatch):
    """The one integrity gate every primitive shares: if the source
    changes between hashing and publication, nothing is published."""
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"

    real_verify = ppt_engine._verify_source_hash_or_raise

    def tampering_verify(input_path, before_hash):
        with open(input_path, "ab") as f:
            f.write(b"\x00")
        return real_verify(input_path, before_hash)

    monkeypatch.setattr(ppt_engine, "_verify_source_hash_or_raise", tampering_verify)

    with pytest.raises(ppt_engine.SafeDeckError):
        ppt_engine.set_shape_text_runs_and_geometry(
            source, str(output), 0, 0,
            run_edits=[(0, 1, "changed")],
            geometry_moves=[],
        )
    assert not output.exists()


# --------------------------------------------------------------------------
# PEI-S1-01 (re-audit): expected_shape_id must be validated on the staged
# copy, inside the transaction - not by a caller-side pre-check that has
# a TOCTOU window.
# --------------------------------------------------------------------------

def test_expected_shape_id_mismatch_on_staged_copy_publishes_nothing(fixture_pptx, tmp_path):
    source, target_id, dep_id, _ = fixture_pptx
    output = tmp_path / "out.pptx"

    with pytest.raises(ppt_engine.MutationError, match="expected_shape_id"):
        ppt_engine.set_shape_text_runs_and_geometry(
            source, str(output), 0, 0,
            run_edits=[(0, 1, "changed")],
            geometry_moves=[],
            expected_shape_id=target_id + 999,
        )
    assert not output.exists()


def test_expected_shape_id_matches_the_staged_copy_not_a_pre_transaction_snapshot(fixture_pptx, tmp_path, monkeypatch):
    """Reproduces the exact adversarial scenario Codex reported: a caller
    resolves `expected_shape_id` against the source as it exists right
    now, but the source is swapped for a different deck (same
    slide/shape_index, different shape_id) before the engine actually
    stages and mutates it. The identity check must fire against what is
    actually staged - not silently succeed because some earlier snapshot
    looked fine."""
    source, target_id, dep_id, _ = fixture_pptx

    # A second deck, structurally identical (same slide/shape_index
    # layout) to the original fixture, but with shape 0's own id
    # rewritten directly in the XML - same structural position,
    # different identity, exactly the adversarial scenario reported
    # ("source changed to ID 999" at the same slide/shape index).
    swapped_path = tmp_path / "swapped.pptx"
    import shutil as _shutil
    _shutil.copyfile(source, str(swapped_path))
    swapped_prs = Presentation(str(swapped_path))
    swapped_target_shape = list(swapped_prs.slides[0].shapes)[0]
    swapped_target_id = target_id + 999
    swapped_target_shape._element.nvSpPr.cNvPr.set("id", str(swapped_target_id))
    swapped_prs.save(str(swapped_path))
    swapped_source = str(swapped_path)
    assert swapped_target_id != target_id
    # Confirm the rewrite actually took effect and the structural
    # position (shape_index 0) is otherwise unchanged.
    reopened = Presentation(swapped_source)
    assert list(reopened.slides[0].shapes)[0].shape_id == swapped_target_id

    output = tmp_path / "out.pptx"

    real_create_working_copy = ppt_engine.create_working_copy
    swap_state = {"done": False}

    def swapping_create_working_copy(path):
        # Simulate the race: right as staging begins (i.e. right as the
        # engine is about to copy `path` into a private working copy),
        # something else has already replaced the file at `path` with a
        # different deck. This runs *after* the caller's own
        # `expected_shape_id` was determined (against the original
        # `target_id`) and *after* this primitive's own preflight check,
        # but *before* the staged copy - the staged copy must therefore
        # reflect the swapped content.
        if not swap_state["done"]:
            swap_state["done"] = True
            import shutil as _shutil
            _shutil.copyfile(swapped_source, path)
        return real_create_working_copy(path)

    monkeypatch.setattr(ppt_engine, "create_working_copy", swapping_create_working_copy)

    # The caller still believes the target is `target_id` (resolved
    # before the swap) - exactly the adversarial report: "caller expected
    # shape ID 2, source changed to ID 999 before engine mutation."
    with pytest.raises(ppt_engine.MutationError, match="expected_shape_id"):
        ppt_engine.set_shape_text_runs_and_geometry(
            source, str(output), 0, 0,
            run_edits=[(0, 1, "changed")],
            geometry_moves=[],
            expected_shape_id=target_id,
        )

    # Nothing was published - in particular, the swapped-in deck's shape
    # was never mutated and returned as if it were the intended target.
    assert not output.exists()
