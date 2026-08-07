"""Safe PPT Engine tests.

All fixtures are small synthetic .pptx files built programmatically with
python-pptx. Nothing here stands in for, or pretends to be, the real
production decks (which remain unavailable - see docs/CURRENT.md).
"""

import io
import os

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from lib import ppt_engine


def _tiny_png_bytes():
    buf = io.BytesIO()
    Image.new("RGB", (2, 2), color=(255, 0, 0)).save(buf, format="PNG")
    return buf.getvalue()


def _build_fixture(path, n_slides=2, with_picture=False):
    prs = Presentation()
    for i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        box.text_frame.text = "slide %d text" % i
        if with_picture and i == 0:
            slide.shapes.add_picture(io.BytesIO(_tiny_png_bytes()), Inches(0.2), Inches(0.2))
    prs.save(path)
    return path


@pytest.fixture
def fixture_pptx(tmp_path):
    return _build_fixture(str(tmp_path / "fixture.pptx"))


@pytest.fixture
def fixture_pptx_with_picture(tmp_path):
    return _build_fixture(str(tmp_path / "fixture_pic.pptx"), with_picture=True)


# --------------------------------------------------------------------------
# loading / DeckSourceError
# --------------------------------------------------------------------------

def test_load_missing_file(tmp_path):
    with pytest.raises(ppt_engine.DeckSourceError, match="not found"):
        ppt_engine.load_deck(str(tmp_path / "does-not-exist.pptx"))


def test_load_empty_file(tmp_path):
    path = tmp_path / "empty.pptx"
    path.write_bytes(b"")
    with pytest.raises(ppt_engine.DeckSourceError, match="empty"):
        ppt_engine.load_deck(str(path))


def test_load_corrupt_file(tmp_path):
    path = tmp_path / "corrupt.pptx"
    path.write_bytes(b"this is not a zip or pptx file at all")
    with pytest.raises(ppt_engine.DeckSourceError):
        ppt_engine.load_deck(str(path))


def test_load_valid_fixture_succeeds(fixture_pptx):
    prs = ppt_engine.load_deck(fixture_pptx)
    assert len(prs.slides) == 2


# --------------------------------------------------------------------------
# inspection
# --------------------------------------------------------------------------

def test_inspect_valid_deck_metadata(fixture_pptx):
    result = ppt_engine.inspect_deck(fixture_pptx)
    assert result["slide_count"] == 2
    assert result["slide_width"] > 0
    assert result["slide_height"] > 0
    for i, slide in enumerate(result["slides"]):
        assert slide["slide_index"] == i
        assert slide["shape_count"] == 1
        shape = slide["shapes"][0]
        assert shape["text"] == "slide %d text" % i
        assert shape["has_text_frame"] is True
        assert shape["has_image"] is False
        assert isinstance(shape["left"], int)
        assert isinstance(shape["top"], int)
        assert isinstance(shape["width"], int)
        assert isinstance(shape["height"], int)


def test_inspect_reports_image_presence(fixture_pptx_with_picture):
    result = ppt_engine.inspect_deck(fixture_pptx_with_picture)
    slide0_shapes = result["slides"][0]["shapes"]
    assert len(slide0_shapes) == 2
    has_image_flags = sorted(s["has_image"] for s in slide0_shapes)
    assert has_image_flags == [False, True]


def test_inspection_does_not_modify_source(fixture_pptx):
    before = ppt_engine.file_sha256(fixture_pptx)
    ppt_engine.inspect_deck(fixture_pptx)
    after = ppt_engine.file_sha256(fixture_pptx)
    assert before == after


def test_repeated_inspection_is_deterministic(fixture_pptx):
    first = ppt_engine.inspect_deck(fixture_pptx)
    second = ppt_engine.inspect_deck(fixture_pptx)
    assert first == second


# --------------------------------------------------------------------------
# safe output path / working copy
# --------------------------------------------------------------------------

def test_output_path_must_differ_from_input(fixture_pptx):
    with pytest.raises(ppt_engine.OutputPathError, match="must not resolve"):
        ppt_engine.set_shape_text(fixture_pptx, fixture_pptx, 0, 0, "x")


def test_no_accidental_overwrite(fixture_pptx, tmp_path):
    output = tmp_path / "output.pptx"
    output.write_bytes(b"pre-existing content")
    with pytest.raises(ppt_engine.OutputPathError, match="already exists"):
        ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "new text")
    # Untouched: the pre-existing file at the output path was not replaced.
    assert output.read_bytes() == b"pre-existing content"


def test_overwrite_true_allows_replacing_existing_output(fixture_pptx, tmp_path):
    output = tmp_path / "output.pptx"
    output.write_bytes(b"stale placeholder")
    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "new text",
                              overwrite=True)
    result = ppt_engine.inspect_deck(str(output))
    assert result["slides"][0]["shapes"][0]["text"] == "new text"


def test_create_working_copy_is_a_distinct_file(fixture_pptx):
    staged = ppt_engine.create_working_copy(fixture_pptx)
    try:
        assert os.path.abspath(staged) != os.path.abspath(fixture_pptx)
        assert os.path.isfile(staged)
        assert ppt_engine.file_sha256(staged) == ppt_engine.file_sha256(fixture_pptx)
    finally:
        os.remove(staged)


# --------------------------------------------------------------------------
# controlled mutation + save/reopen validation
# --------------------------------------------------------------------------

def test_successful_mutation_persists_and_reopens(fixture_pptx, tmp_path):
    output = tmp_path / "mutated.pptx"
    ppt_engine.set_shape_text(fixture_pptx, str(output), 1, 0, "updated text")

    result = ppt_engine.inspect_deck(str(output))
    assert result["slide_count"] == 2
    assert result["slides"][1]["shapes"][0]["text"] == "updated text"
    # Unrelated fixture content (slide 0) is unchanged.
    assert result["slides"][0]["shapes"][0]["text"] == "slide 0 text"


def test_source_preserved_after_mutation(fixture_pptx, tmp_path):
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    output = tmp_path / "mutated.pptx"
    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "changed")
    after_hash = ppt_engine.file_sha256(fixture_pptx)
    assert before_hash == after_hash


def test_mutation_invalid_slide_index_raises_before_writing_output(fixture_pptx, tmp_path):
    output = tmp_path / "mutated.pptx"
    with pytest.raises(ppt_engine.MutationError, match="slide_index"):
        ppt_engine.set_shape_text(fixture_pptx, str(output), 99, 0, "x")
    assert not output.exists()


def test_mutation_invalid_shape_index_raises_before_writing_output(fixture_pptx, tmp_path):
    output = tmp_path / "mutated.pptx"
    with pytest.raises(ppt_engine.MutationError, match="shape_index"):
        ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 99, "x")
    assert not output.exists()


def test_mutation_shape_without_text_frame_raises(fixture_pptx_with_picture, tmp_path):
    # Shape 0 on slide 0 of this fixture is the picture (added first),
    # shape 1 is the text box - target the picture explicitly.
    output = tmp_path / "mutated.pptx"
    inspected = ppt_engine.inspect_deck(fixture_pptx_with_picture)
    picture_index = next(i for i, s in enumerate(inspected["slides"][0]["shapes"])
                         if s["has_image"])
    with pytest.raises(ppt_engine.MutationError, match="no text frame"):
        ppt_engine.set_shape_text(fixture_pptx_with_picture, str(output),
                                  0, picture_index, "x")
    assert not output.exists()


def test_validate_mutated_output_detects_slide_count_mismatch(fixture_pptx, tmp_path):
    output = tmp_path / "mutated.pptx"
    ppt_engine.set_shape_text(fixture_pptx, str(output), 0, 0, "y")
    ok, detail = ppt_engine.validate_mutated_output(
        str(output), expected_slide_count=99, slide_index=0, shape_index=0,
        expected_text="y")
    assert ok is False
    assert "slide count" in detail
