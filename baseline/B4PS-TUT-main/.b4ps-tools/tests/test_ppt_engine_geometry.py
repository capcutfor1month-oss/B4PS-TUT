"""Safe PPT Engine Build 2: visual mutation primitives - move, resize,
atomic geometry update. Synthetic fixtures only; no production PPTX
content is used or fabricated."""

import os

import pytest
from pptx import Presentation
from pptx.util import Inches

from lib import ppt_engine

EMU_PER_INCH = 914400


def _build_fixture(path, n_slides=2):
    prs = Presentation()
    for i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        box.text_frame.text = "slide %d text" % i
    prs.save(path)
    return path


@pytest.fixture
def fixture_pptx(tmp_path):
    return _build_fixture(str(tmp_path / "fixture.pptx"))


# --------------------------------------------------------------------------
# move_shape
# --------------------------------------------------------------------------

def test_move_shape_success_and_reopened_coordinates_match(fixture_pptx, tmp_path):
    output = tmp_path / "moved.pptx"
    new_left, new_top = Inches(2), Inches(2.5)
    ppt_engine.move_shape(fixture_pptx, str(output), 0, 0, new_left, new_top)

    result = ppt_engine.inspect_deck(str(output))
    shape = result["slides"][0]["shapes"][0]
    assert shape["left"] == new_left
    assert shape["top"] == new_top


def test_move_shape_does_not_change_width_or_height(fixture_pptx, tmp_path):
    before = ppt_engine.inspect_deck(fixture_pptx)["slides"][0]["shapes"][0]
    output = tmp_path / "moved.pptx"
    ppt_engine.move_shape(fixture_pptx, str(output), 0, 0, Inches(2), Inches(2))

    after = ppt_engine.inspect_deck(str(output))["slides"][0]["shapes"][0]
    assert after["width"] == before["width"]
    assert after["height"] == before["height"]


def test_move_shape_source_unchanged(fixture_pptx, tmp_path):
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    output = tmp_path / "moved.pptx"
    ppt_engine.move_shape(fixture_pptx, str(output), 0, 0, Inches(2), Inches(2))
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash


def test_move_shape_invalid_slide_index(fixture_pptx, tmp_path):
    output = tmp_path / "moved.pptx"
    with pytest.raises(ppt_engine.MutationError, match="slide_index"):
        ppt_engine.move_shape(fixture_pptx, str(output), 99, 0, Inches(1), Inches(1))
    assert not output.exists()


def test_move_shape_invalid_shape_index(fixture_pptx, tmp_path):
    output = tmp_path / "moved.pptx"
    with pytest.raises(ppt_engine.MutationError, match="shape_index"):
        ppt_engine.move_shape(fixture_pptx, str(output), 0, 99, Inches(1), Inches(1))
    assert not output.exists()


def test_move_shape_rejects_negative_coordinates(fixture_pptx, tmp_path):
    output = tmp_path / "moved.pptx"
    with pytest.raises(ppt_engine.MutationError, match="left"):
        ppt_engine.move_shape(fixture_pptx, str(output), 0, 0, -1, Inches(1))
    assert not output.exists()


def test_move_shape_rejects_non_integer_coordinates(fixture_pptx, tmp_path):
    output = tmp_path / "moved.pptx"
    with pytest.raises(ppt_engine.MutationError):
        ppt_engine.move_shape(fixture_pptx, str(output), 0, 0, 1.5, Inches(1))
    assert not output.exists()


# --------------------------------------------------------------------------
# resize_shape
# --------------------------------------------------------------------------

def test_resize_shape_success_and_reopened_dimensions_match(fixture_pptx, tmp_path):
    output = tmp_path / "resized.pptx"
    new_width, new_height = Inches(4), Inches(1.5)
    ppt_engine.resize_shape(fixture_pptx, str(output), 1, 0, new_width, new_height)

    result = ppt_engine.inspect_deck(str(output))
    shape = result["slides"][1]["shapes"][0]
    assert shape["width"] == new_width
    assert shape["height"] == new_height


def test_resize_shape_does_not_change_position(fixture_pptx, tmp_path):
    before = ppt_engine.inspect_deck(fixture_pptx)["slides"][0]["shapes"][0]
    output = tmp_path / "resized.pptx"
    ppt_engine.resize_shape(fixture_pptx, str(output), 0, 0, Inches(5), Inches(2))

    after = ppt_engine.inspect_deck(str(output))["slides"][0]["shapes"][0]
    assert after["left"] == before["left"]
    assert after["top"] == before["top"]


def test_resize_shape_source_unchanged(fixture_pptx, tmp_path):
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    output = tmp_path / "resized.pptx"
    ppt_engine.resize_shape(fixture_pptx, str(output), 0, 0, Inches(4), Inches(2))
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash


def test_resize_shape_rejects_zero_dimension(fixture_pptx, tmp_path):
    output = tmp_path / "resized.pptx"
    with pytest.raises(ppt_engine.MutationError, match="width"):
        ppt_engine.resize_shape(fixture_pptx, str(output), 0, 0, 0, Inches(1))
    assert not output.exists()


def test_resize_shape_rejects_negative_dimension(fixture_pptx, tmp_path):
    output = tmp_path / "resized.pptx"
    with pytest.raises(ppt_engine.MutationError, match="height"):
        ppt_engine.resize_shape(fixture_pptx, str(output), 0, 0, Inches(1), -100)
    assert not output.exists()


# --------------------------------------------------------------------------
# set_shape_geometry (atomic left+top+width+height)
# --------------------------------------------------------------------------

def test_set_shape_geometry_success(fixture_pptx, tmp_path):
    output = tmp_path / "geom.pptx"
    geometry = dict(left=Inches(0.5), top=Inches(0.5), width=Inches(2), height=Inches(1))
    ppt_engine.set_shape_geometry(fixture_pptx, str(output), 0, 0, **geometry)

    shape = ppt_engine.inspect_deck(str(output))["slides"][0]["shapes"][0]
    assert shape["left"] == geometry["left"]
    assert shape["top"] == geometry["top"]
    assert shape["width"] == geometry["width"]
    assert shape["height"] == geometry["height"]


def test_set_shape_geometry_source_unchanged(fixture_pptx, tmp_path):
    before_hash = ppt_engine.file_sha256(fixture_pptx)
    output = tmp_path / "geom.pptx"
    ppt_engine.set_shape_geometry(fixture_pptx, str(output), 0, 0,
                                  Inches(1), Inches(1), Inches(2), Inches(1))
    assert ppt_engine.file_sha256(fixture_pptx) == before_hash


def test_set_shape_geometry_invalid_value_writes_nothing_and_leaves_no_partial_output(
        fixture_pptx, tmp_path):
    output = tmp_path / "geom.pptx"
    # width is invalid (zero) - must fail before any staging/writing happens,
    # and must not leave a partially-applied geometry anywhere.
    with pytest.raises(ppt_engine.MutationError, match="width"):
        ppt_engine.set_shape_geometry(fixture_pptx, str(output), 0, 0,
                                      Inches(1), Inches(1), 0, Inches(1))
    assert not output.exists()
    # Nothing was left behind in the temp directory used for staging either.
    import tempfile
    leftover = [f for f in os.listdir(tempfile.gettempdir())
               if f.startswith("b4ps_engine_") or f.startswith(".b4ps_publish_")]
    assert leftover == []


def test_set_shape_geometry_all_four_persist_together_not_partially(fixture_pptx, tmp_path):
    output = tmp_path / "geom.pptx"
    ppt_engine.set_shape_geometry(fixture_pptx, str(output), 0, 0,
                                  Inches(1.1), Inches(1.2), Inches(1.3), Inches(1.4))
    shape = ppt_engine.inspect_deck(str(output))["slides"][0]["shapes"][0]
    assert (shape["left"], shape["top"], shape["width"], shape["height"]) == (
        Inches(1.1), Inches(1.2), Inches(1.3), Inches(1.4))


# --------------------------------------------------------------------------
# atomic publish / no leftover partial output on failure
# --------------------------------------------------------------------------

def test_output_path_error_leaves_no_output_file(fixture_pptx, tmp_path):
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"placeholder")
    with pytest.raises(ppt_engine.OutputPathError):
        ppt_engine.move_shape(fixture_pptx, str(output), 0, 0, Inches(1), Inches(1))
    assert output.read_bytes() == b"placeholder"


def test_no_publish_temp_files_left_behind_on_success(fixture_pptx, tmp_path):
    output = tmp_path / "moved.pptx"
    ppt_engine.move_shape(fixture_pptx, str(output), 0, 0, Inches(1), Inches(1))
    leftovers = [f for f in os.listdir(str(tmp_path)) if f.startswith(".b4ps_publish_")]
    assert leftovers == []
