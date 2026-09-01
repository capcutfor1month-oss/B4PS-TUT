"""Safe PPT Engine: deterministic, non-destructive primitives for reading,
inspecting, and controlled-mutating .pptx files, built on python-pptx.

Transaction contract every mutation primitive follows:

    validate source
    -> validate target
    -> validate source/output alias safety
    -> hash source
    -> stage privately
    -> mutate
    -> save private result
    -> reopen and validate result
    -> re-verify source hash unchanged (the ONE integrity gate, immediately
       before publication - not after; a source deletion/read error here
       is itself a typed failure)
    -> publish atomically according to overwrite policy
    -> cleanup (a cleanup failure is reported, never silently discarded -
       see `_staged_copy`)

Once `_publish_atomically` returns successfully, the transaction is
successful - there is no post-publication integrity check. No expected
failure path may alter the source, overwrite an existing destination when
overwrite=False, expose a partial final destination, leak a temp file
without reporting it, or leak a raw (untyped) exception through the CLI.

This module describes deck *structure* only - it assigns no semantic
meaning to any shape (e.g. "this is the instructional red box"). That is
future Documentation Intelligence work, not this layer.

Every public function that touches a real path raises a typed SafeDeckError
subclass on failure rather than letting a bare library exception surface, so
callers (including the CLI) get one clear, actionable reason.
"""

import contextlib
import hashlib
import os
import shutil
import sys
import tempfile

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError
from pptx.oxml.ns import qn


class SafeDeckError(Exception):
    """Base for all typed Safe PPT Engine failures."""


class DeckSourceError(SafeDeckError):
    """Input .pptx is missing, empty, or not a valid package."""


class OutputPathError(SafeDeckError):
    """Requested output path is unsafe: aliases the input (directly, via a
    symlink, or via any other filesystem alias), or already exists without
    explicit overwrite, or was created by something else during the
    operation."""


class MutationError(SafeDeckError):
    """A controlled mutation could not be applied as requested (bad target,
    unsupported shape, etc.) - raised before anything is written out."""


class ValidationError(SafeDeckError):
    """Generated output failed post-save reopen/content validation."""


class TransactionIOError(SafeDeckError):
    """A filesystem or library I/O failure occurred inside a mutation
    transaction (working-copy creation, save, temporary-output
    preparation, or publication) - normalized here so it never leaks as a
    raw OSError/library exception. Distinct from DeckSourceError (bad
    *input*) and OutputPathError (unsafe *destination*): this is for
    failures during the transaction's own mechanics, not bad arguments."""


# --------------------------------------------------------------------------
# loading
# --------------------------------------------------------------------------

def file_sha256(path):
    h = hashlib.sha256()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(1 << 20), b""):
            h.update(chunk)
    return h.hexdigest()


def load_deck(path):
    """Open `path` read-only into memory as a python-pptx Presentation.

    Never modifies `path`. Never falls back to any other file - a bad path
    always fails, it is never silently substituted. Raises DeckSourceError
    with a clear reason on any failure - including a readable ZIP archive
    that simply isn't a PPTX package (python-pptx raises a plain KeyError
    for that case; it is normalized here like every other open failure).
    This is the typed validation boundary every mutation primitive must
    pass *before* hashing, staging, or any other filesystem access on the
    source.
    """
    if not os.path.exists(path):
        raise DeckSourceError("PPTX not found: %s" % path)
    if os.path.getsize(path) == 0:
        raise DeckSourceError("PPTX is empty (0 bytes): %s" % path)
    try:
        return Presentation(path)
    except PackageNotFoundError as exc:
        raise DeckSourceError("Not a valid .pptx package: %s (%s)" % (path, exc))
    except Exception as exc:
        # python-pptx / the underlying zip layer can raise a variety of
        # exception types on a corrupt or non-PPTX package (BadZipFile,
        # KeyError for a missing required part, etc.); normalize them all
        # to one typed, actionable failure rather than leaking the
        # internal type.
        raise DeckSourceError("Could not open .pptx: %s (%s: %s)"
                               % (path, type(exc).__name__, exc))


def _validate_and_hash_source(input_path):
    """Validate `input_path` through the typed loading boundary, then hash
    it. Validation always runs first, so a missing/empty/corrupt input
    raises DeckSourceError here rather than `file_sha256` raising a bare
    OSError first."""
    load_deck(input_path)
    return file_sha256(input_path)


# --------------------------------------------------------------------------
# inspection
# --------------------------------------------------------------------------

def _shape_type_name(shape):
    try:
        return str(shape.shape_type)
    except Exception:
        return "UNKNOWN"


def _shape_text(shape):
    if not shape.has_text_frame:
        return None
    return shape.text_frame.text


def _shape_info(shape, z_order):
    return {
        "shape_index": z_order,
        "z_order": z_order,
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": _shape_type_name(shape),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
        "has_text_frame": shape.has_text_frame,
        "text": _shape_text(shape),
        "has_image": shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
    }


def inspect_deck(path):
    """Deterministic structural inspection of the deck at `path`.

    Describes what exists - slide/shape counts, dimensions, position, text
    presence, image presence, z-order, and stable identifiers where the
    format provides them (shape_id). Assigns no semantic meaning.

    Canonical and deterministic: two byte-identical .pptx files at
    different filesystem paths produce an *equal* result, and repeated
    calls on the same file produce an equal result (no timestamps, temp
    paths, random IDs, or the source path itself are included - the path is
    provenance metadata, not deck structure; see `describe_deck` for a
    variant that adds it).
    """
    prs = load_deck(path)
    slides = []
    for slide_index, slide in enumerate(prs.slides):
        shapes = [_shape_info(shape, z) for z, shape in enumerate(slide.shapes)]
        slides.append({
            "slide_index": slide_index,
            "shape_count": len(shapes),
            "shapes": shapes,
        })
    return {
        "slide_count": len(slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": slides,
    }


def describe_deck(path):
    """`inspect_deck(path)` plus non-canonical source provenance metadata,
    for callers (such as the CLI) that want to know *which file* produced a
    given structural result. `result["structure"]` is always exactly what
    `inspect_deck(path)` would return - safe to compare across paths;
    `result["source_path"]` is deliberately excluded from that comparison.
    """
    return {
        "source_path": os.path.abspath(path),
        "structure": inspect_deck(path),
    }


# --------------------------------------------------------------------------
# safe working copy / output-path safety
# --------------------------------------------------------------------------

def _assert_safe_output_path(input_path, output_path, overwrite):
    """Reject an output path that aliases the input, through any of:
    identical path strings, a symlink (in either direction) resolving to
    the same file, or any other filesystem-level alias `os.path.samefile`
    can detect (hardlinks, bind mounts, etc.). `os.path.realpath` resolves
    symlinks in every existing path component even when the final
    component itself does not yet exist, so this is safe to call
    regardless of whether `output_path` exists.

    This is the primary protection against source aliasing - it does not
    depend on the source-hash check callers also perform.
    """
    in_real = os.path.realpath(input_path)
    out_real = os.path.realpath(output_path)
    if in_real == out_real:
        raise OutputPathError(
            "Output path resolves to the same file as the input path "
            "(directly, or via a symlink/alias): %s" % output_path)
    if os.path.exists(input_path) and os.path.exists(output_path):
        try:
            if os.path.samefile(input_path, output_path):
                raise OutputPathError(
                    "Output path is the same underlying file as the input "
                    "path (same device/inode): %s" % output_path)
        except OSError:
            # Both existed a moment ago per os.path.exists above, but a
            # concurrent removal could still race the stat inside samefile.
            # The realpath comparison above already covers the common case;
            # this is defense-in-depth, not the sole protection.
            pass
    if os.path.exists(out_real) and not overwrite:
        raise OutputPathError(
            "Output path already exists (pass overwrite=True to replace it "
            "deliberately): %s" % output_path)


def create_working_copy(path):
    """A private staging copy in a fresh temp location. All mutation work
    happens here - never on the original, and never directly at the
    caller's requested output path - so a failed or partial mutation can
    never corrupt either. Caller is responsible for removing the returned
    path when done (see `_staged_copy`, which does this for every
    mutation primitive). Exception-safe: allocating the temp path
    (`tempfile.mkstemp`), closing its file descriptor, and copying into it
    are all wrapped, and any failure - including a failure while cleaning
    up after a failed close or copy - is normalized into a typed
    `TransactionIOError` rather than a raw OSError or a silently discarded
    cleanup failure. A close failure is not retried: once `os.close(fd)`
    has reported an error, that fd number's state is not safely
    recoverable from this code - see the note at the `except` block below
    for why - so we only clean up what we still know how to clean up (the
    temp pathname) and raise a typed error.
    """
    load_deck(path)  # validate before copying anything
    try:
        fd, staged = tempfile.mkstemp(prefix="b4ps_engine_", suffix=".pptx")
    except OSError as exc:
        raise TransactionIOError(
            "Could not create a temporary file for staging %s: %s" % (path, exc))
    try:
        os.close(fd)
    except OSError as exc:
        # Once close(fd) has reported an error, fd's state is ambiguous:
        # POSIX says the descriptor slot is deallocated on return
        # regardless of the error, but real platforms and filesystems are
        # not all POSIX-compliant enough to trust that uniformly - and even
        # where it holds, by the time we get back control someone else in
        # the process (another thread, a signal handler, a library, the
        # interpreter's own GC finalizers) may already have opened a new
        # file and been handed this exact fd number back. There is no way
        # from here to distinguish "still our open fd" from "already
        # reused by something else" without racing that possibility, so we
        # deliberately do not call fstat(fd) or close(fd) again - retrying
        # would risk querying or closing a file descriptor this function no
        # longer owns. We only clean up what does not require touching fd
        # again: the temp pathname itself.
        detail = _remove_reporting_failure(staged)
        raise TransactionIOError(
            "Could not close the temporary file allocated for staging %s: %s%s"
            % (path, exc, detail))
    try:
        shutil.copy2(path, staged)
    except OSError as exc:
        detail = _remove_reporting_failure(staged)
        raise TransactionIOError(
            "Could not stage a working copy of %s: %s%s" % (path, exc, detail))
    return staged


def _remove_reporting_failure(path):
    """Remove `path` if it exists. Returns an empty string on success (or
    if there was nothing to remove) or a parenthetical describing the
    removal failure otherwise - callers fold this into the message of the
    typed error they are already raising, so a cleanup failure is always
    visible to the caller rather than silently discarded. Never itself
    raises - it is called while another failure is already being
    reported, and must not replace or hide that failure.
    """
    try:
        if path and os.path.exists(path):
            os.remove(path)
    except OSError as exc:
        return " (additionally, could not remove temporary file %s: %s)" % (path, exc)
    return ""


@contextlib.contextmanager
def _staged_copy(input_path):
    """Context manager wrapping `create_working_copy`: yields the staged
    path, and on exit removes it. A cleanup failure is never silently
    swallowed:

    - If the wrapped block itself raised, that is the primary failure -
      a secondary cleanup failure must not replace or mask it, but it is
      still surfaced, as a warning on stderr, rather than discarded.
    - If the wrapped block succeeded, there is no primary failure to
      protect - a cleanup failure here is itself the only problem, so it
      is raised as a typed `TransactionIOError` rather than ignored,
      since a leaked temp file is a real, reportable outcome.
    """
    staged = create_working_copy(input_path)
    try:
        yield staged
    except BaseException:
        detail = _remove_reporting_failure(staged)
        if detail:
            print("warning: cleanup after a failed operation also failed%s"
                  % detail, file=sys.stderr)
        raise
    else:
        detail = _remove_reporting_failure(staged)
        if detail:
            raise TransactionIOError(
                "Operation succeeded but temporary working copy cleanup "
                "failed%s" % detail)


def _publish_atomically(staged_path, output_path, overwrite):
    """Place a validated, already-saved staged file at `output_path`
    according to the overwrite policy. Never emulates atomicity with a
    separate exists-check followed by a write - that has a race window
    between the check and the write.

    Both branches first copy the complete staged content into a same-
    directory temp file and fsync it, so the *entire* file exists on disk
    before any filesystem operation makes it visible at `output_path` - the
    final path can never be observed holding a partial file.

    overwrite=True: `os.replace(tmp, output_path)` - an atomic rename on
    the same filesystem, allowed to replace whatever is already there.

    overwrite=False: `os.link(tmp, output_path)` - an atomic hard link
    that fails with `FileExistsError` if `output_path` already exists,
    without ever opening or writing to `output_path` itself first (unlike
    an O_CREAT|O_EXCL stream, which creates the destination filename
    before any bytes are written into it). The temp file is unlinked
    afterward in both the success and failure case, so it is never left
    behind; the underlying data survives via the new hard link on success.
    """
    out_dir = os.path.dirname(os.path.abspath(output_path)) or "."
    try:
        os.makedirs(out_dir, exist_ok=True)
    except OSError as exc:
        raise TransactionIOError(
            "Could not prepare output directory %s: %s" % (out_dir, exc))

    try:
        fd, tmp_out = tempfile.mkstemp(prefix=".b4ps_publish_", suffix=".pptx", dir=out_dir)
    except OSError as exc:
        raise TransactionIOError(
            "Could not create a temporary file for publication in %s: %s"
            % (out_dir, exc))
    try:
        with os.fdopen(fd, "wb") as dst:
            with open(staged_path, "rb") as src:
                shutil.copyfileobj(src, dst)
            dst.flush()
            os.fsync(dst.fileno())
    except OSError as exc:
        detail = _remove_reporting_failure(tmp_out)
        raise TransactionIOError(
            "Could not prepare output for publication: %s%s" % (exc, detail))

    if overwrite:
        try:
            os.replace(tmp_out, output_path)
        except OSError as exc:
            detail = _remove_reporting_failure(tmp_out)
            raise TransactionIOError("Could not publish output: %s%s" % (exc, detail))
        return

    try:
        os.link(tmp_out, output_path)
    except FileExistsError:
        detail = _remove_reporting_failure(tmp_out)
        if detail:
            print("warning: cleanup after a failed publish also failed%s"
                  % detail, file=sys.stderr)
        raise OutputPathError(
            "Output path was created by something else during the operation "
            "(pass overwrite=True to replace it deliberately): %s" % output_path)
    except OSError as exc:
        detail = _remove_reporting_failure(tmp_out)
        raise TransactionIOError("Could not publish output: %s%s" % (exc, detail))
    else:
        detail = _remove_reporting_failure(tmp_out)
        if detail:
            # Publication itself succeeded (the hard link exists) - a
            # failure to remove the now-redundant temp name is a real,
            # reportable problem in its own right, not one to discard.
            raise TransactionIOError(
                "Output published, but cleanup of the publication temp "
                "file failed%s" % detail)


def _verify_source_hash_or_raise(input_path, before_hash):
    """Re-read and re-hash `input_path` for the pre-publish integrity
    check, and compare it against `before_hash`. A source that has been
    deleted, or become unreadable, since it was first hashed is itself a
    typed failure (`DeckSourceError`) distinct from a source that is
    merely readable but *changed* (`SafeDeckError`, the safety-invariant
    violation) - both are reported clearly rather than surfacing as a
    bare OSError from `file_sha256`.
    """
    try:
        current_hash = file_sha256(input_path)
    except OSError as exc:
        raise DeckSourceError(
            "Could not re-read the source to verify its integrity "
            "immediately before publication: %s (%s)" % (input_path, exc))
    if current_hash != before_hash:
        raise SafeDeckError(
            "input file changed during the transaction - refusing to publish "
            "(safety invariant violated): %s" % input_path)


def _finalize_transaction(input_path, output_path, staged, overwrite, before_hash):
    """Shared tail of every mutation transaction, run only after the
    mutation has been saved to `staged` and validated by reopening it:

        verify source hash unchanged (the one gate - runs immediately
        before publication is attempted; a source deletion/read error at
        this point is itself a typed failure)
        -> publish atomically according to overwrite policy

    If the source changed - or became unreadable - after it was hashed
    and before this runs, nothing is ever published: the integrity
    failure is raised and `output_path` is left exactly as it was. Once
    `_publish_atomically` returns successfully, the transaction is
    successful - there is no check afterward. Re-hashing the source after
    a successful publish cannot undo that publish, and treating a
    post-hoc difference as a transaction failure would be misleading: the
    output was correctly and safely produced from the source exactly as
    it existed at the moment of publication, which is the guarantee this
    function actually makes.
    """
    _verify_source_hash_or_raise(input_path, before_hash)
    _publish_atomically(staged, output_path, overwrite)


# --------------------------------------------------------------------------
# controlled mutation
# --------------------------------------------------------------------------

def _resolve_shape(prs, slide_index, shape_index):
    if not 0 <= slide_index < len(prs.slides):
        raise MutationError("slide_index %d out of range (deck has %d slide(s))"
                             % (slide_index, len(prs.slides)))
    slide = prs.slides[slide_index]
    shapes = list(slide.shapes)
    if not 0 <= shape_index < len(shapes):
        raise MutationError("shape_index %d out of range (slide %d has %d shape(s))"
                             % (shape_index, slide_index, len(shapes)))
    return shapes[shape_index]


def _preflight_resolve_shape(input_path, slide_index, shape_index):
    """Resolve the target shape against the *original* input before any
    working copy is staged, so a bad slide/shape index (or, where the
    caller checks shape_type itself, a wrong shape kind) fails with no
    temp file ever created. Returns the shape from a throwaway load - it
    is not used for mutation, only for preflight validation."""
    prs = load_deck(input_path)
    return _resolve_shape(prs, slide_index, shape_index)


def validate_mutated_output(output_path, expected_slide_count, slide_index,
                             shape_index, expected_text):
    """Reopen `output_path` from scratch and confirm: it is still a valid
    PPTX package, slide count is unchanged, and the targeted shape's text
    now matches what was requested. Returns (ok, detail)."""
    try:
        prs = load_deck(output_path)
    except DeckSourceError as exc:
        return False, "generated output failed to reopen: %s" % exc
    if len(prs.slides) != expected_slide_count:
        return False, ("slide count changed: expected %d, got %d"
                        % (expected_slide_count, len(prs.slides)))
    try:
        shape = _resolve_shape(prs, slide_index, shape_index)
    except MutationError as exc:
        return False, "targeted shape missing after save: %s" % exc
    if not shape.has_text_frame or shape.text_frame.text != expected_text:
        return False, "mutation did not persist as expected"
    return True, "ok"


def _save_staged(prs, staged):
    """Save the in-memory mutated presentation to the staged path,
    normalizing any filesystem/library failure into a typed error."""
    try:
        prs.save(staged)
    except OSError as exc:
        raise TransactionIOError(
            "Could not save the mutated deck to a temporary file: %s" % exc)


def set_shape_text(input_path, output_path, slide_index, shape_index, new_text,
                    overwrite=False):
    """Controlled mutation primitive: set one explicitly targeted shape's
    text-frame text. Targeting is deterministic by (slide_index,
    shape_index) only - no fuzzy or semantic matching.

    Never mutates `input_path`. Stages the change on a private working
    copy, validates the saved result by reopening it, verifies the source
    is still unchanged, and only then places the result at `output_path`.
    On any failure, nothing is written to `output_path` and `input_path`
    is left untouched.
    """
    _assert_safe_output_path(input_path, output_path, overwrite)
    target = _preflight_resolve_shape(input_path, slide_index, shape_index)
    if not target.has_text_frame:
        raise MutationError(
            "shape %d on slide %d has no text frame - cannot set text"
            % (shape_index, slide_index))
    before_hash = _validate_and_hash_source(input_path)

    with _staged_copy(input_path) as staged:
        prs = Presentation(staged)
        shape = _resolve_shape(prs, slide_index, shape_index)
        shape.text_frame.text = new_text
        expected_slide_count = len(prs.slides)
        _save_staged(prs, staged)

        ok, detail = validate_mutated_output(
            staged, expected_slide_count, slide_index, shape_index, new_text)
        if not ok:
            raise ValidationError(detail)

        _finalize_transaction(input_path, output_path, staged, overwrite, before_hash)

    return output_path


def validate_text_runs_and_geometry_output(output_path, expected_slide_count, slide_index,
                                            shape_index, run_edits, geometry_moves):
    """Reopen `output_path` from scratch and confirm: it is still a valid
    PPTX package, slide count is unchanged, every requested run edit
    persisted exactly, and every requested geometry move landed on the
    exact shape it was aimed at (by `shape_id`, not merely by index) with
    the exact requested `top`. Returns (ok, detail)."""
    try:
        prs = load_deck(output_path)
    except DeckSourceError as exc:
        return False, "generated output failed to reopen: %s" % exc
    if len(prs.slides) != expected_slide_count:
        return False, ("slide count changed: expected %d, got %d"
                        % (expected_slide_count, len(prs.slides)))
    try:
        shape = _resolve_shape(prs, slide_index, shape_index)
    except MutationError as exc:
        return False, "targeted shape missing after save: %s" % exc
    for p_idx, r_idx, expected_text in run_edits:
        try:
            actual_text = shape.text_frame.paragraphs[p_idx].runs[r_idx].text
        except IndexError:
            return False, "run edit target missing after save: paragraph %d run %d" % (p_idx, r_idx)
        if actual_text != expected_text:
            return False, ("run edit did not persist: paragraph %d run %d expected %r, got %r"
                            % (p_idx, r_idx, expected_text, actual_text))
    shapes = list(prs.slides[slide_index].shapes)
    for move_shape_index, move_shape_id, expected_top in geometry_moves:
        if not 0 <= move_shape_index < len(shapes):
            return False, "geometry move target missing after save: shape_index %d" % move_shape_index
        moved = shapes[move_shape_index]
        if moved.shape_id != move_shape_id:
            return False, ("geometry move target identity changed after save at shape_index %d: "
                            "expected shape_id %r, found %r" % (move_shape_index, move_shape_id, moved.shape_id))
        if moved.top != expected_top:
            return False, ("geometry move did not persist at shape_index %d: expected top=%r, got %r"
                            % (move_shape_index, expected_top, moved.top))
    return True, "ok"


def set_shape_text_runs_and_geometry(
    input_path, output_path, slide_index, shape_index,
    run_edits, geometry_moves=None, overwrite=False, expected_shape_id=None,
):
    """Controlled mutation primitive built for Presentation Editing
    Intelligence's scene-aware edits: applies a set of already-planned,
    formatting-preserving run-level text edits to one target shape, plus
    a set of already-planned geometry moves to other shapes on the same
    slide, as **one** atomic transaction.

    This exists because a text edit and its dependent local reflow must
    land in a single save - this file's other per-primitive functions are
    each individually built for exactly one shape's one kind of change
    per file-to-file call, so composing two of them for one edit would
    require two separate mutation passes (two separate publish steps,
    each independently observable, each its own no-clobber race window)
    for what the caller intends as one edit. This primitive follows the
    identical transaction discipline documented at the top of this file -
    the same `_staged_copy`/`_finalize_transaction`/`_publish_atomically`
    machinery every other primitive uses, including the same race-free,
    `os.link`-based `overwrite=False` no-clobber publication (Codex
    finding PEI-S1-01: no second, weaker publication protocol exists
    anywhere in Presentation Editing Intelligence - every mutation this
    repository performs against a real `.pptx` funnels through this one
    audited boundary).

    `run_edits`: iterable of `(paragraph_index, run_index, new_text)` -
    only that run's own `<a:t>` text node is touched; `Run.text`'s setter
    leaves every other run's formatting (`<a:rPr>`: bold, italic,
    underline, size, color, font name) completely untouched, and every
    run not named here is left completely alone.

    `geometry_moves`: iterable of `(shape_index, shape_id, new_top)` -
    `new_top` is the shape's complete new `top` in EMU (absolute, not a
    delta - matching `move_shape`'s own convention), validated the same
    way `move_shape` validates its own `top` argument. The supplied
    `shape_id` must match the shape actually found at `shape_index` or
    the whole transaction fails before anything is staged for that move -
    this primitive never guesses or retargets.

    `expected_shape_id`, if given, is the text target's own identity
    guard (Codex finding PEI-S1-01, re-audit): it is validated against
    the shape actually resolved on the **staged copy**, immediately
    before any run edit is applied - not against a separate read-only
    load taken before staging. A caller-side pre-transaction check has a
    TOCTOU window (the source could be swapped between that check and
    this call actually staging/mutating it); validating on the staged
    copy, inside this same transaction, closes that window - the staged
    copy is what gets mutated, so checking its identity is checking the
    identity of the thing that is actually about to change.

    Never mutates `input_path`. On any failure (bad target, invalid
    coordinate, an `expected_shape_id` mismatch on the text target, a
    `shape_id` mismatch on any geometry move, save/validation/integrity
    failure, or the output path already existing with `overwrite=False`),
    nothing is written to `output_path` and `input_path` is left
    untouched.
    """
    run_edits = list(run_edits)
    geometry_moves = list(geometry_moves or [])
    for _, _, new_top in geometry_moves:
        _validate_coordinate(new_top, "top")

    _assert_safe_output_path(input_path, output_path, overwrite)
    target = _preflight_resolve_shape(input_path, slide_index, shape_index)
    if run_edits and not target.has_text_frame:
        raise MutationError(
            "shape %d on slide %d has no text frame - cannot edit text"
            % (shape_index, slide_index))
    before_hash = _validate_and_hash_source(input_path)

    with _staged_copy(input_path) as staged:
        prs = Presentation(staged)
        shape_target = _resolve_shape(prs, slide_index, shape_index)

        if expected_shape_id is not None and shape_target.shape_id != expected_shape_id:
            raise MutationError(
                "expected_shape_id guard mismatch at slide %d shape %d: expected %r, found %r "
                "on the staged copy - refusing to mutate the wrong shape"
                % (slide_index, shape_index, expected_shape_id, shape_target.shape_id)
            )

        paragraphs = shape_target.text_frame.paragraphs if shape_target.has_text_frame else []
        for p_idx, r_idx, new_text in run_edits:
            if not 0 <= p_idx < len(paragraphs):
                raise MutationError(
                    "paragraph_index %d out of range (shape has %d paragraph(s))"
                    % (p_idx, len(paragraphs)))
            runs = paragraphs[p_idx].runs
            if not 0 <= r_idx < len(runs):
                raise MutationError(
                    "run_index %d out of range in paragraph %d (has %d run(s))"
                    % (r_idx, p_idx, len(runs)))
            runs[r_idx].text = new_text

        shapes = list(prs.slides[slide_index].shapes)
        for move_shape_index, move_shape_id, new_top in geometry_moves:
            if not 0 <= move_shape_index < len(shapes):
                raise MutationError(
                    "geometry move shape_index %d out of range (slide %d has %d shape(s))"
                    % (move_shape_index, slide_index, len(shapes)))
            moved = shapes[move_shape_index]
            if moved.shape_id != move_shape_id:
                raise MutationError(
                    "geometry move shape_id mismatch at shape_index %d: expected %r, found %r "
                    "- refusing to move the wrong shape"
                    % (move_shape_index, move_shape_id, moved.shape_id))
            moved.top = new_top

        expected_slide_count = len(prs.slides)
        _save_staged(prs, staged)

        ok, detail = validate_text_runs_and_geometry_output(
            staged, expected_slide_count, slide_index, shape_index, run_edits, geometry_moves)
        if not ok:
            raise ValidationError(detail)

        _finalize_transaction(input_path, output_path, staged, overwrite, before_hash)

    return output_path


# --------------------------------------------------------------------------
# geometry mutation (move / resize / atomic left+top+width+height)
# --------------------------------------------------------------------------
#
# Coordinates and dimensions are all in EMU (English Metric Units) - the
# native integer unit python-pptx and the underlying OOXML format use
# throughout (914400 EMU = 1 inch). No other unit is accepted or converted
# here; callers convert if they have inches/points/pixels.

def _validate_coordinate(value, name):
    if not isinstance(value, int) or isinstance(value, bool):
        raise MutationError("%s must be an integer EMU value, got %r" % (name, value))
    if value < 0:
        raise MutationError("%s must be >= 0 EMU, got %d" % (name, value))


def _validate_dimension(value, name):
    if not isinstance(value, int) or isinstance(value, bool):
        raise MutationError("%s must be an integer EMU value, got %r" % (name, value))
    if value <= 0:
        raise MutationError("%s must be > 0 EMU, got %d" % (name, value))


def _apply_geometry(shape, geometry):
    """Apply an already-validated {attr: value} geometry dict to `shape`.
    Only the given attributes are touched - e.g. a move (left/top only)
    never assigns width/height, so it cannot change them even indirectly."""
    for key, value in geometry.items():
        setattr(shape, key, value)


def validate_geometry_output(output_path, expected_slide_count, slide_index,
                              shape_index, expected_geometry):
    """Reopen `output_path` and confirm slide count is unchanged, the
    targeted shape is still resolvable, and every geometry attribute that
    was requested now matches exactly. Returns (ok, detail)."""
    try:
        prs = load_deck(output_path)
    except DeckSourceError as exc:
        return False, "generated output failed to reopen: %s" % exc
    if len(prs.slides) != expected_slide_count:
        return False, ("slide count changed: expected %d, got %d"
                        % (expected_slide_count, len(prs.slides)))
    try:
        shape = _resolve_shape(prs, slide_index, shape_index)
    except MutationError as exc:
        return False, "targeted shape missing after save: %s" % exc
    for key, expected_value in expected_geometry.items():
        actual = getattr(shape, key)
        if actual != expected_value:
            return False, ("%s mismatch after save: expected %r, got %r"
                            % (key, expected_value, actual))
    return True, "ok"


def _run_geometry_mutation(input_path, output_path, slide_index, shape_index,
                            geometry, overwrite):
    """Shared preflight -> staged-copy -> mutate -> save -> validate ->
    verify-source -> atomic-publish flow for every geometry primitive
    below. `geometry` values must already be validated by the caller
    before this runs.
    """
    _assert_safe_output_path(input_path, output_path, overwrite)
    _preflight_resolve_shape(input_path, slide_index, shape_index)
    before_hash = _validate_and_hash_source(input_path)

    with _staged_copy(input_path) as staged:
        prs = Presentation(staged)
        shape = _resolve_shape(prs, slide_index, shape_index)
        _apply_geometry(shape, geometry)
        expected_slide_count = len(prs.slides)
        expected_geometry = {
            "left": shape.left, "top": shape.top,
            "width": shape.width, "height": shape.height,
        }
        _save_staged(prs, staged)

        ok, detail = validate_geometry_output(
            staged, expected_slide_count, slide_index, shape_index, expected_geometry)
        if not ok:
            raise ValidationError(detail)

        _finalize_transaction(input_path, output_path, staged, overwrite, before_hash)

    return output_path


def move_shape(input_path, output_path, slide_index, shape_index, left, top,
               overwrite=False):
    """Controlled mutation: set one explicitly targeted shape's position
    (left, top), in EMU. Width and height are never touched - verified by
    the post-save validation, which re-checks left/top only, leaving
    width/height to whatever the file already had.

    Never mutates `input_path`. On any failure (bad index, invalid
    coordinate, save/validation/integrity failure) nothing is written to
    `output_path`.
    """
    _validate_coordinate(left, "left")
    _validate_coordinate(top, "top")
    return _run_geometry_mutation(input_path, output_path, slide_index, shape_index,
                                  {"left": left, "top": top}, overwrite)


def resize_shape(input_path, output_path, slide_index, shape_index, width, height,
                 overwrite=False):
    """Controlled mutation: set one explicitly targeted shape's size
    (width, height), in EMU. Setting width/height on a python-pptx shape
    does not itself move left/top - position is left untouched, and a
    dedicated regression test confirms that behavior explicitly rather
    than assuming it.

    Never mutates `input_path`. On any failure nothing is written to
    `output_path`.
    """
    _validate_dimension(width, "width")
    _validate_dimension(height, "height")
    return _run_geometry_mutation(input_path, output_path, slide_index, shape_index,
                                  {"width": width, "height": height}, overwrite)


def set_shape_geometry(input_path, output_path, slide_index, shape_index,
                       left, top, width, height, overwrite=False):
    """Atomic geometry update: left, top, width, and height are all set
    together as one operation - the primitive future red-box relocation
    will use. All four values are validated *before* the working copy is
    even staged, so an invalid value fails cleanly with nothing written
    anywhere and no shape attribute touched, even in memory. Either the
    complete requested geometry is applied, saved, verified, and
    published, or the whole operation raises and `output_path` is left
    untouched.
    """
    _validate_coordinate(left, "left")
    _validate_coordinate(top, "top")
    _validate_dimension(width, "width")
    _validate_dimension(height, "height")
    return _run_geometry_mutation(
        input_path, output_path, slide_index, shape_index,
        {"left": left, "top": top, "width": width, "height": height}, overwrite)


# --------------------------------------------------------------------------
# picture replacement
# --------------------------------------------------------------------------

def _validate_image_file(path):
    if not os.path.exists(path):
        raise MutationError("replacement image not found: %s" % path)
    if os.path.getsize(path) == 0:
        raise MutationError("replacement image is empty (0 bytes): %s" % path)
    try:
        with Image.open(path) as img:
            img.verify()
    except Exception as exc:
        raise MutationError(
            "replacement image is not a valid/readable image: %s (%s: %s)"
            % (path, type(exc).__name__, exc))


def validate_picture_replacement(output_path, expected_slide_count, slide_index,
                                  shape_index, expected_geometry, expected_image_bytes):
    """Reopen `output_path` and confirm: slide count unchanged, the target
    shape is still resolvable and is still a picture, its frame geometry
    exactly matches what was expected, and its image payload now matches
    the replacement image's bytes exactly. Returns (ok, detail)."""
    try:
        prs = load_deck(output_path)
    except DeckSourceError as exc:
        return False, "generated output failed to reopen: %s" % exc
    if len(prs.slides) != expected_slide_count:
        return False, ("slide count changed: expected %d, got %d"
                        % (expected_slide_count, len(prs.slides)))
    try:
        shape = _resolve_shape(prs, slide_index, shape_index)
    except MutationError as exc:
        return False, "targeted shape missing after save: %s" % exc
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False, "targeted shape is no longer a picture after save"
    for key, expected_value in expected_geometry.items():
        actual = getattr(shape, key)
        if actual != expected_value:
            return False, ("%s mismatch after save: expected %r, got %r"
                            % (key, expected_value, actual))
    if shape.image.blob != expected_image_bytes:
        return False, "image payload does not match the replacement image after save"
    return True, "ok"


def _drop_relationship_if_unreferenced(slide, part, rId):
    """Remove the relationship `rId` from `part`'s relationships collection
    only if no remaining `a:blip` element *anywhere in the slide's XML
    tree* still references it - not only inside the shape tree (`p:spTree`).
    Image relationships are not used only inside shapes: a slide background
    (`p:cSld/p:bg/p:bgPr/a:blipFill/a:blip`) is a sibling of `p:spTree`
    under `p:cSld`, and any other valid slide element that carries a blip
    would be reachable the same way. Scanning the whole `<p:sld>` element
    (rather than just its `spTree` child) is what makes this correct for
    all of those cases, not only ordinary picture shapes.

    python-pptx's own package writer serializes only parts still reachable
    via some relationship (a graph walk from the package root), so once
    the last relationship referencing an image part is dropped, that image
    part is automatically absent from the next save - no separate manual
    part deletion is needed or attempted here.

    Safe for the shared-image case: if another picture, the slide
    background, or any other element still uses `rId` (python-pptx assigns
    the same relationship id when the same image content is related to the
    same part more than once, which is exactly how a picture and a
    same-image background end up sharing one), the relationship - and
    therefore the underlying image part - is left exactly as-is.
    """
    still_referenced = any(
        blip.get(qn("r:embed")) == rId for blip in slide._element.iter(qn("a:blip")))
    if not still_referenced:
        part.rels.pop(rId)


def replace_picture(input_path, output_path, slide_index, shape_index, image_path,
                    overwrite=False):
    """Controlled mutation: replace a picture shape's embedded image.

    Default (and currently only) geometry behavior: the shape's existing
    left/top/width/height frame is preserved exactly, regardless of the
    replacement image's own pixel dimensions or aspect ratio - the new
    image is placed into the unchanged frame exactly as PowerPoint's own
    blipFill model renders it. This primitive does not compute, apply, or
    guess any new cropping; it does not touch `a:srcRect` at all, so any
    crop already present on the shape is left exactly as-is (preserved by
    construction, not by special-casing it).

    Only the picture's image relationship (`p:blipFill/a:blip/@r:embed`) is
    repointed at a newly-added image part - this is the smallest reliable
    OOXML-level operation for this, isolated entirely behind this function
    (including the stale-relationship cleanup, which inspects the whole
    slide XML tree, not just shapes - see `_drop_relationship_if_unreferenced`);
    no other code in this module or its callers touches raw OOXML for
    image handling. After repointing, the old relationship is dropped if
    nothing else on the slide still references it, so a replaced
    screenshot does not linger in the saved package as unused (and
    potentially sensitive) media; a relationship still shared by another
    picture or the slide background is left untouched.

    Rejects a target shape that is not a picture, and a replacement image
    that is missing, empty, or unreadable - both before anything is
    staged. Never mutates `input_path`. On any failure nothing is written
    to `output_path`.
    """
    _validate_image_file(image_path)
    _assert_safe_output_path(input_path, output_path, overwrite)
    target = _preflight_resolve_shape(input_path, slide_index, shape_index)
    if target.shape_type != MSO_SHAPE_TYPE.PICTURE:
        raise MutationError(
            "shape %d on slide %d is not a picture (shape_type=%s) - "
            "cannot replace its image"
            % (shape_index, slide_index, _shape_type_name(target)))
    before_hash = _validate_and_hash_source(input_path)
    with open(image_path, "rb") as fh:
        new_image_bytes = fh.read()

    with _staged_copy(input_path) as staged:
        prs = Presentation(staged)
        slide = prs.slides[slide_index]
        shape = _resolve_shape(prs, slide_index, shape_index)

        expected_geometry = {
            "left": shape.left, "top": shape.top,
            "width": shape.width, "height": shape.height,
        }
        slide_part = shape.part
        old_rId = shape._pic.blipFill.blip.rEmbed
        _new_image_part, new_rId = slide_part.get_or_add_image_part(image_path)
        shape._pic.blipFill.blip.rEmbed = new_rId
        if old_rId is not None and old_rId != new_rId:
            _drop_relationship_if_unreferenced(slide, slide_part, old_rId)

        expected_slide_count = len(prs.slides)
        _save_staged(prs, staged)

        ok, detail = validate_picture_replacement(
            staged, expected_slide_count, slide_index, shape_index,
            expected_geometry, new_image_bytes)
        if not ok:
            raise ValidationError(detail)

        _finalize_transaction(input_path, output_path, staged, overwrite, before_hash)

    return output_path
